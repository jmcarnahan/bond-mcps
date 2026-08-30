"""
Text extraction from Office documents (docx, pptx, xlsx, pdf).

All functions accept raw bytes and return extracted text strings.
No temp files are created — all operations use in-memory BytesIO streams.
Libraries are imported lazily to keep server startup fast.
"""

import io
import logging
from typing import Any

logger = logging.getLogger(__name__)

MAX_DOCUMENT_DOWNLOAD_BYTES = 50_000_000  # 50 MB download limit
MAX_EXTRACTED_TEXT_CHARS = 200_000  # ~200K chars returned to Claude

_EXTRACTABLE_MIMES: dict[str, str] = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/pdf": "pdf",
}

_EXTRACTABLE_EXTENSIONS: dict[str, str] = {
    ".docx": "docx",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
    ".pdf": "pdf",
}

_LEGACY_EXTENSIONS = frozenset({".doc", ".xls", ".ppt"})

_MAX_ROWS_PER_SHEET = 500


def is_extractable_document(item: dict[str, Any]) -> bool:
    """Return True if the driveItem is a supported Office document."""
    mime = item.get("file", {}).get("mimeType", "")
    if mime in _EXTRACTABLE_MIMES:
        return True
    name = item.get("name", "")
    dot_idx = name.rfind(".")
    if dot_idx >= 0:
        ext = name[dot_idx:].lower()
        return ext in _EXTRACTABLE_EXTENSIONS
    return False


def _get_format(mime_type: str, filename: str) -> str | None:
    """Determine the document format from MIME type or filename extension."""
    fmt = _EXTRACTABLE_MIMES.get(mime_type)
    if fmt:
        return fmt
    dot_idx = filename.rfind(".")
    if dot_idx >= 0:
        ext = filename[dot_idx:].lower()
        if ext in _LEGACY_EXTENSIONS:
            return "legacy"
        return _EXTRACTABLE_EXTENSIONS.get(ext)
    return None


def extract_document_text(data: bytes, mime_type: str, filename: str) -> str | None:
    """Extract text content from a document.

    Returns extracted text, or None if the format is unsupported or extraction fails.
    Handles errors internally — never raises to the caller.
    """
    fmt = _get_format(mime_type, filename)
    if fmt is None:
        return None
    if fmt == "legacy":
        return (
            "Legacy Office formats (.doc, .xls, .ppt) are not supported for text extraction. "
            "Please convert to the modern format (.docx, .xlsx, .pptx) in Microsoft 365."
        )

    try:
        if fmt == "docx":
            text = _extract_docx(data)
        elif fmt == "pptx":
            text = _extract_pptx(data)
        elif fmt == "xlsx":
            text = _extract_xlsx(data)
        elif fmt == "pdf":
            text = _extract_pdf(data)
        else:
            return None
    except Exception as e:
        err_str = str(e).lower()
        if "password" in err_str or "encrypted" in err_str or "decrypt" in err_str:
            return "This document is password-protected and cannot be read."
        if "bad zip" in err_str or "not a zip" in err_str:
            return (
                "Could not extract text from this document. "
                "The file may be corrupted or in an unsupported format."
            )
        logger.warning("Document extraction failed for %s: %s", filename, e)
        return None

    if not text or not text.strip():
        return (
            "Document was parsed but contains no extractable text (may be image-only or scanned)."
        )

    if len(text) > MAX_EXTRACTED_TEXT_CHARS:
        truncated = text[:MAX_EXTRACTED_TEXT_CHARS]
        return (
            f"{truncated}\n\n"
            f"... [Content truncated. Showing first {MAX_EXTRACTED_TEXT_CHARS:,} of {len(text):,} characters.]"
        )

    return text


def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            # Check for inline images in otherwise-empty paragraphs
            inline_images = para._element.findall(
                ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing"
            )
            if inline_images:
                parts.append("[Image]")
            continue
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            level_str = para.style.name.replace("Heading ", "").strip()
            try:
                level = int(level_str)
            except ValueError:
                level = 1
            parts.append(f"{'#' * level} {text}")
        else:
            parts.append(text)

    for table in doc.tables:
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append("\n[Table]\n" + "\n".join(rows))

    # Count total inline images in the document
    all_drawings = doc.element.findall(
        ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing"
    )
    if all_drawings:
        parts.append(f"\n[Document contains {len(all_drawings)} embedded image(s) — not shown]")

    return "\n\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_text: list[str] = []
        image_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_text.append(" | ".join(cells))

        if image_count:
            slide_text.append(f"[{image_count} image(s)]")

        if slide_text:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Speaker Notes]: {notes}")

    return "\n\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []

    try:
        for sheet in wb.worksheets:
            rows: list[str] = []
            empty_row_count = 0
            for row in sheet.iter_rows(values_only=True):
                if all(cell is None for cell in row):
                    empty_row_count += 1
                    if empty_row_count > 3:
                        break
                    continue
                empty_row_count = 0
                cells = [str(cell) if cell is not None else "" for cell in row]
                rows.append(" | ".join(cells))
                if len(rows) >= _MAX_ROWS_PER_SHEET:
                    break

            if rows:
                header = f"--- Sheet: {sheet.title} ---"
                parts.append(f"{header}\n" + "\n".join(rows))
                total_rows = sheet.max_row or len(rows)
                if total_rows > _MAX_ROWS_PER_SHEET:
                    parts.append(f"[... showing first {_MAX_ROWS_PER_SHEET} of ~{total_rows} rows]")
    finally:
        wb.close()

    return "\n\n".join(parts)


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []

    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            parts.append(f"--- Page {i} ---\n{text.strip()}")

    return "\n\n".join(parts)
