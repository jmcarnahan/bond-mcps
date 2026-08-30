"""Tests for document text extraction (docx, pptx, xlsx, pdf)."""

import io

import pytest
from ms_graph.document_extract import (
    _extract_docx,
    _extract_pdf,
    _extract_pptx,
    _extract_xlsx,
    extract_document_text,
    is_extractable_document,
)

# ---------------------------------------------------------------------------
# Fixtures — create minimal valid Office files programmatically
# ---------------------------------------------------------------------------


@pytest.fixture
def sample_docx_bytes():
    from docx import Document

    doc = Document()
    doc.add_heading("Test Document", level=1)
    doc.add_paragraph("First paragraph with content.")
    doc.add_paragraph("Second paragraph.")
    doc.add_heading("Section Two", level=2)
    doc.add_paragraph("Under section two.")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_docx_with_table():
    from docx import Document

    doc = Document()
    doc.add_paragraph("Before table.")
    table = doc.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Age"
    table.cell(0, 2).text = "City"
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "30"
    table.cell(1, 2).text = "NYC"
    doc.add_paragraph("After table.")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_bytes():
    from pptx import Presentation

    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Presentation Title"
    slide1.placeholders[1].text = "Subtitle text"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Slide Two"
    slide2.placeholders[1].text = "Bullet point content"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_with_notes():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Title Slide"
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "These are speaker notes."

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_xlsx_bytes():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Score", "Grade"])
    ws.append(["Alice", 95, "A"])
    ws.append(["Bob", 82, "B"])
    ws.append(["Charlie", 71, "C"])

    ws2 = wb.create_sheet("Summary")
    ws2.append(["Total", 248])
    ws2.append(["Average", 82.67])

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_pdf_bytes():
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    # pypdf doesn't easily create text content from scratch,
    # so we test with a minimal valid PDF that has no extractable text
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# is_extractable_document tests
# ---------------------------------------------------------------------------


class TestIsExtractableDocument:
    def test_docx_by_mime(self):
        item = {
            "name": "report.docx",
            "file": {
                "mimeType": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            },
        }
        assert is_extractable_document(item) is True

    def test_pptx_by_mime(self):
        item = {
            "name": "deck.pptx",
            "file": {
                "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            },
        }
        assert is_extractable_document(item) is True

    def test_xlsx_by_mime(self):
        item = {
            "name": "data.xlsx",
            "file": {
                "mimeType": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            },
        }
        assert is_extractable_document(item) is True

    def test_pdf_by_mime(self):
        item = {"name": "doc.pdf", "file": {"mimeType": "application/pdf"}}
        assert is_extractable_document(item) is True

    def test_extension_fallback_docx(self):
        item = {"name": "report.docx", "file": {"mimeType": "application/octet-stream"}}
        assert is_extractable_document(item) is True

    def test_extension_fallback_pdf(self):
        item = {"name": "scan.pdf", "file": {"mimeType": ""}}
        assert is_extractable_document(item) is True

    def test_text_file_not_extractable(self):
        item = {"name": "readme.md", "file": {"mimeType": "text/markdown"}}
        assert is_extractable_document(item) is False

    def test_image_not_extractable(self):
        item = {"name": "photo.png", "file": {"mimeType": "image/png"}}
        assert is_extractable_document(item) is False

    def test_folder_not_extractable(self):
        item = {"name": "Documents", "folder": {"childCount": 5}}
        assert is_extractable_document(item) is False


# ---------------------------------------------------------------------------
# DOCX extraction tests
# ---------------------------------------------------------------------------


class TestDocxExtraction:
    def test_extract_paragraphs(self, sample_docx_bytes):
        text = _extract_docx(sample_docx_bytes)
        assert "First paragraph with content." in text
        assert "Second paragraph." in text
        assert "Under section two." in text

    def test_extract_headings_as_markdown(self, sample_docx_bytes):
        text = _extract_docx(sample_docx_bytes)
        assert "# Test Document" in text
        assert "## Section Two" in text

    def test_extract_tables(self, sample_docx_with_table):
        text = _extract_docx(sample_docx_with_table)
        assert "Before table." in text
        assert "After table." in text
        assert "Name | Age | City" in text
        assert "Alice | 30 | NYC" in text

    def test_empty_document(self):
        from docx import Document

        doc = Document()
        buf = io.BytesIO()
        doc.save(buf)
        text = _extract_docx(buf.getvalue())
        assert text == ""

    def test_corrupted_bytes(self):
        result = extract_document_text(
            b"not a valid docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "bad.docx",
        )
        assert result is not None
        assert "corrupted" in result.lower() or "not a zip" in result.lower()


# ---------------------------------------------------------------------------
# PPTX extraction tests
# ---------------------------------------------------------------------------


class TestPptxExtraction:
    def test_extract_slides(self, sample_pptx_bytes):
        text = _extract_pptx(sample_pptx_bytes)
        assert "--- Slide 1 ---" in text
        assert "Presentation Title" in text
        assert "--- Slide 2 ---" in text
        assert "Slide Two" in text

    def test_extract_speaker_notes(self, sample_pptx_with_notes):
        text = _extract_pptx(sample_pptx_with_notes)
        assert "[Speaker Notes]" in text
        assert "These are speaker notes." in text

    def test_empty_presentation(self):
        from pptx import Presentation

        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)
        text = _extract_pptx(buf.getvalue())
        assert text == ""


# ---------------------------------------------------------------------------
# XLSX extraction tests
# ---------------------------------------------------------------------------


class TestXlsxExtraction:
    def test_extract_multi_sheet(self, sample_xlsx_bytes):
        text = _extract_xlsx(sample_xlsx_bytes)
        assert "--- Sheet: Data ---" in text
        assert "Name | Score | Grade" in text
        assert "Alice | 95 | A" in text
        assert "--- Sheet: Summary ---" in text
        assert "Total | 248" in text

    def test_empty_workbook(self):
        from openpyxl import Workbook

        wb = Workbook()
        buf = io.BytesIO()
        wb.save(buf)
        text = _extract_xlsx(buf.getvalue())
        # Empty workbook has a sheet but no data
        assert text == ""

    def test_row_truncation(self):
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Big"
        for i in range(600):
            ws.append([f"row{i}", i, i * 2])
        buf = io.BytesIO()
        wb.save(buf)

        text = _extract_xlsx(buf.getvalue())
        assert "--- Sheet: Big ---" in text
        assert "row0" in text
        assert "row499" in text
        # Row 500+ should be truncated
        assert "row500" not in text
        assert "showing first 500" in text.lower() or "truncated" in text.lower()


# ---------------------------------------------------------------------------
# PDF extraction tests
# ---------------------------------------------------------------------------


class TestPdfExtraction:
    def test_blank_pdf(self, sample_pdf_bytes):
        # A blank page with no text content
        text = _extract_pdf(sample_pdf_bytes)
        assert text == ""

    def test_corrupted_pdf(self):
        result = extract_document_text(b"not a pdf", "application/pdf", "bad.pdf")
        assert result is None or "corrupted" in (result or "").lower()


# ---------------------------------------------------------------------------
# Dispatcher and truncation tests
# ---------------------------------------------------------------------------


class TestDispatcher:
    def test_routes_by_mime_docx(self, sample_docx_bytes):
        result = extract_document_text(
            sample_docx_bytes,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "test.docx",
        )
        assert result is not None
        assert "First paragraph" in result

    def test_routes_by_extension(self, sample_docx_bytes):
        result = extract_document_text(sample_docx_bytes, "application/octet-stream", "test.docx")
        assert result is not None
        assert "First paragraph" in result

    def test_unsupported_format(self):
        result = extract_document_text(b"data", "image/png", "photo.png")
        assert result is None

    def test_legacy_format_message(self):
        result = extract_document_text(b"data", "application/msword", "old.doc")
        assert result is not None
        assert "Legacy Office formats" in result

    def test_truncation(self, sample_docx_bytes):
        from unittest.mock import patch

        import ms_graph.document_extract as mod

        with patch.object(mod, "MAX_EXTRACTED_TEXT_CHARS", 20):
            result = extract_document_text(
                sample_docx_bytes,
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "test.docx",
            )
        assert result is not None
        assert "truncated" in result.lower()

    def test_empty_extraction_message(self):
        from docx import Document

        doc = Document()
        buf = io.BytesIO()
        doc.save(buf)
        result = extract_document_text(
            buf.getvalue(),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "empty.docx",
        )
        assert result is not None
        assert "no extractable text" in result.lower()
