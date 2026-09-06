"""Tests for the FastMCP server using in-process client.

In-process FastMCP clients don't have HTTP request context, so we mock
get_graph_token() directly instead of get_http_headers().
"""

import base64
import json
from unittest.mock import MagicMock, patch
from urllib.parse import parse_qs, quote, urlparse

import httpx
import pytest
import respx
from ms_graph import mail, mail_policy
from ms_graph.files import _encode_sharing_url
from ms_graph.graph_client import GRAPH_BASE_URL
from ms_graph.power_bi import POWERBI_BASE_URL

from auth.exceptions import MissingProviderConnection

from .conftest import (
    GRAPH_ERROR_400,
    GRAPH_ERROR_403,
    GRAPH_ERROR_404,
    GRAPH_ERROR_410,
    SAMPLE_ATTACHMENTS_RESPONSE,
    SAMPLE_CHANNEL_FILES_FOLDER,
    SAMPLE_CHANNEL_MESSAGES_RESPONSE,
    SAMPLE_CHANNELS_RESPONSE,
    SAMPLE_CHAT_MEMBERS_RESPONSE,
    SAMPLE_CHAT_MESSAGE_FULL,
    SAMPLE_CHAT_MESSAGE_SENT,
    SAMPLE_CHAT_MESSAGE_WITH_CARD,
    SAMPLE_CHAT_MESSAGE_WITH_FILE,
    SAMPLE_CHAT_MESSAGE_WITH_IMAGE,
    SAMPLE_CHAT_MESSAGE_WITH_JUNK_ATTACHMENTS,
    SAMPLE_CHAT_MESSAGES_PAGE,
    SAMPLE_CHAT_MESSAGES_PAGE_WITH_ATTACHMENTS,
    SAMPLE_CHAT_MESSAGES_RESPONSE,
    SAMPLE_CHATS_PAGE,
    SAMPLE_CHATS_PAGE_NEXT_LINK,
    SAMPLE_CHATS_RESPONSE,
    SAMPLE_COPY_COMPLETED,
    SAMPLE_COPY_FAILED,
    SAMPLE_CREATED_ATTACHMENT,
    SAMPLE_DELTA_LINK,
    SAMPLE_DELTA_MESSAGE,
    SAMPLE_DELTA_NEXT_LINK,
    SAMPLE_DELTA_PAGE_FINAL,
    SAMPLE_DELTA_PAGE_NEXT,
    SAMPLE_DELTA_TOMBSTONE,
    SAMPLE_DRAFT_FOR_SEND,
    SAMPLE_DRAFT_MESSAGE,
    SAMPLE_DRIVE_CHILDREN_RESPONSE,
    SAMPLE_DRIVE_ITEM_BINARY,
    SAMPLE_DRIVE_ITEM_FILE,
    SAMPLE_DRIVE_ITEM_FOLDER,
    SAMPLE_DRIVE_ITEM_LARGE_TEXT,
    SAMPLE_DRIVE_ITEM_WORD,
    SAMPLE_EXTERNAL_DELTA_MESSAGE,
    SAMPLE_EXTERNAL_ITEM_ATTACHMENT,
    SAMPLE_EXTERNAL_ITEM_ATTACHMENT_META,
    SAMPLE_EXTERNAL_MESSAGE,
    SAMPLE_EXTERNAL_MESSAGE_DETAIL,
    SAMPLE_FILE_ATTACHMENT,
    SAMPLE_FORWARDING_RULE,
    SAMPLE_INLINE_ATTACHMENT,
    SAMPLE_INVITE_RESPONSE,
    SAMPLE_ITEM_ATTACHMENT,
    SAMPLE_MAIL_FOLDER,
    SAMPLE_MAIL_FOLDERS_RESPONSE,
    SAMPLE_MAILBOX_SETTINGS,
    SAMPLE_MESSAGE,
    SAMPLE_MESSAGE_2,
    SAMPLE_MESSAGE_DETAIL,
    SAMPLE_MESSAGE_DETAIL_NO_BODY,
    SAMPLE_MESSAGE_RULE,
    SAMPLE_MESSAGE_RULES_RESPONSE,
    SAMPLE_MESSAGE_WITH_ATTACHMENTS,
    SAMPLE_MESSAGES_PAGE1,
    SAMPLE_MESSAGES_PAGE2,
    SAMPLE_MESSAGES_RESPONSE,
    SAMPLE_NEW_DRAFT,
    SAMPLE_ONBEHALF_MESSAGE,
    SAMPLE_PBI_DASHBOARDS_RESPONSE,
    SAMPLE_PBI_DATASETS_RESPONSE,
    SAMPLE_PBI_DAX_RESULT,
    SAMPLE_PBI_EXPORT_SUCCEEDED,
    SAMPLE_PBI_REPORTS_RESPONSE,
    SAMPLE_PBI_WORKSPACES_RESPONSE,
    SAMPLE_REFERENCE_ATTACHMENT,
    SAMPLE_REPLY_DRAFT,
    SAMPLE_SEARCH_RESPONSE,
    SAMPLE_SEARCH_RESPONSE_EMPTY,
    SAMPLE_SENDER_ONLY_EXTERNAL,
    SAMPLE_SENDER_ONLY_INTERNAL,
    SAMPLE_SHARED_TEXT_FILE,
    SAMPLE_SITES_RESPONSE,
    SAMPLE_TEAMS_DRIVE_ITEM,
    SAMPLE_TEAMS_RESPONSE,
    SAMPLE_TEAMS_UPLOAD_RESPONSE,
    SAMPLE_TEAMS_UPLOADED_ITEM,
    SAMPLE_UNSENT_DRAFT,
    SAMPLE_UPLOADED_FILE,
    SAMPLE_USER_PROFILE,
    TEAMS_FILE_ATTACHMENT_ID,
    TEAMS_FILE_URL,
    TEAMS_HOSTED_ID,
    TEAMS_HOSTED_URL,
    TEAMS_PPTX_MIME,
    TEAMS_UPLOAD_GUID,
    TEAMS_WEBDAV_URL,
)

MONITOR_URL = "https://api.onedrive.com/v1.0/monitor/copy-op-token"
CONNECT_URL = "https://auth.example.com/connect/microsoft?ticket=t"
SOURCE_DRIVE_ID = SAMPLE_DRIVE_ITEM_WORD["parentReference"]["driveId"]
PBI_EXPORT_MONITOR_URL = (
    f"{POWERBI_BASE_URL}/groups/ws-id-001/reports/rpt-id-001/exports/export-id-001"
)

# Attachment URLs. aget_message interpolates the message ID raw (legacy), while
# every attachment path percent-encodes it — hence the two spellings.
ATT_MSG_ID = SAMPLE_MESSAGE["id"]
ATT_BASE = f"{GRAPH_BASE_URL}/me/messages/{quote(ATT_MSG_ID, safe='')}/attachments"
ATT_FILE_URL = f"{ATT_BASE}/{quote(SAMPLE_FILE_ATTACHMENT['id'], safe='')}"
ATT_ITEM_URL = f"{ATT_BASE}/{quote(SAMPLE_ITEM_ATTACHMENT['id'], safe='')}"
ATT_REF_URL = f"{ATT_BASE}/{quote(SAMPLE_REFERENCE_ATTACHMENT['id'], safe='')}"
DRAFT_BASE = f"{GRAPH_BASE_URL}/me/messages/{quote(SAMPLE_DRAFT_MESSAGE['id'], safe='')}"
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

# Teams attachment routes. The sharing URL resolves through /shares/{token};
# the more specific content and thumbnail routes are exact URLs so registration
# order cannot shadow the metadata route.
TEAMS_CHAT_ID = "chat-1on1-001"
TEAMS_CHAT_MSGS = f"{GRAPH_BASE_URL}/chats/{TEAMS_CHAT_ID}/messages"
TEAMS_FILE_MSG_URL = f"{TEAMS_CHAT_MSGS}/chat-msg-file-001"
TEAMS_IMAGE_MSG_URL = f"{TEAMS_CHAT_MSGS}/chat-msg-image-001"
TEAMS_CARD_MSG_URL = f"{TEAMS_CHAT_MSGS}/chat-msg-card-001"
TEAMS_JUNK_MSG_URL = f"{TEAMS_CHAT_MSGS}/chat-msg-junk-001"
TEAMS_HOSTED_VALUE_URL = f"{TEAMS_IMAGE_MSG_URL}/hostedContents/{TEAMS_HOSTED_ID}/$value"
TEAMS_SHARE_BASE = f"{GRAPH_BASE_URL}/shares/{_encode_sharing_url(TEAMS_FILE_URL)}/driveItem"
TEAMS_SHARE_CONTENT_URL = f"{TEAMS_SHARE_BASE}/content"
TEAMS_SHARE_THUMB_URL = f"{TEAMS_SHARE_BASE}/thumbnails/0/medium/content"
PNG_BYTES = b"\x89PNG\r\n\x1a\n-fake-png"

# Sending a file into a chat: OneDrive upload, re-fetch, roster, share.
TEAMS_SEND_UPLOAD_URL = (
    f"{GRAPH_BASE_URL}/me/drive/root:/Microsoft%20Teams%20Chat%20Files/notes.txt:/content"
)
TEAMS_SEND_ITEM_URL = f"{GRAPH_BASE_URL}/me/drive/items/teams-upload-001"
TEAMS_SEND_MEMBERS_URL = f"{GRAPH_BASE_URL}/chats/{TEAMS_CHAT_ID}/members"
TEAMS_SEND_INVITE_URL = f"{TEAMS_SEND_ITEM_URL}/invite"


def _token_with_oid(oid: str) -> str:
    """An unsigned JWT whose payload carries the given oid claim."""
    payload = base64.urlsafe_b64encode(json.dumps({"oid": oid}).encode()).rstrip(b"=")
    return f"eyJhbGciOiJub25lIn0.{payload.decode()}.sig"


def _invite_recipients() -> list:
    """The recipients of the one /invite call respx saw."""
    for call in respx.calls:
        if call.request.url.path.endswith("/invite"):
            return json.loads(call.request.content)["recipients"]
    raise AssertionError("no /invite request was made")


def _mock_chat_file_upload() -> None:
    """The four requests that put a file where a chat message can point at it."""
    respx.put(url__startswith=TEAMS_SEND_UPLOAD_URL).mock(
        return_value=httpx.Response(201, json=SAMPLE_TEAMS_UPLOAD_RESPONSE)
    )
    respx.get(url__startswith=TEAMS_SEND_ITEM_URL).mock(
        return_value=httpx.Response(200, json=SAMPLE_TEAMS_UPLOADED_ITEM)
    )
    respx.get(TEAMS_SEND_MEMBERS_URL).mock(
        return_value=httpx.Response(200, json=SAMPLE_CHAT_MEMBERS_RESPONSE)
    )
    respx.post(TEAMS_SEND_INVITE_URL).mock(
        return_value=httpx.Response(200, json=SAMPLE_INVITE_RESPONSE)
    )


def _graph_trail() -> list[tuple[str, str]]:
    """(method, path) for every request respx saw, in order."""
    return [(c.request.method, c.request.url.path) for c in respx.calls]


def _docx_bytes() -> bytes:
    """A real docx so the text sink exercises the extractor, not a stub."""
    from ms_graph import document_create

    return document_create.markdown_to_docx("# Quarterly Title\n\nBody text here.")


def _mock_token(token: str = "test-ms-token"):
    """Patch get_graph_token to return a test token."""
    return patch("ms_graph_mcp.get_graph_token", return_value=token)


def _mock_pbi_token(token: str = "test-pbi-token"):
    """Patch get_powerbi_token to return a test token."""
    return patch("ms_graph_mcp.get_powerbi_token", return_value=token)


def _get_text(result) -> str:
    """Extract text from FastMCP CallToolResult."""
    return result.content[0].text


@pytest.fixture
def mcp_server():
    """Import and return the MCP server instance."""
    from ms_graph_mcp import mcp

    return mcp


# ---------------------------------------------------------------------------
# Profile
# ---------------------------------------------------------------------------


class TestMCPProfileTools:
    """Test user profile MCP tools via in-process FastMCP client."""

    @respx.mock
    async def test_get_user_profile_with_mailbox_address(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me").mock(
            return_value=httpx.Response(200, json=SAMPLE_USER_PROFILE)
        )
        respx.get(f"{GRAPH_BASE_URL}/me/mailboxSettings").mock(
            return_value=httpx.Response(200, json=SAMPLE_MAILBOX_SETTINGS)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("get_user_profile", {})

        text = _get_text(result)
        assert "Test User" in text
        assert "user@example.com" in text
        assert "Mailbox Address" in text
        assert "mailbox@example.com" in text

    @respx.mock
    async def test_get_user_profile_without_mailbox_scope(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me").mock(
            return_value=httpx.Response(200, json=SAMPLE_USER_PROFILE)
        )
        respx.get(f"{GRAPH_BASE_URL}/me/mailboxSettings").mock(
            return_value=httpx.Response(
                403, json={"error": {"code": "ErrorAccessDenied", "message": "Access denied"}}
            )
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("get_user_profile", {})

        text = _get_text(result)
        assert "Test User" in text
        assert "Mailbox Address" not in text


# ---------------------------------------------------------------------------
# Email
# ---------------------------------------------------------------------------


class TestMCPEmailTools:
    """Test consolidated email MCP tools."""

    @respx.mock
    async def test_list_emails_no_query(self, mcp_server):
        """No query → lists inbox as pipe-delimited CSV."""
        respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGES_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_emails", {"top": 10})

        text = _get_text(result)
        assert "2 message(s) in inbox" in text
        assert (
            "date|from_name|from_address|to|subject|is_read|body_preview|has_attachments|id" in text
        )
        assert "Weekly Report" in text
        assert "alice@example.com" in text

    @respx.mock
    async def test_list_emails_with_query(self, mcp_server):
        """query set → search mode, CSV output."""
        respx.get(f"{GRAPH_BASE_URL}/me/messages").mock(
            return_value=httpx.Response(200, json={"value": [SAMPLE_MESSAGE]})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_emails", {"query": "weekly"})

        text = _get_text(result)
        assert '1 result(s) for "weekly"' in text
        assert (
            "date|from_name|from_address|to|subject|is_read|body_preview|has_attachments|id" in text
        )
        assert "Weekly Report" in text

    @respx.mock
    async def test_list_emails_search_no_results(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/messages").mock(
            return_value=httpx.Response(200, json={"value": []})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_emails", {"query": "nonexistent"})

        text = _get_text(result)
        assert "No messages found" in text
        assert "nonexistent" in text

    @respx.mock
    async def test_list_emails_empty_inbox(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages").mock(
            return_value=httpx.Response(200, json={"value": []})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_emails", {})

        assert "No messages found" in _get_text(result)

    @respx.mock
    async def test_list_emails_custom_folder(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/sentitems/messages").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGES_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_emails", {"folder": "sentitems"})

        text = _get_text(result)
        assert "2 message(s) in sentitems" in text
        assert "Weekly Report" in text

    @respx.mock
    async def test_list_emails_custom_display_name_resolves(self, mcp_server):
        """A custom folder display name resolves to its ID, then lists messages (#54)."""
        folder_id = SAMPLE_MAIL_FOLDER["id"]  # "AQMkAGfolder-001", displayName "Projects"
        respx.get(f"{GRAPH_BASE_URL}/me/mailFolders").mock(
            return_value=httpx.Response(200, json=SAMPLE_MAIL_FOLDERS_RESPONSE)
        )
        respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/{folder_id}/messages").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGES_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_emails", {"folder": "projects"})

        text = _get_text(result)
        assert "2 message(s) in projects" in text  # header keeps the display name
        assert "Weekly Report" in text

    @respx.mock
    async def test_list_emails_unknown_folder_returns_clear_error(self, mcp_server):
        """An unresolvable folder name returns a human-readable error, not a raw 400 (#54)."""
        respx.get(f"{GRAPH_BASE_URL}/me/mailFolders").mock(
            return_value=httpx.Response(200, json=SAMPLE_MAIL_FOLDERS_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_emails", {"folder": "ghost"})

        assert _get_text(result) == "Folder 'ghost' not found."

    @respx.mock
    async def test_list_emails_well_known_skips_folder_lookup(self, mcp_server):
        """Well-known folders resolve without hitting the mailFolders list endpoint (#54)."""
        list_route = respx.get(f"{GRAPH_BASE_URL}/me/mailFolders").mock(
            return_value=httpx.Response(200, json=SAMPLE_MAIL_FOLDERS_RESPONSE)
        )
        respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/sentitems/messages").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGES_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                await client.call_tool("list_emails", {"folder": "sentitems"})

        assert not list_route.called

    @respx.mock
    async def test_list_emails_query_scopes_to_custom_folder(self, mcp_server):
        """query + explicit custom folder scopes the search to that folder (#54)."""
        folder_id = SAMPLE_MAIL_FOLDER["id"]
        respx.get(f"{GRAPH_BASE_URL}/me/mailFolders").mock(
            return_value=httpx.Response(200, json=SAMPLE_MAIL_FOLDERS_RESPONSE)
        )
        scoped = respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/{folder_id}/messages").mock(
            return_value=httpx.Response(200, json={"value": [SAMPLE_MESSAGE]})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "list_emails", {"query": "report", "folder": "projects"}
                )

        assert scoped.called
        scoped_url = str(scoped.calls[0].request.url)
        assert "$search" in scoped_url or "%24search" in scoped_url
        assert '1 result(s) for "report"' in _get_text(result)

    @respx.mock
    async def test_list_emails_query_default_inbox_stays_global(self, mcp_server):
        """query with the default inbox folder searches globally, no folder lookup (#54)."""
        list_route = respx.get(f"{GRAPH_BASE_URL}/me/mailFolders").mock(
            return_value=httpx.Response(200, json=SAMPLE_MAIL_FOLDERS_RESPONSE)
        )
        global_route = respx.get(f"{GRAPH_BASE_URL}/me/messages").mock(
            return_value=httpx.Response(200, json={"value": [SAMPLE_MESSAGE]})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                await client.call_tool("list_emails", {"query": "report"})

        assert global_route.called
        assert not list_route.called

    @respx.mock
    async def test_list_emails_pagination(self, mcp_server):
        """Pagination follows @odata.nextLink to fetch all pages."""
        responses = iter(
            [
                httpx.Response(200, json=SAMPLE_MESSAGES_PAGE1),
                httpx.Response(200, json=SAMPLE_MESSAGES_PAGE2),
            ]
        )
        respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages").mock(
            side_effect=lambda req: next(responses)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_emails", {"top": 1000})

        text = _get_text(result)
        assert "2 message(s) in inbox" in text
        assert "Weekly Report" in text
        assert "Re: Project Update" in text

    @respx.mock
    async def test_list_emails_csv_structure(self, mcp_server):
        """Verify CSV output is parseable and has correct field values."""
        import csv
        import io

        respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGES_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_emails", {"top": 10})

        text = _get_text(result)
        # Skip the header comment line and blank line
        csv_text = text.split("\n\n", 1)[1]
        reader = csv.reader(io.StringIO(csv_text), delimiter="|")
        rows = list(reader)

        # Header row
        assert rows[0] == [
            "date",
            "from_name",
            "from_address",
            "to",
            "subject",
            "is_read",
            "body_preview",
            "has_attachments",
            "id",
        ]
        # First data row (SAMPLE_MESSAGE)
        assert rows[1][0] == "2025-12-15T10:30:00Z"  # date
        assert rows[1][1] == "Alice Smith"  # from_name
        assert rows[1][2] == "alice@example.com"  # from_address
        assert rows[1][3] == "bob@example.com"  # to
        assert rows[1][4] == "Weekly Report"  # subject
        assert rows[1][5] == "False"  # is_read
        assert rows[1][6] == "Here is the weekly report. Best, Alice"  # body_preview
        assert rows[1][7] == ""  # has_attachments — absent on this sample
        assert rows[1][8] == "AAMkAGI2TG93AAA="  # id
        # Second data row (SAMPLE_MESSAGE_2)
        assert rows[2][4] == "Re: Project Update"
        assert rows[2][5] == "True"

    @respx.mock
    async def test_list_emails_requests_has_attachments(self, mcp_server):
        """The attachment flag has to be $selected or Graph omits it (#Phase 2)."""
        route = respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages").mock(
            return_value=httpx.Response(
                200, json={"value": [{**SAMPLE_MESSAGE, "hasAttachments": True}]}
            )
        )
        with _mock_token():
            result = await _call(mcp_server, "list_emails", {"top": 10})

        select = parse_qs(urlparse(str(route.calls[0].request.url)).query)["$select"][0]
        assert "hasAttachments" in select
        assert "|True|AAMkAGI2TG93AAA=" in _get_text(result)

    @respx.mock
    async def test_list_emails_pipe_in_subject(self, mcp_server):
        """Subjects containing pipe characters are properly quoted in CSV."""
        import csv
        import io

        msg_with_pipe = {
            **SAMPLE_MESSAGE,
            "subject": "Re: Q4 | Budget Review",
            "bodyPreview": "Preview with | pipe char",
        }
        respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages").mock(
            return_value=httpx.Response(200, json={"value": [msg_with_pipe]})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_emails", {"top": 10})

        text = _get_text(result)
        csv_text = text.split("\n\n", 1)[1]
        reader = csv.reader(io.StringIO(csv_text), delimiter="|")
        rows = list(reader)
        assert rows[1][4] == "Re: Q4 | Budget Review"
        assert rows[1][6] == "Preview with | pipe char"

    @respx.mock
    async def test_list_emails_with_mark_as_read(self, mcp_server):
        msg_id = "AAMkAGI2TG93AAA="
        respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages").mock(
            return_value=httpx.Response(200, json={"value": [SAMPLE_MESSAGE]})
        )
        patch_route = respx.patch(f"{GRAPH_BASE_URL}/me/messages/{msg_id}").mock(
            return_value=httpx.Response(200, json={**SAMPLE_MESSAGE, "isRead": True})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "list_emails",
                    {"top": 10, "options": f'{{"mark_as_read": ["{msg_id}"]}}'},
                )

        text = _get_text(result)
        assert "Weekly Report" in text
        assert patch_route.called
        payload = json.loads(patch_route.calls[0].request.content)
        assert payload == {"isRead": True}
        assert "1 message(s) marked as read" in text

    @respx.mock
    async def test_list_emails_mark_as_read_rejects_non_array(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages").mock(
            return_value=httpx.Response(200, json={"value": [SAMPLE_MESSAGE]})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "list_emails",
                    {"top": 10, "options": '{"mark_as_read": "single-id"}'},
                )

        text = _get_text(result)
        assert "must be a JSON array" in text

    @respx.mock
    async def test_read_email_plain_body(self, mcp_server):
        msg_id = "AAMkAGI2TG93AAA="
        respx.get(f"{GRAPH_BASE_URL}/me/messages/{msg_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("read_email", {"message_id": msg_id})

        text = _get_text(result)
        assert "Weekly Report" in text
        assert "Alice Smith" in text
        assert "Here is the weekly report" in text

    @respx.mock
    async def test_read_email_html_body(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/messages/{SAMPLE_MESSAGE_2['id']}").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE_2)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "read_email", {"message_id": SAMPLE_MESSAGE_2["id"]}
                )

        text = _get_text(result)
        assert "HTML content" in text
        assert "Charlie Brown" in text

    @respx.mock
    async def test_read_email_with_mark_as_read(self, mcp_server):
        msg_id = "AAMkAGI2TG93AAA="
        respx.get(f"{GRAPH_BASE_URL}/me/messages/{msg_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE)
        )
        patch_route = respx.patch(f"{GRAPH_BASE_URL}/me/messages/{msg_id}").mock(
            return_value=httpx.Response(200, json={**SAMPLE_MESSAGE, "isRead": True})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "read_email",
                    {"message_id": msg_id, "options": '{"mark_as_read": true}'},
                )

        text = _get_text(result)
        assert "Weekly Report" in text
        assert patch_route.called
        payload = json.loads(patch_route.calls[0].request.content)
        assert payload == {"isRead": True}
        assert "marked as read" in text.lower()

    @respx.mock
    async def test_read_email_mark_as_unread(self, mcp_server):
        msg_id = "AAMkAGI2TG93AAA="
        respx.get(f"{GRAPH_BASE_URL}/me/messages/{msg_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE)
        )
        patch_route = respx.patch(f"{GRAPH_BASE_URL}/me/messages/{msg_id}").mock(
            return_value=httpx.Response(200, json={**SAMPLE_MESSAGE, "isRead": False})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "read_email",
                    {"message_id": msg_id, "options": '{"mark_as_read": false}'},
                )

        text = _get_text(result)
        assert "Weekly Report" in text
        assert patch_route.called
        payload = json.loads(patch_route.calls[0].request.content)
        assert payload == {"isRead": False}
        assert "marked as unread" in text.lower()

    @respx.mock
    async def test_send_email_plain_text(self, mcp_server):
        route = respx.post(f"{GRAPH_BASE_URL}/me/sendMail").mock(return_value=httpx.Response(202))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "send_email",
                    {"to": "alice@example.com", "subject": "Hi", "body": "Hello!"},
                )

        assert "sent" in _get_text(result).lower()
        assert route.called
        payload = json.loads(route.calls[0].request.content)
        assert payload["message"]["body"]["contentType"] == "Text"

    @respx.mock
    async def test_send_email_html_body_auto_detected(self, mcp_server):
        route = respx.post(f"{GRAPH_BASE_URL}/me/sendMail").mock(return_value=httpx.Response(202))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                await client.call_tool(
                    "send_email",
                    {
                        "to": "alice@example.com",
                        "subject": "HTML test",
                        "body": "<p>Hello <strong>Alice</strong>! Click <a href='https://example.com'>here</a>.</p>",
                    },
                )

        payload = json.loads(route.calls[0].request.content)
        assert payload["message"]["body"]["contentType"] == "HTML"

    @respx.mock
    async def test_send_email_placeholder_not_mistaken_for_html(self, mcp_server):
        """'Dear <FirstName>,' must stay Text."""
        route = respx.post(f"{GRAPH_BASE_URL}/me/sendMail").mock(return_value=httpx.Response(202))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                await client.call_tool(
                    "send_email",
                    {
                        "to": "alice@example.com",
                        "subject": "Template",
                        "body": "Dear <FirstName>, thanks for reaching out.",
                    },
                )

        payload = json.loads(route.calls[0].request.content)
        assert payload["message"]["body"]["contentType"] == "Text"

    @respx.mock
    async def test_send_email_with_cc(self, mcp_server):
        route = respx.post(f"{GRAPH_BASE_URL}/me/sendMail").mock(return_value=httpx.Response(202))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "send_email",
                    {
                        "to": "alice@example.com",
                        "subject": "Hi",
                        "body": "Hello!",
                        "options": '{"cc": "bob@example.com"}',
                    },
                )

        text = _get_text(result)
        assert "CC" in text
        payload = json.loads(route.calls[0].request.content)
        assert len(payload["message"]["ccRecipients"]) == 1

    @respx.mock
    async def test_send_email_with_bcc(self, mcp_server):
        route = respx.post(f"{GRAPH_BASE_URL}/me/sendMail").mock(return_value=httpx.Response(202))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "send_email",
                    {
                        "to": "alice@example.com",
                        "subject": "Hi",
                        "body": "Hello!",
                        "options": '{"bcc": "hidden@example.com,secret@example.com"}',
                    },
                )

        text = _get_text(result)
        assert "BCC" in text
        assert "2 recipients" in text
        # BCC addresses should NOT appear in the return message
        assert "hidden@example.com" not in text
        payload = json.loads(route.calls[0].request.content)
        assert len(payload["message"]["bccRecipients"]) == 2
        assert (
            payload["message"]["bccRecipients"][0]["emailAddress"]["address"]
            == "hidden@example.com"
        )

    @respx.mock
    async def test_send_email_multiple_recipients(self, mcp_server):
        route = respx.post(f"{GRAPH_BASE_URL}/me/sendMail").mock(return_value=httpx.Response(202))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                await client.call_tool(
                    "send_email",
                    {"to": "alice@example.com, bob@example.com", "subject": "Hi", "body": "Hello!"},
                )

        payload = json.loads(route.calls[0].request.content)
        assert len(payload["message"]["toRecipients"]) == 2

    @respx.mock
    async def test_send_email_explicit_text_overrides_html(self, mcp_server):
        route = respx.post(f"{GRAPH_BASE_URL}/me/sendMail").mock(return_value=httpx.Response(202))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                await client.call_tool(
                    "send_email",
                    {
                        "to": "alice@example.com",
                        "subject": "S",
                        "body": "<p>HTML</p>",
                        "options": '{"body_type": "Text"}',
                    },
                )

        payload = json.loads(route.calls[0].request.content)
        assert payload["message"]["body"]["contentType"] == "Text"

    @respx.mock
    async def test_send_email_no_from_by_default(self, mcp_server):
        route = respx.post(f"{GRAPH_BASE_URL}/me/sendMail").mock(return_value=httpx.Response(202))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                await client.call_tool(
                    "send_email",
                    {"to": "alice@example.com", "subject": "Hi", "body": "Hello!"},
                )

        payload = json.loads(route.calls[0].request.content)
        assert "from" not in payload["message"]

    @respx.mock
    async def test_send_email_with_from_address(self, mcp_server):
        route = respx.post(f"{GRAPH_BASE_URL}/me/sendMail").mock(return_value=httpx.Response(202))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                await client.call_tool(
                    "send_email",
                    {
                        "to": "alice@example.com",
                        "subject": "Hi",
                        "body": "Hello!",
                        "options": '{"from_address": "mailbox@example.com"}',
                    },
                )

        payload = json.loads(route.calls[0].request.content)
        assert payload["message"]["from"]["emailAddress"]["address"] == "mailbox@example.com"

    @respx.mock
    async def test_graph_error_propagates_from_email_tool(self, mcp_server):
        from fastmcp.exceptions import ToolError

        respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages").mock(
            return_value=httpx.Response(
                401,
                json={"error": {"code": "InvalidAuthenticationToken", "message": "Token expired"}},
            )
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                with pytest.raises(ToolError, match="InvalidAuthenticationToken"):
                    await client.call_tool("list_emails", {})

    @respx.mock
    async def test_read_email_lists_attachments(self, mcp_server):
        """A message with attachments gets a section under the body."""
        respx.get(f"{GRAPH_BASE_URL}/me/messages/{ATT_MSG_ID}").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE_WITH_ATTACHMENTS)
        )
        respx.get(ATT_BASE).mock(return_value=httpx.Response(200, json=SAMPLE_ATTACHMENTS_RESPONSE))
        with _mock_token():
            result = await _call(mcp_server, "read_email", {"message_id": ATT_MSG_ID})

        text = _get_text(result)
        assert "Here is the weekly report" in text
        # The inline logo is counted, not listed.
        assert "**Attachments (2):**" in text
        assert "- report.pdf — application/pdf, 1.2 MB — id: `AAMkAttachFile001=`" in text
        assert (
            "- Q4 Plan.docx — link: https://contoso.sharepoint.com/:w:/s/team/Q4Plan"
            " — id: `AAMkAttachRef004=`" in text
        )
        assert "logo.png" not in text
        assert '(+1 inline image not shown; pass options {"include_inline": true})' in text

    @respx.mock
    async def test_read_email_include_inline_lists_them(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/messages/{ATT_MSG_ID}").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE_WITH_ATTACHMENTS)
        )
        respx.get(ATT_BASE).mock(return_value=httpx.Response(200, json=SAMPLE_ATTACHMENTS_RESPONSE))
        with _mock_token():
            result = await _call(
                mcp_server,
                "read_email",
                {"message_id": ATT_MSG_ID, "options": '{"include_inline": true}'},
            )

        text = _get_text(result)
        assert "**Attachments (3):**" in text
        assert "- logo.png — image/png, 4.0 KB (inline) — id: `AAMkAttachInline002=`" in text
        assert "not shown" not in text

    @respx.mock
    async def test_read_email_attachment_listing_failure_is_a_note(self, mcp_server):
        """A failed listing must not cost the caller the body they asked for."""
        respx.get(f"{GRAPH_BASE_URL}/me/messages/{ATT_MSG_ID}").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE_WITH_ATTACHMENTS)
        )
        respx.get(ATT_BASE).mock(return_value=httpx.Response(403, json=GRAPH_ERROR_403))
        with _mock_token():
            result = await _call(mcp_server, "read_email", {"message_id": ATT_MSG_ID})

        text = _get_text(result)
        assert "Here is the weekly report" in text
        assert "*(could not list attachments: Graph API error 403" in text
        assert "**Attachments" not in text

    @respx.mock
    async def test_read_email_without_attachments_has_no_section(self, mcp_server):
        """hasAttachments false means no attachment request at all."""
        respx.get(f"{GRAPH_BASE_URL}/me/messages/{ATT_MSG_ID}").mock(
            return_value=httpx.Response(200, json={**SAMPLE_MESSAGE, "hasAttachments": False})
        )
        list_route = respx.get(ATT_BASE).mock(
            return_value=httpx.Response(200, json=SAMPLE_ATTACHMENTS_RESPONSE)
        )
        with _mock_token():
            result = await _call(mcp_server, "read_email", {"message_id": ATT_MSG_ID})

        assert "**Attachments" not in _get_text(result)
        assert not list_route.called

    @respx.mock
    async def test_read_email_all_inline_still_prints_the_header(self, mcp_server):
        """Zero listed but attachments present — say so rather than staying silent."""
        respx.get(f"{GRAPH_BASE_URL}/me/messages/{ATT_MSG_ID}").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE_WITH_ATTACHMENTS)
        )
        respx.get(ATT_BASE).mock(
            return_value=httpx.Response(200, json={"value": [SAMPLE_INLINE_ATTACHMENT]})
        )
        with _mock_token():
            result = await _call(mcp_server, "read_email", {"message_id": ATT_MSG_ID})

        text = _get_text(result)
        assert "**Attachments (0):**" in text
        assert '(+1 inline image not shown; pass options {"include_inline": true})' in text

    @respx.mock
    async def test_send_email_with_text_and_base64_attachments(self, mcp_server):
        """Attachments switch the send to the draft path: create, attach, send."""
        send_mail = respx.post(f"{GRAPH_BASE_URL}/me/sendMail").mock(
            return_value=httpx.Response(202)
        )
        create = respx.post(f"{GRAPH_BASE_URL}/me/messages").mock(
            return_value=httpx.Response(201, json=SAMPLE_DRAFT_MESSAGE)
        )
        attach = respx.post(f"{DRAFT_BASE}/attachments").mock(
            return_value=httpx.Response(201, json=SAMPLE_CREATED_ATTACHMENT)
        )
        send = respx.post(f"{DRAFT_BASE}/send").mock(return_value=httpx.Response(202))

        specs = [
            {"name": "notes.txt", "text": "hello"},
            {"name": "img.png", "base64": base64.b64encode(b"\x89PNG").decode("ascii")},
        ]
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_email",
                {
                    "to": "alice@example.com",
                    "subject": "Hi",
                    "body": "Hello!",
                    "options": json.dumps({"attachments": specs}),
                },
            )

        assert create.called and send.called
        assert not send_mail.called
        assert attach.call_count == 2
        payloads = [json.loads(call.request.content) for call in attach.calls]
        assert [p["name"] for p in payloads] == ["notes.txt", "img.png"]
        assert all(p["@odata.type"] == "#microsoft.graph.fileAttachment" for p in payloads)
        assert base64.b64decode(payloads[0]["contentBytes"]) == b"hello"

        text = _get_text(result)
        assert "Email sent to alice@example.com with 2 attachment(s):" in text
        assert "notes.txt (5 B)" in text
        assert "img.png (4 B)" in text

    @respx.mock
    async def test_send_email_bad_attachment_spec_sends_nothing(self, mcp_server):
        send_mail = respx.post(f"{GRAPH_BASE_URL}/me/sendMail").mock(
            return_value=httpx.Response(202)
        )
        create = respx.post(f"{GRAPH_BASE_URL}/me/messages").mock(
            return_value=httpx.Response(201, json=SAMPLE_DRAFT_MESSAGE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_email",
                {
                    "to": "alice@example.com",
                    "subject": "Hi",
                    "body": "Hello!",
                    "options": '{"attachments": [{"name": "x"}]}',
                },
            )

        assert "attachments[0]:" in _get_text(result)
        assert not send_mail.called
        assert not create.called

    @respx.mock
    async def test_send_email_attachments_must_be_an_array(self, mcp_server):
        send_mail = respx.post(f"{GRAPH_BASE_URL}/me/sendMail").mock(
            return_value=httpx.Response(202)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_email",
                {
                    "to": "alice@example.com",
                    "subject": "Hi",
                    "body": "Hello!",
                    "options": '{"attachments": "nope"}',
                },
            )

        assert "attachments must be a JSON array" in _get_text(result)
        assert not send_mail.called

    @respx.mock
    async def test_send_email_empty_attachment_list_uses_send_mail(self, mcp_server):
        """An empty list is not "attachments" — keep the cheap one-request path."""
        send_mail = respx.post(f"{GRAPH_BASE_URL}/me/sendMail").mock(
            return_value=httpx.Response(202)
        )
        create = respx.post(f"{GRAPH_BASE_URL}/me/messages").mock(
            return_value=httpx.Response(201, json=SAMPLE_DRAFT_MESSAGE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_email",
                {
                    "to": "alice@example.com",
                    "subject": "Hi",
                    "body": "Hello!",
                    "options": '{"attachments": []}',
                },
            )

        assert send_mail.called
        assert not create.called
        assert _get_text(result) == "Email sent to alice@example.com."


class TestMCPGetEmailAttachment:
    """get_email_attachment — the markdown attachment reader."""

    @respx.mock
    async def test_text_mode_extracts_a_word_document(self, mcp_server):
        docx = _docx_bytes()
        meta = {
            **SAMPLE_FILE_ATTACHMENT,
            "name": "report.docx",
            "contentType": DOCX_MIME,
            "size": len(docx),
        }
        respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(200, content=docx, headers={"Content-Type": DOCX_MIME})
        )
        respx.get(ATT_FILE_URL).mock(return_value=httpx.Response(200, json=meta))
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_FILE_ATTACHMENT["id"]},
            )

        text = _get_text(result)
        assert "**Name:** report.docx" in text
        assert f"**ID:** `{SAMPLE_FILE_ATTACHMENT['id']}`" in text
        assert "Quarterly Title" in text
        assert "Body text here." in text

    @respx.mock
    async def test_text_mode_decodes_a_plain_text_attachment(self, mcp_server):
        meta = {
            **SAMPLE_FILE_ATTACHMENT,
            "name": "notes.txt",
            "contentType": "text/plain",
            "size": 11,
        }
        respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(
                200, content=b"hello there", headers={"Content-Type": "text/plain"}
            )
        )
        respx.get(ATT_FILE_URL).mock(return_value=httpx.Response(200, json=meta))
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_FILE_ATTACHMENT["id"]},
            )

        text = _get_text(result)
        assert "**Type:** text/plain (11 B)" in text
        assert text.endswith("---\nhello there")

    @respx.mock
    async def test_text_mode_on_a_binary_says_so(self, mcp_server):
        respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(
                200, content=b"\x89PNG\r\n\x1a\n", headers={"Content-Type": "image/png"}
            )
        )
        respx.get(ATT_FILE_URL).mock(
            return_value=httpx.Response(
                200,
                json={
                    **SAMPLE_FILE_ATTACHMENT,
                    "name": "logo.png",
                    "contentType": "image/png",
                    "size": 8,
                },
            )
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_FILE_ATTACHMENT["id"]},
            )

        text = _get_text(result)
        assert "binary file and cannot be displayed as text" in text
        assert 'Use mode "onedrive" to save it, or "base64" if it is under 1 MB.' in text

    @respx.mock
    async def test_base64_mode_returns_the_bytes(self, mcp_server):
        payload = b"\x89PNG\r\n\x1a\n"
        respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(200, content=payload, headers={"Content-Type": "image/png"})
        )
        respx.get(ATT_FILE_URL).mock(
            return_value=httpx.Response(
                200,
                json={
                    **SAMPLE_FILE_ATTACHMENT,
                    "name": "logo.png",
                    "contentType": "image/png",
                    "size": len(payload),
                },
            )
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_FILE_ATTACHMENT["id"],
                    "mode": "base64",
                },
            )

        text = _get_text(result)
        assert f"**Base64 ({len(payload)} bytes):**" in text
        assert base64.b64encode(payload).decode("ascii") in text

    @respx.mock
    async def test_base64_refuses_oversize_without_downloading(self, mcp_server):
        """The metadata size decides, so the bytes never cross the wire."""
        value_route = respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(200, content=b"x")
        )
        respx.get(ATT_FILE_URL).mock(
            return_value=httpx.Response(200, json={**SAMPLE_FILE_ATTACHMENT, "size": 2_000_000})
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_FILE_ATTACHMENT["id"],
                    "mode": "base64",
                },
            )

        assert "Too large to return as base64 (limit 976.6 KB)." in _get_text(result)
        assert not value_route.called

    @respx.mock
    async def test_text_mode_refuses_oversize_without_downloading(self, mcp_server):
        value_route = respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(200, content=b"x")
        )
        respx.get(ATT_FILE_URL).mock(
            return_value=httpx.Response(200, json={**SAMPLE_FILE_ATTACHMENT, "size": 60_000_000})
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_FILE_ATTACHMENT["id"]},
            )

        assert "too large for text extraction (limit: 50 MB)" in _get_text(result)
        assert not value_route.called

    @respx.mock
    async def test_onedrive_mode_uploads_and_returns_the_link(self, mcp_server):
        respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(
                200, content=b"%PDF-1.7", headers={"Content-Type": "application/pdf"}
            )
        )
        respx.get(ATT_FILE_URL).mock(
            return_value=httpx.Response(200, json={**SAMPLE_FILE_ATTACHMENT, "size": 8})
        )
        upload = respx.put(f"{GRAPH_BASE_URL}/me/drive/root:/Attachments/report.pdf:/content").mock(
            return_value=httpx.Response(201, json=SAMPLE_UPLOADED_FILE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_FILE_ATTACHMENT["id"],
                    "mode": "onedrive",
                },
            )

        assert upload.called
        assert upload.calls[0].request.content == b"%PDF-1.7"
        text = _get_text(result)
        assert f"**Saved to OneDrive:** {SAMPLE_UPLOADED_FILE['webUrl']}" in text
        assert f"**Item ID:** `{SAMPLE_UPLOADED_FILE['id']}`" in text

    @respx.mock
    async def test_onedrive_mode_honors_folder_path_option(self, mcp_server):
        respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(200, content=b"%PDF-1.7")
        )
        respx.get(ATT_FILE_URL).mock(
            return_value=httpx.Response(200, json={**SAMPLE_FILE_ATTACHMENT, "size": 8})
        )
        upload = respx.put(f"{GRAPH_BASE_URL}/me/drive/root:/Inbox/Files/report.pdf:/content").mock(
            return_value=httpx.Response(201, json=SAMPLE_UPLOADED_FILE)
        )
        with _mock_token():
            await _call(
                mcp_server,
                "get_email_attachment",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_FILE_ATTACHMENT["id"],
                    "mode": "onedrive",
                    "options": '{"folder_path": "Inbox/Files"}',
                },
            )

        assert upload.called

    @respx.mock
    async def test_item_attachment_renders_the_inner_message(self, mcp_server):
        inner = {
            "subject": "Budget draft",
            "from": {"emailAddress": {"name": "Dana Lee", "address": "dana@example.com"}},
            "receivedDateTime": "2025-12-01T09:00:00Z",
            "bodyPreview": "Numbers attached",
            "body": {"contentType": "text", "content": "Numbers attached, see inside."},
        }

        def _respond(request):
            # Graph query strings arrive percent-encoded: "$" is "%24".
            if "expand" in str(request.url):
                return httpx.Response(200, json={**SAMPLE_ITEM_ATTACHMENT, "item": inner})
            return httpx.Response(200, json=SAMPLE_ITEM_ATTACHMENT)

        respx.get(url__startswith=ATT_ITEM_URL).mock(side_effect=_respond)
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_ITEM_ATTACHMENT["id"]},
            )

        text = _get_text(result)
        assert "**Attached message**" in text
        assert "**Subject:** Budget draft" in text
        assert "**From:** Dana Lee <dana@example.com>" in text
        assert "**Date:** 2025-12-01T09:00:00Z" in text
        assert "Numbers attached, see inside." in text

    @respx.mock
    async def test_item_attachment_html_body_falls_back_to_the_preview(self, mcp_server):
        inner = {
            "subject": "Budget draft",
            "bodyPreview": "Numbers attached",
            "body": {"contentType": "html", "content": "<p>Numbers attached</p>"},
        }

        def _respond(request):
            # Graph query strings arrive percent-encoded: "$" is "%24".
            if "expand" in str(request.url):
                return httpx.Response(200, json={**SAMPLE_ITEM_ATTACHMENT, "item": inner})
            return httpx.Response(200, json=SAMPLE_ITEM_ATTACHMENT)

        respx.get(url__startswith=ATT_ITEM_URL).mock(side_effect=_respond)
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_ITEM_ATTACHMENT["id"]},
            )

        text = _get_text(result)
        assert "Numbers attached" in text
        assert "[HTML body, 23 chars — preview only]" in text

    @respx.mock
    async def test_item_attachment_downloads_as_eml_in_base64_mode(self, mcp_server):
        respx.get(f"{ATT_ITEM_URL}/$value").mock(
            return_value=httpx.Response(200, content=b"From: dana@example.com")
        )
        respx.get(ATT_ITEM_URL).mock(
            return_value=httpx.Response(200, json={**SAMPLE_ITEM_ATTACHMENT, "size": 22})
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_ITEM_ATTACHMENT["id"],
                    "mode": "base64",
                },
            )

        text = _get_text(result)
        assert "**Base64 (22 bytes):**" in text
        assert base64.b64encode(b"From: dana@example.com").decode("ascii") in text

    @respx.mock
    async def test_item_attachment_base64_refuses_oversize_without_downloading(self, mcp_server):
        """The base64 ceiling applies to attached messages too — before any download."""
        value_route = respx.get(f"{ATT_ITEM_URL}/$value").mock(
            return_value=httpx.Response(200, content=b"x")
        )
        respx.get(ATT_ITEM_URL).mock(
            return_value=httpx.Response(200, json={**SAMPLE_ITEM_ATTACHMENT, "size": 2_000_000})
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_ITEM_ATTACHMENT["id"],
                    "mode": "base64",
                },
            )

        assert "Too large to return as base64" in _get_text(result)
        assert not value_route.called

    @respx.mock
    async def test_reference_attachment_returns_its_link_in_every_mode(self, mcp_server):
        respx.get(ATT_REF_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_REFERENCE_ATTACHMENT)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_REFERENCE_ATTACHMENT["id"],
                    "mode": "onedrive",
                },
            )

        text = _get_text(result)
        assert f"**Link:** {SAMPLE_REFERENCE_ATTACHMENT['sourceUrl']}" in text
        assert "This is a link attachment" in text

    async def test_bad_mode_is_rejected_before_any_request(self, mcp_server):
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {"message_id": ATT_MSG_ID, "attachment_id": "a", "mode": "pdf"},
            )

        assert _get_text(result) == "mode must be one of: text, base64, onedrive; got 'pdf'"

    @respx.mock
    async def test_shared_mailbox_uses_the_users_path(self, mcp_server):
        base = (
            f"{GRAPH_BASE_URL}/users/support@example.com/messages/"
            f"{quote(ATT_MSG_ID, safe='')}/attachments/"
            f"{quote(SAMPLE_FILE_ATTACHMENT['id'], safe='')}"
        )
        respx.get(f"{base}/$value").mock(
            return_value=httpx.Response(
                200, content=b"hello", headers={"Content-Type": "text/plain"}
            )
        )
        meta_route = respx.get(base).mock(
            return_value=httpx.Response(
                200,
                json={
                    **SAMPLE_FILE_ATTACHMENT,
                    "name": "notes.txt",
                    "contentType": "text/plain",
                    "size": 5,
                },
            )
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_FILE_ATTACHMENT["id"],
                    "mailbox": "support@example.com",
                },
            )

        assert meta_route.called
        assert _get_text(result).endswith("---\nhello")


# ---------------------------------------------------------------------------
# Teams
# ---------------------------------------------------------------------------


class TestMCPTeamsTools:
    """Test consolidated Teams MCP tools."""

    @respx.mock
    async def test_list_teams_all(self, mcp_server):
        """No team_id → returns all joined teams."""
        respx.get(f"{GRAPH_BASE_URL}/me/joinedTeams").mock(
            return_value=httpx.Response(200, json=SAMPLE_TEAMS_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_teams", {})

        text = _get_text(result)
        assert "Engineering" in text
        assert "2 team(s)" in text

    @respx.mock
    async def test_list_teams_channels(self, mcp_server):
        """team_id set → returns channels for that team."""
        team_id = "team-id-001"
        respx.get(f"{GRAPH_BASE_URL}/teams/{team_id}/channels").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHANNELS_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_teams", {"team_id": team_id})

        text = _get_text(result)
        assert "General" in text
        assert "2 channel(s)" in text
        assert team_id in text

    @respx.mock
    async def test_list_teams_empty(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/joinedTeams").mock(
            return_value=httpx.Response(200, json={"value": []})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_teams", {})

        assert "No teams found" in _get_text(result)

    @respx.mock
    async def test_list_teams_not_available(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/joinedTeams").mock(
            return_value=httpx.Response(403, json=GRAPH_ERROR_403)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_teams", {})

        assert "not available" in _get_text(result).lower()

    @respx.mock
    async def test_list_chats(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/chats").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHATS_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_chats", {})

        text = _get_text(result)
        assert "3 chat(s)" in text
        assert "Alice Smith" in text
        assert "unread|type|name" in text
        assert "False|oneOnOne" in text
        assert "True|group" in text
        assert "True|meeting" in text

    @respx.mock
    async def test_list_chats_unread_edge_cases(self, mcp_server):
        """Verifies unread column: no messages = not unread, null read timestamp = unread."""
        chat_no_messages = {
            "id": "chat-empty-001",
            "chatType": "oneOnOne",
            "topic": None,
            "members": [{"displayName": "Alice"}],
            "lastMessagePreview": None,
            "viewpoint": None,
        }
        chat_null_read = {
            "id": "chat-null-read-001",
            "chatType": "group",
            "topic": "Test",
            "members": [{"displayName": "Bob"}],
            "lastMessagePreview": {
                "createdDateTime": "2025-12-15T10:00:00Z",
                "body": {"content": "hello"},
                "from": {"user": {"displayName": "Bob"}},
            },
            "viewpoint": {"lastMessageReadDateTime": None},
        }
        respx.get(f"{GRAPH_BASE_URL}/me/chats").mock(
            return_value=httpx.Response(200, json={"value": [chat_no_messages, chat_null_read]})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_chats", {})

        text = _get_text(result)
        assert "False|oneOnOne" in text
        assert "True|group" in text

    @respx.mock
    async def test_list_chats_empty(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/chats").mock(
            return_value=httpx.Response(200, json={"value": []})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_chats", {})

        assert "No chats found" in _get_text(result)

    async def test_list_chats_invalid_type(self, mcp_server):
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_chats", {"chat_type": "invalid"})

        assert "Invalid chat_type" in _get_text(result)

    @respx.mock
    async def test_read_teams_messages_channel(self, mcp_server):
        """team_id + channel_id → reads channel messages."""
        team_id = "team-id-001"
        channel_id = "channel-id-001"
        respx.get(f"{GRAPH_BASE_URL}/teams/{team_id}/channels/{channel_id}/messages").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHANNEL_MESSAGES_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "read_teams_messages",
                    {"team_id": team_id, "channel_id": channel_id, "since": "2025-01-01"},
                )

        text = _get_text(result)
        assert "2 message(s)" in text
        assert channel_id in text

    @respx.mock
    async def test_read_teams_messages_chat(self, mcp_server):
        """chat_id → reads chat messages."""
        chat_id = "chat-1on1-001"
        respx.get(f"{GRAPH_BASE_URL}/chats/{chat_id}/messages").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGES_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "read_teams_messages",
                    {"chat_id": chat_id, "since": "2025-01-01"},
                )

        text = _get_text(result)
        assert "1 message(s)" in text
        assert chat_id in text

    @respx.mock
    async def test_read_teams_messages_chat_takes_priority(self, mcp_server):
        """When both chat_id and team_id+channel_id are set, chat_id takes priority."""
        chat_id = "chat-1on1-001"
        respx.get(f"{GRAPH_BASE_URL}/chats/{chat_id}/messages").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGES_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "read_teams_messages",
                    {
                        "chat_id": chat_id,
                        "team_id": "team-id-001",
                        "channel_id": "channel-id-001",
                        "since": "2025-01-01",
                    },
                )

        text = _get_text(result)
        assert chat_id in text

    async def test_read_teams_messages_no_ids(self, mcp_server):
        """No IDs provided → returns helpful error."""
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("read_teams_messages", {})

        assert "Provide either chat_id" in _get_text(result)

    async def test_read_teams_messages_only_team_id(self, mcp_server):
        """Only team_id (no channel_id) → returns helpful error."""
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("read_teams_messages", {"team_id": "t1"})

        assert "Provide either chat_id" in _get_text(result)

    @respx.mock
    async def test_read_teams_messages_max_content_length(self, mcp_server):
        """max_content_length option truncates message bodies."""
        chat_id = "chat-1on1-001"
        long_msg = {
            "id": "msg-long-001",
            "messageType": "message",
            "createdDateTime": "2025-12-15T12:00:00Z",
            "from": {"user": {"displayName": "Tim"}, "application": None},
            "body": {"contentType": "text", "content": "SELECT " + "x" * 2000},
            "attachments": [],
        }
        respx.get(f"{GRAPH_BASE_URL}/chats/{chat_id}/messages").mock(
            return_value=httpx.Response(200, json={"value": [long_msg]})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result_full = await client.call_tool(
                    "read_teams_messages",
                    {"chat_id": chat_id, "since": "2025-01-01"},
                )
                result_truncated = await client.call_tool(
                    "read_teams_messages",
                    {
                        "chat_id": chat_id,
                        "since": "2025-01-01",
                        "options": '{"max_content_length": 50}',
                    },
                )

        full_text = _get_text(result_full)
        truncated_text = _get_text(result_truncated)
        assert "x" * 100 in full_text
        assert "..." in truncated_text
        assert "x" * 100 not in truncated_text

    @respx.mock
    async def test_read_email_max_content_length(self, mcp_server):
        """max_content_length option truncates email body."""
        long_email = {
            **SAMPLE_MESSAGE,
            "id": "msg-long-email",
            "body": {"contentType": "text", "content": "Report: " + "z" * 3000},
        }
        respx.get(f"{GRAPH_BASE_URL}/me/messages/{long_email['id']}").mock(
            return_value=httpx.Response(200, json=long_email)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result_full = await client.call_tool("read_email", {"message_id": long_email["id"]})
                result_truncated = await client.call_tool(
                    "read_email",
                    {
                        "message_id": long_email["id"],
                        "options": '{"max_content_length": 100}',
                    },
                )

        full_text = _get_text(result_full)
        truncated_text = _get_text(result_truncated)
        assert "z" * 200 in full_text
        assert "z" * 200 not in truncated_text

    @respx.mock
    async def test_send_teams_message_to_channel(self, mcp_server):
        """team_id + channel_id → sends to channel."""
        team_id = "team-id-001"
        channel_id = "channel-id-001"
        route = respx.post(f"{GRAPH_BASE_URL}/teams/{team_id}/channels/{channel_id}/messages").mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "send_teams_message",
                    {"message": "Hello!", "team_id": team_id, "channel_id": channel_id},
                )

        text = _get_text(result)
        assert "channel" in text.lower()
        assert route.called

    @respx.mock
    async def test_send_teams_message_to_chat(self, mcp_server):
        """chat_id → sends to chat."""
        chat_id = "chat-1on1-001"
        route = respx.post(f"{GRAPH_BASE_URL}/chats/{chat_id}/messages").mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "send_teams_message",
                    {"message": "Hello!", "chat_id": chat_id},
                )

        text = _get_text(result)
        assert "chat" in text.lower()
        assert route.called

    async def test_send_teams_message_no_ids(self, mcp_server):
        """No IDs → helpful error, no API call made."""
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("send_teams_message", {"message": "Hello!"})

        assert "Provide either chat_id" in _get_text(result)

    @respx.mock
    async def test_send_teams_message_newlines_converted(self, mcp_server):
        """Auto content_type converts newlines to <br> in plain text."""
        chat_id = "chat-1on1-001"
        route = respx.post(f"{GRAPH_BASE_URL}/chats/{chat_id}/messages").mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "send_teams_message",
                    {"message": "Hello\nWorld", "chat_id": chat_id},
                )

        text = _get_text(result)
        assert "chat" in text.lower()
        payload = json.loads(route.calls[0].request.content)
        assert payload["body"]["contentType"] == "html"
        assert payload["body"]["content"] == "Hello<br>World"

    @respx.mock
    async def test_send_teams_message_html_hyperlink(self, mcp_server):
        """HTML hyperlinks pass through in auto mode."""
        chat_id = "chat-1on1-001"
        html_msg = '<a href="https://example.com">Click here</a>'
        route = respx.post(f"{GRAPH_BASE_URL}/chats/{chat_id}/messages").mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "send_teams_message",
                    {"message": html_msg, "chat_id": chat_id},
                )

        payload = json.loads(route.calls[0].request.content)
        assert payload["body"]["contentType"] == "html"
        assert payload["body"]["content"] == html_msg

    @respx.mock
    async def test_send_teams_message_with_user_mention(self, mcp_server):
        """Mentions option builds Graph API mentions payload."""
        team_id = "team-id-001"
        channel_id = "channel-id-001"
        route = respx.post(f"{GRAPH_BASE_URL}/teams/{team_id}/channels/{channel_id}/messages").mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "send_teams_message",
                    {
                        "message": "Hey check this out",
                        "team_id": team_id,
                        "channel_id": channel_id,
                        "options": '{"mentions": [{"user_id": "aad-123", "name": "Alice"}]}',
                    },
                )

        text = _get_text(result)
        assert "channel" in text.lower()
        payload = json.loads(route.calls[0].request.content)
        assert "mentions" in payload
        assert len(payload["mentions"]) == 1
        assert payload["mentions"][0]["mentionText"] == "Alice"
        assert payload["mentions"][0]["mentioned"]["user"]["id"] == "aad-123"
        body_content = payload["body"]["content"]
        assert '<at id="0">Alice</at>' in body_content
        assert "Hey check this out" in body_content

    @respx.mock
    async def test_send_teams_message_mention_everyone(self, mcp_server):
        """mention_everyone builds channel-wide mention."""
        team_id = "team-id-001"
        channel_id = "channel-id-001"
        route = respx.post(f"{GRAPH_BASE_URL}/teams/{team_id}/channels/{channel_id}/messages").mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "send_teams_message",
                    {
                        "message": "Important update",
                        "team_id": team_id,
                        "channel_id": channel_id,
                        "options": '{"mention_everyone": true}',
                    },
                )

        text = _get_text(result)
        assert "channel" in text.lower()
        payload = json.loads(route.calls[0].request.content)
        assert "mentions" in payload
        assert payload["mentions"][0]["mentionText"] == "Everyone"
        assert (
            payload["mentions"][0]["mentioned"]["conversation"]["conversationIdentityType"]
            == "channel"
        )
        body_content = payload["body"]["content"]
        assert '<at id="0">Everyone</at>' in body_content
        assert "Important update" in body_content

    @respx.mock
    async def test_send_teams_message_content_type_in_options(self, mcp_server):
        """content_type in options works like the old positional param."""
        chat_id = "chat-1on1-001"
        route = respx.post(f"{GRAPH_BASE_URL}/chats/{chat_id}/messages").mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                await client.call_tool(
                    "send_teams_message",
                    {
                        "message": "plain text only",
                        "chat_id": chat_id,
                        "options": '{"content_type": "text"}',
                    },
                )

        payload = json.loads(route.calls[0].request.content)
        assert payload["body"]["contentType"] == "text"

    @respx.mock
    async def test_get_teams_activity(self, mcp_server):
        # Wire up all the calls the activity scanner makes
        respx.get(f"{GRAPH_BASE_URL}/me/joinedTeams").mock(
            return_value=httpx.Response(200, json={"value": []})
        )
        respx.get(f"{GRAPH_BASE_URL}/me/chats").mock(
            return_value=httpx.Response(200, json={"value": []})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("get_teams_activity", {"hours": 1})

        text = _get_text(result)
        assert "No Teams activity" in text

    @respx.mock
    async def test_get_teams_activity_not_available(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/joinedTeams").mock(
            return_value=httpx.Response(403, json=GRAPH_ERROR_403)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("get_teams_activity", {})

        assert "not available" in _get_text(result).lower()

    @respx.mock
    async def test_list_chats_with_mark_as_read(self, mcp_server):
        """mark_as_read option triggers POST markChatReadForUser for each ID."""
        import base64

        payload_data = {"oid": "user-obj-id", "tid": "tenant-id-123"}
        header = base64.urlsafe_b64encode(b'{"alg":"RS256"}').rstrip(b"=").decode()
        payload = base64.urlsafe_b64encode(json.dumps(payload_data).encode()).rstrip(b"=").decode()
        sig = base64.urlsafe_b64encode(b"s").rstrip(b"=").decode()
        fake_token = f"{header}.{payload}.{sig}"

        respx.get(f"{GRAPH_BASE_URL}/me/chats").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHATS_RESPONSE)
        )
        mark_route = respx.post(f"{GRAPH_BASE_URL}/chats/chat-1on1-001/markChatReadForUser").mock(
            return_value=httpx.Response(204)
        )

        with _mock_token(fake_token):
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "list_chats",
                    {"options": '{"mark_as_read": ["chat-1on1-001"]}'},
                )

        text = _get_text(result)
        assert "3 chat(s)" in text
        assert mark_route.called
        sent_payload = json.loads(mark_route.calls[0].request.content)
        assert sent_payload == {"user": {"id": "user-obj-id", "tenantId": "tenant-id-123"}}
        assert "1 chat(s) marked as read" in text

    @respx.mock
    async def test_list_chats_mark_as_read_rejects_non_array(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/chats").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHATS_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "list_chats",
                    {"options": '{"mark_as_read": true}'},
                )

        assert "must be a JSON array" in _get_text(result)

    @respx.mock
    async def test_read_teams_messages_with_mark_as_read(self, mcp_server):
        """mark_as_read: true marks the chat as read after reading messages."""
        import base64

        payload_data = {"oid": "user-obj-id", "tid": "tenant-id-123"}
        header = base64.urlsafe_b64encode(b'{"alg":"RS256"}').rstrip(b"=").decode()
        payload = base64.urlsafe_b64encode(json.dumps(payload_data).encode()).rstrip(b"=").decode()
        sig = base64.urlsafe_b64encode(b"s").rstrip(b"=").decode()
        fake_token = f"{header}.{payload}.{sig}"

        chat_id = "chat-1on1-001"
        respx.get(f"{GRAPH_BASE_URL}/chats/{chat_id}/messages").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGES_RESPONSE)
        )
        mark_route = respx.post(f"{GRAPH_BASE_URL}/chats/{chat_id}/markChatReadForUser").mock(
            return_value=httpx.Response(204)
        )

        with _mock_token(fake_token):
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "read_teams_messages",
                    {
                        "chat_id": chat_id,
                        "since": "2025-01-01",
                        "options": '{"mark_as_read": true}',
                    },
                )

        text = _get_text(result)
        assert "1 message(s)" in text
        assert mark_route.called
        sent_payload = json.loads(mark_route.calls[0].request.content)
        assert sent_payload == {"user": {"id": "user-obj-id", "tenantId": "tenant-id-123"}}
        assert "marked as read" in text.lower()

    @respx.mock
    async def test_read_teams_messages_mark_as_read_channel_rejected(self, mcp_server):
        """mark_as_read with team_id+channel_id returns error."""
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "read_teams_messages",
                    {
                        "team_id": "team-id-001",
                        "channel_id": "channel-id-001",
                        "since": "2025-01-01",
                        "options": '{"mark_as_read": true}',
                    },
                )

        assert "only supported for chats" in _get_text(result)


class TestMCPSendTeamsMessageFiles:
    """send_teams_message with attachments and images in its options."""

    @respx.mock
    async def test_a_chat_file_is_uploaded_shared_then_posted(self, mcp_server):
        _mock_chat_file_upload()
        post = respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_teams_message",
                {
                    "message": "here you go",
                    "chat_id": TEAMS_CHAT_ID,
                    "options": '{"attachments": [{"name": "notes.txt", "text": "hello"}]}',
                },
            )

        assert _graph_trail() == [
            ("PUT", "/v1.0/me/drive/root:/Microsoft Teams Chat Files/notes.txt:/content"),
            ("GET", "/v1.0/me/drive/items/teams-upload-001"),
            ("GET", f"/v1.0/chats/{TEAMS_CHAT_ID}/members"),
            ("POST", "/v1.0/me/drive/items/teams-upload-001/invite"),
            ("POST", f"/v1.0/chats/{TEAMS_CHAT_ID}/messages"),
        ]
        payload = json.loads(post.calls[0].request.content)
        assert payload["attachments"] == [
            {
                "id": TEAMS_UPLOAD_GUID,
                "contentType": "reference",
                "contentUrl": TEAMS_WEBDAV_URL,
                "name": "notes.txt",
            }
        ]
        assert f'<attachment id="{TEAMS_UPLOAD_GUID}"></attachment>' in payload["body"]["content"]
        assert _get_text(result) == ("Message sent to Teams chat with 1 file(s): notes.txt (5 B).")

    @respx.mock
    async def test_an_image_rides_inside_the_message_with_no_upload(self, mcp_server):
        post = respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        encoded = base64.b64encode(PNG_BYTES).decode()
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_teams_message",
                {
                    "message": "look",
                    "chat_id": TEAMS_CHAT_ID,
                    "options": json.dumps({"images": [{"name": "pic.png", "base64": encoded}]}),
                },
            )

        assert [method for method, _ in _graph_trail()] == ["POST"]
        payload = json.loads(post.calls[0].request.content)
        assert payload["hostedContents"][0]["@microsoft.graph.temporaryId"] == "1"
        assert payload["hostedContents"][0]["contentType"] == "image/png"
        assert base64.b64decode(payload["hostedContents"][0]["contentBytes"]) == PNG_BYTES
        assert '<img src="../hostedContents/1/$value">' in payload["body"]["content"]
        assert _get_text(result) == "Message sent to Teams chat with 1 inline image(s)."

    @respx.mock
    async def test_the_sender_from_the_token_is_not_invited(self, mcp_server):
        """The token's oid identifies the sender; they own the file already."""
        _mock_chat_file_upload()
        respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        with _mock_token(_token_with_oid("user-id-001")):
            await _call(
                mcp_server,
                "send_teams_message",
                {
                    "message": "here you go",
                    "chat_id": TEAMS_CHAT_ID,
                    "options": '{"attachments": [{"name": "notes.txt", "text": "hello"}]}',
                },
            )

        assert _invite_recipients() == [{"email": "alice@example.com"}]

    @respx.mock
    async def test_a_failed_post_removes_the_uploaded_file(self, mcp_server):
        _mock_chat_file_upload()
        respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(503, json={"error": {"code": "ServiceUnavailable"}})
        )
        removed = respx.delete(f"{GRAPH_BASE_URL}/drives/drive-001/items/teams-upload-001").mock(
            return_value=httpx.Response(204)
        )
        from fastmcp.exceptions import ToolError

        with _mock_token():
            with pytest.raises(ToolError, match="503"):
                await _call(
                    mcp_server,
                    "send_teams_message",
                    {
                        "message": "here you go",
                        "chat_id": TEAMS_CHAT_ID,
                        "options": '{"attachments": [{"name": "notes.txt", "text": "hello"}]}',
                    },
                )

        assert removed.call_count == 1

    @respx.mock
    async def test_files_and_images_are_both_named_in_the_confirmation(self, mcp_server):
        _mock_chat_file_upload()
        respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        encoded = base64.b64encode(PNG_BYTES).decode()
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_teams_message",
                {
                    "message": "both",
                    "chat_id": TEAMS_CHAT_ID,
                    "options": json.dumps(
                        {
                            "attachments": [{"name": "notes.txt", "text": "hello"}],
                            "images": [{"name": "pic.png", "base64": encoded}],
                        }
                    ),
                },
            )

        assert _get_text(result) == (
            "Message sent to Teams chat with 1 file(s): notes.txt (5 B) and 1 inline image(s)."
        )

    @respx.mock
    async def test_a_non_image_in_images_is_refused_before_anything_is_sent(self, mcp_server):
        post = respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_teams_message",
                {
                    "message": "look",
                    "chat_id": TEAMS_CHAT_ID,
                    "options": '{"images": [{"name": "notes.txt", "text": "hi"}]}',
                },
            )

        assert "not an image" in _get_text(result)
        assert not post.called

    @respx.mock
    async def test_a_bad_image_spec_is_reported_against_images_not_attachments(self, mcp_server):
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_teams_message",
                {
                    "message": "look",
                    "chat_id": TEAMS_CHAT_ID,
                    "options": '{"images": [{"text": "no name"}]}',
                },
            )

        text = _get_text(result)
        assert text.startswith("images[0]:")
        assert "attachments[" not in text

    @respx.mock
    async def test_a_bad_attachment_spec_names_its_index(self, mcp_server):
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_teams_message",
                {
                    "message": "hi",
                    "chat_id": TEAMS_CHAT_ID,
                    "options": '{"attachments": [{"name": "a.txt", "text": "ok"}, {}]}',
                },
            )

        assert _get_text(result).startswith("attachments[1]:")

    @respx.mock
    async def test_a_403_on_the_upload_explains_the_missing_permission(self, mcp_server):
        respx.put(url__startswith=TEAMS_SEND_UPLOAD_URL).mock(
            return_value=httpx.Response(403, json=GRAPH_ERROR_403)
        )
        post = respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_teams_message",
                {
                    "message": "here",
                    "chat_id": TEAMS_CHAT_ID,
                    "options": '{"attachments": [{"name": "notes.txt", "text": "hello"}]}',
                },
            )

        text = _get_text(result)
        assert "Files permission missing" in text
        assert "Files.ReadWrite" in text
        assert not post.called

    @respx.mock
    async def test_a_channel_file_lands_in_the_channel_drive(self, mcp_server):
        team_id = "team-id-001"
        channel_id = "channel-id-001"
        respx.get(f"{GRAPH_BASE_URL}/teams/{team_id}/channels/{channel_id}/filesFolder").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHANNEL_FILES_FOLDER)
        )
        respx.put(
            url__startswith=(
                f"{GRAPH_BASE_URL}/drives/drive-team-001/items/"
                "folder-channel-001:/notes.txt:/content"
            )
        ).mock(return_value=httpx.Response(201, json=SAMPLE_TEAMS_UPLOAD_RESPONSE))
        respx.get(
            url__startswith=f"{GRAPH_BASE_URL}/drives/drive-team-001/items/teams-upload-001"
        ).mock(return_value=httpx.Response(200, json=SAMPLE_TEAMS_UPLOADED_ITEM))
        post = respx.post(f"{GRAPH_BASE_URL}/teams/{team_id}/channels/{channel_id}/messages").mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )

        with _mock_token():
            result = await _call(
                mcp_server,
                "send_teams_message",
                {
                    "message": "deck",
                    "team_id": team_id,
                    "channel_id": channel_id,
                    "options": '{"attachments": [{"name": "notes.txt", "text": "hello"}]}',
                },
            )

        assert _graph_trail() == [
            ("GET", f"/v1.0/teams/{team_id}/channels/{channel_id}/filesFolder"),
            ("PUT", "/v1.0/drives/drive-team-001/items/folder-channel-001:/notes.txt:/content"),
            ("GET", "/v1.0/drives/drive-team-001/items/teams-upload-001"),
            ("POST", f"/v1.0/teams/{team_id}/channels/{channel_id}/messages"),
        ]
        assert json.loads(post.calls[0].request.content)["attachments"][0]["name"] == "notes.txt"
        assert _get_text(result).startswith("Message sent to Teams channel with 1 file(s)")

    @respx.mock
    async def test_an_empty_attachments_list_still_takes_the_file_path(self, mcp_server):
        """An explicit [] means "no files", not "old code path" — and must still send."""
        post = respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_teams_message",
                {
                    "message": "plain",
                    "chat_id": TEAMS_CHAT_ID,
                    "options": '{"attachments": []}',
                },
            )

        assert json.loads(post.calls[0].request.content) == {
            "body": {"contentType": "html", "content": "plain"}
        }
        assert _get_text(result) == "Message sent to Teams chat."


class TestMCPReadTeamsMessagesAttachments:
    """The attachments column, and the markers inside the content column."""

    @respx.mock
    async def test_attachments_column(self, mcp_server):
        page = {
            "value": [
                *SAMPLE_CHAT_MESSAGES_PAGE_WITH_ATTACHMENTS["value"][:3],
                SAMPLE_CHAT_MESSAGE_FULL,
            ]
        }
        respx.get(TEAMS_CHAT_MSGS).mock(return_value=httpx.Response(200, json=page))
        with _mock_token():
            result = await _call(
                mcp_server,
                "read_teams_messages",
                {"chat_id": TEAMS_CHAT_ID, "since": "2025-01-01"},
            )

        lines = _get_text(result).splitlines()
        assert lines[1] == "timestamp|sender|content|attachments|id"
        rows = {line.split("|")[-1]: line.split("|") for line in lines[2:] if line}

        assert rows["chat-msg-file-001"][3] == f"roadmap.pptx [file:{TEAMS_FILE_ATTACHMENT_ID}]"
        assert "[File: roadmap.pptx]" in rows["chat-msg-file-001"][2]
        assert rows["chat-msg-image-001"][3] == f"[image:{TEAMS_HOSTED_ID}]"
        assert rows["chat-msg-card-001"][3] == "[card]"
        # A message with nothing attached leaves the column empty.
        assert rows["chat-msg-001"][2:4] == ["Sounds good!", ""]


class TestMCPGetTeamsAttachment:
    """get_teams_attachment: files, inline images, cards, and every refusal."""

    async def test_invalid_mode(self, mcp_server):
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    "chat_id": TEAMS_CHAT_ID,
                    "mode": "pdf",
                },
            )

        assert _get_text(result) == "mode must be one of: text, base64, onedrive; got 'pdf'"

    async def test_missing_ids(self, mcp_server):
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {"message_id": "m1", "attachment_id": "a1"},
            )

        assert _get_text(result) == "Provide either chat_id, or both team_id and channel_id."

    @respx.mock
    async def test_unknown_id_lists_what_the_message_has(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-file-001",
                    "attachment_id": "nope",
                    "chat_id": TEAMS_CHAT_ID,
                },
            )

        text = _get_text(result)
        assert "No attachment with id `nope`" in text
        assert f"Available: file: {TEAMS_FILE_ATTACHMENT_ID}" in text

    @respx.mock
    async def test_file_text_mode_extracts_the_document(self, mcp_server):
        docx_item = {
            **SAMPLE_TEAMS_DRIVE_ITEM,
            "name": "notes.docx",
            "file": {"mimeType": DOCX_MIME},
        }
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        respx.get(TEAMS_SHARE_CONTENT_URL).mock(
            return_value=httpx.Response(200, content=_docx_bytes())
        )
        respx.get(TEAMS_SHARE_BASE).mock(return_value=httpx.Response(200, json=docx_item))
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    "chat_id": TEAMS_CHAT_ID,
                },
            )

        text = _get_text(result)
        assert "Quarterly Title" in text
        assert "**Name:** roadmap.pptx" in text

    @respx.mock
    async def test_file_base64_mode(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        respx.get(TEAMS_SHARE_CONTENT_URL).mock(
            return_value=httpx.Response(200, content=b"PPTXBYTES")
        )
        respx.get(TEAMS_SHARE_BASE).mock(
            return_value=httpx.Response(200, json=SAMPLE_TEAMS_DRIVE_ITEM)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    "chat_id": TEAMS_CHAT_ID,
                    "mode": "base64",
                },
            )

        text = _get_text(result)
        assert "**Base64 (" in text
        assert base64.b64encode(b"PPTXBYTES").decode() in text

    @respx.mock
    async def test_file_onedrive_mode_saves_a_copy(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        respx.get(TEAMS_SHARE_CONTENT_URL).mock(
            return_value=httpx.Response(200, content=b"PPTXBYTES")
        )
        respx.get(TEAMS_SHARE_BASE).mock(
            return_value=httpx.Response(200, json=SAMPLE_TEAMS_DRIVE_ITEM)
        )
        upload = respx.put(
            f"{GRAPH_BASE_URL}/me/drive/root:/Attachments/roadmap.pptx:/content"
        ).mock(return_value=httpx.Response(201, json=SAMPLE_UPLOADED_FILE))
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    "chat_id": TEAMS_CHAT_ID,
                    "mode": "onedrive",
                },
            )

        assert upload.called
        assert "**Saved to OneDrive:**" in _get_text(result)

    @respx.mock
    async def test_text_mode_refuses_an_oversized_file_before_downloading(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        content = respx.get(TEAMS_SHARE_CONTENT_URL).mock(
            return_value=httpx.Response(200, content=b"never")
        )
        respx.get(TEAMS_SHARE_BASE).mock(
            return_value=httpx.Response(200, json={**SAMPLE_TEAMS_DRIVE_ITEM, "size": 60_000_000})
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    "chat_id": TEAMS_CHAT_ID,
                },
            )

        assert "too large for text extraction (limit: 50 MB)" in _get_text(result)
        assert not content.called

    @respx.mock
    async def test_base64_mode_refuses_an_oversized_file_before_downloading(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        content = respx.get(TEAMS_SHARE_CONTENT_URL).mock(
            return_value=httpx.Response(200, content=b"never")
        )
        respx.get(TEAMS_SHARE_BASE).mock(
            return_value=httpx.Response(200, json={**SAMPLE_TEAMS_DRIVE_ITEM, "size": 2_000_000})
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    "chat_id": TEAMS_CHAT_ID,
                    "mode": "base64",
                },
            )

        assert "Too large to return as base64" in _get_text(result)
        assert not content.called

    @respx.mock
    async def test_403_on_the_sharing_link_is_access_denied(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        respx.get(TEAMS_SHARE_BASE).mock(return_value=httpx.Response(403, json=GRAPH_ERROR_403))
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    "chat_id": TEAMS_CHAT_ID,
                },
            )

        assert "**Access denied:**" in _get_text(result)

    @respx.mock
    async def test_404_on_the_sharing_link_is_item_not_found(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        respx.get(TEAMS_SHARE_BASE).mock(return_value=httpx.Response(404, json=GRAPH_ERROR_404))
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    "chat_id": TEAMS_CHAT_ID,
                },
            )

        assert "**Item not found**" in _get_text(result)

    @respx.mock
    async def test_a_shared_folder_is_refused(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        respx.get(TEAMS_SHARE_BASE).mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_ITEM_FOLDER)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    "chat_id": TEAMS_CHAT_ID,
                },
            )

        assert "is a folder, not a file" in _get_text(result)

    @respx.mock
    async def test_inline_image_base64_names_itself_from_the_id(self, mcp_server):
        respx.get(TEAMS_IMAGE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_IMAGE)
        )
        respx.get(TEAMS_HOSTED_VALUE_URL).mock(
            return_value=httpx.Response(
                200, content=PNG_BYTES, headers={"Content-Type": "image/png"}
            )
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-image-001",
                    "attachment_id": TEAMS_HOSTED_ID,
                    "chat_id": TEAMS_CHAT_ID,
                    "mode": "base64",
                },
            )

        text = _get_text(result)
        assert "**Name:** image-aWQ9eF8wLWN1.png" in text
        assert "image/png" in text
        assert base64.b64encode(PNG_BYTES).decode() in text

    @respx.mock
    async def test_inline_image_text_mode_reports_a_binary(self, mcp_server):
        respx.get(TEAMS_IMAGE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_IMAGE)
        )
        respx.get(TEAMS_HOSTED_VALUE_URL).mock(
            return_value=httpx.Response(
                200, content=PNG_BYTES, headers={"Content-Type": "image/png"}
            )
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-image-001",
                    "attachment_id": TEAMS_HOSTED_ID,
                    "chat_id": TEAMS_CHAT_ID,
                },
            )

        assert "binary file" in _get_text(result)

    @respx.mock
    async def test_card_renders_its_text_without_downloading(self, mcp_server):
        respx.get(TEAMS_CARD_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_CARD)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-card-001",
                    "attachment_id": "card-att-001",
                    "chat_id": TEAMS_CHAT_ID,
                },
            )

        assert "Deploy finished" in _get_text(result)

    @respx.mock
    async def test_quoted_message_reference_has_nothing_to_download(self, mcp_server):
        respx.get(TEAMS_JUNK_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_JUNK_ATTACHMENTS)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-junk-001",
                    "attachment_id": "ref-001",
                    "chat_id": TEAMS_CHAT_ID,
                },
            )

        assert "quoted message reference" in _get_text(result)

    @respx.mock
    async def test_unknown_kind_cannot_be_fetched(self, mcp_server):
        respx.get(TEAMS_JUNK_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_JUNK_ATTACHMENTS)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-junk-001",
                    "attachment_id": "img-att",
                    "chat_id": TEAMS_CHAT_ID,
                },
            )

        text = _get_text(result)
        assert "Attachment of type image/png cannot be fetched" in text
        assert "URL: https://x/y.png" in text

    @respx.mock
    async def test_file_without_a_content_url_is_explained(self, mcp_server):
        msg = {
            **SAMPLE_CHAT_MESSAGE_WITH_FILE,
            "attachments": [{"id": TEAMS_FILE_ATTACHMENT_ID, "contentType": "reference"}],
        }
        respx.get(TEAMS_FILE_MSG_URL).mock(return_value=httpx.Response(200, json=msg))
        share = respx.get(url__startswith=f"{GRAPH_BASE_URL}/shares/").mock(
            return_value=httpx.Response(200, json=SAMPLE_TEAMS_DRIVE_ITEM)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    "chat_id": TEAMS_CHAT_ID,
                },
            )

        assert "no content URL" in _get_text(result)
        assert not share.called

    @respx.mock
    async def test_403_on_the_message_reports_teams_unavailable(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(return_value=httpx.Response(403, json=GRAPH_ERROR_403))
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    "chat_id": TEAMS_CHAT_ID,
                },
            )

        assert _get_text(result) == "Microsoft Teams is not available for this account."

    @respx.mock
    async def test_channel_form_reads_the_channel_message(self, mcp_server):
        route = respx.get(f"{GRAPH_BASE_URL}/teams/t1/channels/c1/messages/chat-msg-card-001").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_CARD)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-card-001",
                    "attachment_id": "card-att-001",
                    "team_id": "t1",
                    "channel_id": "c1",
                },
            )

        assert route.called
        assert "Deploy finished" in _get_text(result)

    @respx.mock
    async def test_chat_id_wins_over_a_team_and_channel(self, mcp_server):
        chat_route = respx.get(TEAMS_CARD_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_CARD)
        )
        channel_route = respx.get(url__startswith=f"{GRAPH_BASE_URL}/teams/").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_CARD)
        )
        with _mock_token():
            await _call(
                mcp_server,
                "get_teams_attachment",
                {
                    "message_id": "chat-msg-card-001",
                    "attachment_id": "card-att-001",
                    "chat_id": TEAMS_CHAT_ID,
                    "team_id": "t1",
                    "channel_id": "c1",
                },
            )

        assert chat_route.called
        assert not channel_route.called


# ---------------------------------------------------------------------------
# Files
# ---------------------------------------------------------------------------


class TestMCPFileTools:
    """Test consolidated file MCP tools."""

    @respx.mock
    async def test_list_sharepoint_sites(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/sites").mock(
            return_value=httpx.Response(200, json=SAMPLE_SITES_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_sharepoint_sites", {"query": "engineering"})

        text = _get_text(result)
        assert "Engineering Hub" in text
        assert "2" in text

    @respx.mock
    async def test_list_sharepoint_sites_empty(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/sites").mock(
            return_value=httpx.Response(200, json={"value": []})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_sharepoint_sites", {"query": "nope"})

        assert "No SharePoint sites found" in _get_text(result)

    @respx.mock
    async def test_list_files_onedrive_browse(self, mcp_server):
        """No site_id, no query → OneDrive root browse."""
        respx.get(f"{GRAPH_BASE_URL}/me/drive/root/children").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_CHILDREN_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_files", {})

        text = _get_text(result)
        assert "Documents" in text
        assert "report.csv" in text

    @respx.mock
    async def test_list_files_onedrive_subfolder(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/drive/root:/Documents:/children").mock(
            return_value=httpx.Response(200, json={"value": [SAMPLE_DRIVE_ITEM_FILE]})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_files", {"folder_path": "Documents"})

        text = _get_text(result)
        assert "report.csv" in text
        assert "Documents" in text

    @respx.mock
    async def test_list_files_sharepoint_browse(self, mcp_server):
        """site_id set, no query → SharePoint browse."""
        site_id = "site-id-001"
        respx.get(f"{GRAPH_BASE_URL}/sites/{site_id}/drive/root/children").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_CHILDREN_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_files", {"site_id": site_id})

        text = _get_text(result)
        assert "Documents" in text

    @respx.mock
    async def test_list_files_search_query(self, mcp_server):
        """query set → search mode, site_id and folder_path ignored."""
        respx.post(f"{GRAPH_BASE_URL}/search/query").mock(
            return_value=httpx.Response(200, json=SAMPLE_SEARCH_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_files", {"query": "budget"})

        text = _get_text(result)
        assert "Q4-budget.xlsx" in text
        assert "budget-notes.md" in text

    @respx.mock
    async def test_list_files_search_no_results(self, mcp_server):
        respx.post(f"{GRAPH_BASE_URL}/search/query").mock(
            return_value=httpx.Response(200, json=SAMPLE_SEARCH_RESPONSE_EMPTY)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_files", {"query": "nonexistent"})

        assert "No files found" in _get_text(result)

    @respx.mock
    async def test_list_files_empty_folder(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/drive/root/children").mock(
            return_value=httpx.Response(200, json={"value": []})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_files", {})

        assert "No files found" in _get_text(result)

    @respx.mock
    async def test_inspect_file_metadata_only_default(self, mcp_server):
        """Default (no read_content arg) → metadata only, no download."""
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/{SAMPLE_DRIVE_ITEM_FILE['id']}").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_ITEM_FILE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "inspect_file",
                    {"item_id": SAMPLE_DRIVE_ITEM_FILE["id"]},
                )

        text = _get_text(result)
        assert "report.csv" in text
        assert "Alice Smith" in text
        assert "---" not in text

    @respx.mock
    async def test_inspect_file_metadata_only_explicit(self, mcp_server):
        """read_content=False → metadata only, no download."""
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/{SAMPLE_DRIVE_ITEM_FILE['id']}").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_ITEM_FILE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "inspect_file",
                    {"item_id": SAMPLE_DRIVE_ITEM_FILE["id"], "read_content": False},
                )

        text = _get_text(result)
        assert "report.csv" in text
        assert "---" not in text

    @respx.mock
    async def test_inspect_file_with_content(self, mcp_server):
        """read_content=True → downloads text content."""
        item_id = SAMPLE_DRIVE_ITEM_FILE["id"]
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/{item_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_ITEM_FILE)
        )
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/{item_id}/content").mock(
            return_value=httpx.Response(200, content=b"col1,col2\n1,2\n")
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "inspect_file",
                    {"item_id": item_id, "read_content": True},
                )

        text = _get_text(result)
        assert "report.csv" in text
        assert "col1,col2" in text
        assert "---" in text

    @respx.mock
    async def test_inspect_file_binary_returns_message(self, mcp_server):
        """Non-extractable binary files (images) report they cannot be shown as text."""
        image_item = {
            "id": "file-id-img-001",
            "name": "diagram.png",
            "size": 500_000,
            "file": {"mimeType": "image/png"},
            "lastModifiedDateTime": "2025-12-15T10:00:00Z",
            "lastModifiedBy": {"user": {"displayName": "Alice Smith", "id": "user-001"}},
            "webUrl": "https://onedrive.live.com/edit.aspx?resid=file-id-img-001",
            "parentReference": {"driveId": "drive-001", "path": "/drive/root:"},
        }
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/file-id-img-001").mock(
            return_value=httpx.Response(200, json=image_item)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "inspect_file", {"item_id": "file-id-img-001", "read_content": True}
                )

        text = _get_text(result)
        assert "diagram.png" in text
        assert "binary" in text.lower()

    @respx.mock
    async def test_inspect_file_extracts_pptx(self, mcp_server):
        """PPTX files are now extracted via document extraction."""
        import io

        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide Title"
        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        item_id = SAMPLE_DRIVE_ITEM_BINARY["id"]
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/{item_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_ITEM_BINARY)
        )
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/{item_id}/content").mock(
            return_value=httpx.Response(200, content=pptx_bytes)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "inspect_file", {"item_id": item_id, "read_content": True}
                )

        text = _get_text(result)
        assert "presentation.pptx" in text
        assert "Test Slide Title" in text

    @respx.mock
    async def test_inspect_file_too_large(self, mcp_server):
        item_id = SAMPLE_DRIVE_ITEM_LARGE_TEXT["id"]
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/{item_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_ITEM_LARGE_TEXT)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "inspect_file", {"item_id": item_id, "read_content": True}
                )

        text = _get_text(result)
        assert "too large" in text.lower()

    @respx.mock
    async def test_inspect_file_sharepoint(self, mcp_server):
        site_id = "site-id-001"
        item_id = SAMPLE_DRIVE_ITEM_FILE["id"]
        respx.get(f"{GRAPH_BASE_URL}/sites/{site_id}/drive/items/{item_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_ITEM_FILE)
        )
        respx.get(f"{GRAPH_BASE_URL}/sites/{site_id}/drive/items/{item_id}/content").mock(
            return_value=httpx.Response(200, content=b"data")
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "inspect_file",
                    {"item_id": item_id, "site_id": site_id, "read_content": True},
                )

        assert "report.csv" in _get_text(result)

    @respx.mock
    async def test_graph_error_propagates_from_file_tool(self, mcp_server):
        from fastmcp.exceptions import ToolError

        respx.get(f"{GRAPH_BASE_URL}/me/drive/root/children").mock(
            return_value=httpx.Response(
                401,
                json={"error": {"code": "InvalidAuthenticationToken", "message": "Token expired"}},
            )
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                with pytest.raises(ToolError, match="InvalidAuthenticationToken"):
                    await client.call_tool("list_files", {})


class TestMCPInspectFileJson:
    """inspect_file_json: the structured sibling of inspect_file."""

    @respx.mock
    async def test_by_item_id_returns_metadata_only(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/file-id-001").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_ITEM_FILE)
        )
        with _mock_token():
            result = await _call(mcp_server, "inspect_file_json", {"item_id": "file-id-001"})

        assert _structured(result) == {
            "item_id": "file-id-001",
            "name": "report.csv",
            "size": 1024,
            "content_type": "text/csv",
            "web_url": SAMPLE_DRIVE_ITEM_FILE["webUrl"],
            "modified": "2025-12-15T10:30:00Z",
            "is_folder": False,
        }

    @respx.mock
    async def test_by_url_resolves_the_sharing_link(self, mcp_server):
        respx.get(TEAMS_SHARE_BASE).mock(
            return_value=httpx.Response(200, json=SAMPLE_SHARED_TEXT_FILE)
        )
        with _mock_token():
            result = await _call(mcp_server, "inspect_file_json", {"url": TEAMS_FILE_URL})

        assert _structured(result) == {
            "item_id": "shared-file-002",
            "name": "notes.md",
            "size": 256,
            "content_type": "text/markdown",
            "web_url": SAMPLE_SHARED_TEXT_FILE["webUrl"],
            "modified": "2025-12-21T14:00:00Z",
            "is_folder": False,
        }

    @respx.mock
    async def test_an_item_id_that_is_a_url_is_treated_as_one(self, mcp_server):
        route = respx.get(TEAMS_SHARE_BASE).mock(
            return_value=httpx.Response(200, json=SAMPLE_SHARED_TEXT_FILE)
        )
        with _mock_token():
            result = await _call(mcp_server, "inspect_file_json", {"item_id": TEAMS_FILE_URL})

        assert route.called
        assert _structured(result)["name"] == "notes.md"

    @respx.mock
    async def test_read_content_returns_the_text(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/file-id-001/content").mock(
            return_value=httpx.Response(200, content=b"col1,col2\n1,2\n")
        )
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/file-id-001").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_ITEM_FILE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "inspect_file_json",
                {"item_id": "file-id-001", "read_content": "true"},
            )

        assert _structured(result)["text"] == "col1,col2\n1,2\n"

    @respx.mock
    async def test_a_folder_reads_as_a_folder_with_no_text(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/folder-id-001").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_ITEM_FOLDER)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "inspect_file_json",
                {"item_id": "folder-id-001", "read_content": "true"},
            )

        data = _structured(result)
        assert data["is_folder"] is True
        assert data["text"] is None

    async def test_missing_target(self, mcp_server):
        with _mock_token():
            result = await _call(mcp_server, "inspect_file_json", {})

        assert _structured(result) == {"error": "missing_target"}

    @pytest.mark.parametrize(
        ("status", "error"),
        [(403, "access_denied"), (404, "not_found"), (400, "invalid_link")],
    )
    @respx.mock
    async def test_sharing_link_failures_map_to_sentinels(self, mcp_server, status, error):
        respx.get(TEAMS_SHARE_BASE).mock(return_value=httpx.Response(status, json=GRAPH_ERROR_404))
        with _mock_token():
            result = await _call(mcp_server, "inspect_file_json", {"url": TEAMS_FILE_URL})

        assert _structured(result) == {"error": error}

    async def test_not_connected(self, mcp_server):
        with _mock_missing_connection():
            result = await _call(mcp_server, "inspect_file_json", {"item_id": "file-id-001"})

        assert _structured(result) == {"error": "not_connected", "connect_url": CONNECT_URL}

    @respx.mock
    async def test_an_unknown_item_id_propagates_as_a_tool_error(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/nope").mock(
            return_value=httpx.Response(404, json=GRAPH_ERROR_404)
        )
        from fastmcp.exceptions import ToolError

        with _mock_token():
            with pytest.raises(ToolError, match="404"):
                await _call(mcp_server, "inspect_file_json", {"item_id": "nope"})


class TestMCPUploadTool:
    """Tests for the upload_file MCP tool."""

    @respx.mock
    async def test_upload_creates_file_in_root(self, mcp_server):
        route = respx.put(f"{GRAPH_BASE_URL}/me/drive/root:/notes.md:/content").mock(
            return_value=httpx.Response(201, json=SAMPLE_UPLOADED_FILE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "upload_file",
                    {"filename": "notes.md", "content": "# Hello"},
                )

        text = _get_text(result)
        assert "notes.md" in text
        assert "uploaded" in text.lower()
        assert route.called

    @respx.mock
    async def test_upload_to_subfolder(self, mcp_server):
        route = respx.put(f"{GRAPH_BASE_URL}/me/drive/root:/Documents/data.csv:/content").mock(
            return_value=httpx.Response(201, json=SAMPLE_UPLOADED_FILE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "upload_file",
                    {"filename": "data.csv", "content": "a,b\n1,2", "folder_path": "Documents"},
                )

        assert "uploaded" in _get_text(result).lower()
        assert route.calls[0].request.headers["Content-Type"] == "text/csv"

    @respx.mock
    async def test_upload_to_sharepoint(self, mcp_server):
        site_id = "site-id-001"
        route = respx.put(
            f"{GRAPH_BASE_URL}/sites/{site_id}/drive/root:/Shared Documents/report.html:/content"
        ).mock(return_value=httpx.Response(201, json=SAMPLE_UPLOADED_FILE))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                await client.call_tool(
                    "upload_file",
                    {
                        "filename": "report.html",
                        "content": "<h1>Hello</h1>",
                        "folder_path": "Shared Documents",
                        "site_id": site_id,
                    },
                )

        assert route.called
        assert route.calls[0].request.headers["Content-Type"] == "text/html"

    @respx.mock
    async def test_upload_result_includes_id_and_url(self, mcp_server):
        respx.put(f"{GRAPH_BASE_URL}/me/drive/root:/readme.txt:/content").mock(
            return_value=httpx.Response(201, json=SAMPLE_UPLOADED_FILE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "upload_file", {"filename": "readme.txt", "content": "hello"}
                )

        text = _get_text(result)
        assert SAMPLE_UPLOADED_FILE["id"] in text
        assert SAMPLE_UPLOADED_FILE["webUrl"] in text

    @respx.mock
    async def test_upload_docx_converts_markdown(self, mcp_server):
        """Uploading a .docx file converts markdown content to Word format."""
        route = respx.put(f"{GRAPH_BASE_URL}/me/drive/root:/Review.docx:/content").mock(
            return_value=httpx.Response(201, json=SAMPLE_UPLOADED_FILE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "upload_file",
                    {"filename": "Review.docx", "content": "# Title\n\nA paragraph."},
                )

        text = _get_text(result)
        assert "uploaded" in text.lower()
        assert route.called
        req = route.calls[0].request
        assert req.headers["Content-Type"] == (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        assert req.content[:2] == b"PK"  # ZIP magic bytes = valid .docx

    @respx.mock
    async def test_upload_base64_binary(self, mcp_server):
        """Uploading with content_encoding=base64 decodes and uploads binary."""
        import base64

        raw = b"\x89PNG\r\n\x1a\n" + b"\x00" * 100
        encoded = base64.b64encode(raw).decode()
        route = respx.put(f"{GRAPH_BASE_URL}/me/drive/root:/image.png:/content").mock(
            return_value=httpx.Response(201, json=SAMPLE_UPLOADED_FILE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "upload_file",
                    {
                        "filename": "image.png",
                        "content": encoded,
                        "content_encoding": "base64",
                    },
                )

        text = _get_text(result)
        assert "uploaded" in text.lower()
        assert route.called
        assert route.calls[0].request.content == raw

    async def test_upload_base64_invalid_content(self, mcp_server):
        """Invalid base64 returns a user-friendly error."""
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "upload_file",
                    {
                        "filename": "file.bin",
                        "content": "not-valid-base64!!!",
                        "content_encoding": "base64",
                    },
                )

        text = _get_text(result)
        assert "Failed to decode" in text

    @respx.mock
    async def test_upload_docx_to_sharepoint(self, mcp_server):
        """Docx upload works with SharePoint site_id."""
        site_id = "site-id-001"
        route = respx.put(
            f"{GRAPH_BASE_URL}/sites/{site_id}/drive/root:/Contracts/Review.docx:/content"
        ).mock(return_value=httpx.Response(201, json=SAMPLE_UPLOADED_FILE))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "upload_file",
                    {
                        "filename": "Review.docx",
                        "content": "# Contract\n\n- Clause 1\n- Clause 2",
                        "folder_path": "Contracts",
                        "site_id": site_id,
                    },
                )

        assert "uploaded" in _get_text(result).lower()
        assert route.called


class TestMCPCopyOrRenameTool:
    """Tests for the consolidated manage_file MCP tool (copy/rename/delete)."""

    @pytest.fixture(autouse=True)
    def patch_sleep(self, no_sleep):
        pass

    @respx.mock
    async def test_copy_action(self, mcp_server):
        """action='copy' → server-side copy."""
        item_id = SAMPLE_DRIVE_ITEM_WORD["id"]
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/{item_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_ITEM_WORD)
        )
        respx.post(f"{GRAPH_BASE_URL}/drives/{SOURCE_DRIVE_ID}/items/{item_id}/copy").mock(
            return_value=httpx.Response(202, headers={"Location": MONITOR_URL})
        )
        respx.get(MONITOR_URL).mock(return_value=httpx.Response(200, json=SAMPLE_COPY_COMPLETED))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_file",
                    {"item_id": item_id, "new_name": "template-copy.docx", "action": "copy"},
                )

        text = _get_text(result)
        assert "copied" in text.lower()
        assert "template-copy.docx" in text
        assert SAMPLE_COPY_COMPLETED["resourceId"] in text

    @respx.mock
    async def test_copy_with_destination_folder(self, mcp_server):
        item_id = SAMPLE_DRIVE_ITEM_WORD["id"]
        dest = "folder-id-archive"
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/{item_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_ITEM_WORD)
        )
        copy_route = respx.post(
            f"{GRAPH_BASE_URL}/drives/{SOURCE_DRIVE_ID}/items/{item_id}/copy"
        ).mock(return_value=httpx.Response(202, headers={"Location": MONITOR_URL}))
        respx.get(MONITOR_URL).mock(return_value=httpx.Response(200, json=SAMPLE_COPY_COMPLETED))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                await client.call_tool(
                    "manage_file",
                    {
                        "item_id": item_id,
                        "new_name": "archived.docx",
                        "action": "copy",
                        "options": f'{{"destination_folder_id": "{dest}"}}',
                    },
                )

        copy_body = json.loads(copy_route.calls[0].request.content)
        assert copy_body["parentReference"]["id"] == dest

    @respx.mock
    async def test_copy_error_returns_message(self, mcp_server):
        item_id = SAMPLE_DRIVE_ITEM_WORD["id"]
        respx.get(f"{GRAPH_BASE_URL}/me/drive/items/{item_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_ITEM_WORD)
        )
        respx.post(f"{GRAPH_BASE_URL}/drives/{SOURCE_DRIVE_ID}/items/{item_id}/copy").mock(
            return_value=httpx.Response(202, headers={"Location": MONITOR_URL})
        )
        respx.get(MONITOR_URL).mock(return_value=httpx.Response(200, json=SAMPLE_COPY_FAILED))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_file",
                    {"item_id": item_id, "new_name": "copy.docx", "action": "copy"},
                )

        text = _get_text(result)
        assert "Error performing copy" in text
        assert "accessDenied" in text

    @respx.mock
    async def test_rename_action(self, mcp_server):
        """action='rename' (default) → PATCH rename."""
        item_id = SAMPLE_DRIVE_ITEM_FILE["id"]
        renamed = {**SAMPLE_DRIVE_ITEM_FILE, "name": "final-report.csv"}
        route = respx.patch(f"{GRAPH_BASE_URL}/me/drive/items/{item_id}").mock(
            return_value=httpx.Response(200, json=renamed)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_file",
                    {"item_id": item_id, "new_name": "final-report.csv"},
                )

        text = _get_text(result)
        assert "final-report.csv" in text
        assert "renamed" in text.lower()
        body = json.loads(route.calls[0].request.content)
        assert body == {"name": "final-report.csv"}

    @respx.mock
    async def test_rename_on_sharepoint(self, mcp_server):
        site_id = "site-id-001"
        item_id = SAMPLE_DRIVE_ITEM_WORD["id"]
        renamed = {**SAMPLE_DRIVE_ITEM_WORD, "name": "final-doc.docx"}
        respx.patch(f"{GRAPH_BASE_URL}/sites/{site_id}/drive/items/{item_id}").mock(
            return_value=httpx.Response(200, json=renamed)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_file",
                    {
                        "item_id": item_id,
                        "new_name": "final-doc.docx",
                        "options": f'{{"site_id": "{site_id}"}}',
                    },
                )

        assert "final-doc.docx" in _get_text(result)

    async def test_invalid_action(self, mcp_server):
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_file",
                    {"item_id": "x", "new_name": "y", "action": "move"},
                )

        text = _get_text(result)
        assert "Invalid action" in text
        assert "move" in text

    @respx.mock
    async def test_rename_not_found_returns_message(self, mcp_server):
        item_id = SAMPLE_DRIVE_ITEM_FILE["id"]
        respx.patch(f"{GRAPH_BASE_URL}/me/drive/items/{item_id}").mock(
            return_value=httpx.Response(
                404, json={"error": {"code": "ResourceNotFound", "message": "Item not found."}}
            )
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_file",
                    {"item_id": item_id, "new_name": "x.csv"},
                )

        assert "not found" in _get_text(result).lower()

    @respx.mock
    async def test_delete_not_found_returns_message(self, mcp_server):
        item_id = SAMPLE_DRIVE_ITEM_FILE["id"]
        respx.delete(f"{GRAPH_BASE_URL}/me/drive/items/{item_id}").mock(
            return_value=httpx.Response(
                404, json={"error": {"code": "itemNotFound", "message": "Item not found."}}
            )
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_file",
                    {"item_id": item_id, "action": "delete"},
                )

        assert "not found" in _get_text(result).lower()

    @respx.mock
    async def test_delete_action(self, mcp_server):
        """action='delete' → DELETE the item (204 No Content)."""
        item_id = SAMPLE_DRIVE_ITEM_FILE["id"]
        route = respx.delete(f"{GRAPH_BASE_URL}/me/drive/items/{item_id}").mock(
            return_value=httpx.Response(204)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_file",
                    {"item_id": item_id, "action": "delete"},
                )

        assert route.called
        assert "Deleted" in _get_text(result)

    @respx.mock
    async def test_delete_on_sharepoint(self, mcp_server):
        site_id = "site-id-001"
        item_id = SAMPLE_DRIVE_ITEM_WORD["id"]
        route = respx.delete(f"{GRAPH_BASE_URL}/sites/{site_id}/drive/items/{item_id}").mock(
            return_value=httpx.Response(204)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_file",
                    {
                        "item_id": item_id,
                        "action": "delete",
                        "options": f'{{"site_id": "{site_id}"}}',
                    },
                )

        assert route.called
        assert "Deleted" in _get_text(result)

    async def test_rename_requires_new_name(self, mcp_server):
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_file",
                    {"item_id": "x", "action": "rename"},
                )

        assert "new_name is required" in _get_text(result)


# ---------------------------------------------------------------------------
# Power BI
# ---------------------------------------------------------------------------


class TestMCPPowerBITools:
    """Tests for the Power BI MCP tools."""

    @pytest.fixture(autouse=True)
    def patch_sleep(self, no_sleep):
        pass

    @respx.mock
    async def test_list_powerbi_workspaces(self, mcp_server):
        respx.get(f"{POWERBI_BASE_URL}/groups").mock(
            return_value=httpx.Response(200, json=SAMPLE_PBI_WORKSPACES_RESPONSE)
        )
        with _mock_pbi_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_powerbi_workspaces", {})

        text = _get_text(result)
        assert "Analytics Hub" in text
        assert "Finance Reports" in text
        assert "My workspace" in text
        assert "3 workspace(s)" in text

    @respx.mock
    async def test_list_powerbi_workspaces_empty(self, mcp_server):
        respx.get(f"{POWERBI_BASE_URL}/groups").mock(
            return_value=httpx.Response(200, json={"value": []})
        )
        with _mock_pbi_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_powerbi_workspaces", {})

        # Even with no named workspaces, My workspace is always shown
        text = _get_text(result)
        assert "My workspace" in text
        assert "1 workspace(s)" in text

    @respx.mock
    async def test_list_powerbi_content_all(self, mcp_server):
        """content_type=all returns datasets, reports, and dashboards."""
        ws_id = "ws-id-001"
        respx.get(f"{POWERBI_BASE_URL}/groups/{ws_id}/datasets").mock(
            return_value=httpx.Response(200, json=SAMPLE_PBI_DATASETS_RESPONSE)
        )
        respx.get(f"{POWERBI_BASE_URL}/groups/{ws_id}/reports").mock(
            return_value=httpx.Response(200, json=SAMPLE_PBI_REPORTS_RESPONSE)
        )
        respx.get(f"{POWERBI_BASE_URL}/groups/{ws_id}/dashboards").mock(
            return_value=httpx.Response(200, json=SAMPLE_PBI_DASHBOARDS_RESPONSE)
        )
        with _mock_pbi_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_powerbi_content", {"workspace_id": ws_id})

        text = _get_text(result)
        assert "Datasets" in text
        assert "Sales" in text
        assert "Reports" in text
        assert "Q4 Dashboard" in text
        assert "Dashboards" in text
        assert "Executive Overview" in text

    @respx.mock
    async def test_list_powerbi_content_datasets_only(self, mcp_server):
        ws_id = "ws-id-001"
        respx.get(f"{POWERBI_BASE_URL}/groups/{ws_id}/datasets").mock(
            return_value=httpx.Response(200, json=SAMPLE_PBI_DATASETS_RESPONSE)
        )
        with _mock_pbi_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "list_powerbi_content",
                    {"workspace_id": ws_id, "content_type": "datasets"},
                )

        text = _get_text(result)
        assert "Datasets" in text
        assert "Sales" in text
        assert "Reports" not in text

    @respx.mock
    async def test_list_powerbi_workspaces_includes_my_workspace(self, mcp_server):
        """list_powerbi_workspaces always includes My workspace with id='me'."""
        respx.get(f"{POWERBI_BASE_URL}/groups").mock(
            return_value=httpx.Response(200, json=SAMPLE_PBI_WORKSPACES_RESPONSE)
        )
        with _mock_pbi_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_powerbi_workspaces", {})

        text = _get_text(result)
        assert "My workspace" in text
        assert "me" in text  # sentinel ID
        assert "3 workspace(s)" in text  # 1 My + 2 named

    @respx.mock
    async def test_list_powerbi_content_my_workspace(self, mcp_server):
        """workspace_id='me' routes to root /datasets endpoint, not /groups/me/."""
        respx.get(f"{POWERBI_BASE_URL}/datasets").mock(
            return_value=httpx.Response(200, json=SAMPLE_PBI_DATASETS_RESPONSE)
        )
        respx.get(f"{POWERBI_BASE_URL}/reports").mock(
            return_value=httpx.Response(200, json=SAMPLE_PBI_REPORTS_RESPONSE)
        )
        respx.get(f"{POWERBI_BASE_URL}/dashboards").mock(
            return_value=httpx.Response(200, json=SAMPLE_PBI_DASHBOARDS_RESPONSE)
        )
        with _mock_pbi_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("list_powerbi_content", {"workspace_id": "me"})

        text = _get_text(result)
        assert "Sales" in text
        assert "Q4 Dashboard" in text

    async def test_list_powerbi_content_invalid_type(self, mcp_server):
        with _mock_pbi_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "list_powerbi_content",
                    {"workspace_id": "ws-id-001", "content_type": "tiles"},
                )

        assert "Invalid content_type" in _get_text(result)

    @respx.mock
    async def test_query_dataset(self, mcp_server):
        ws_id = "ws-id-001"
        ds_id = "ds-id-001"
        route = respx.post(
            f"{POWERBI_BASE_URL}/groups/{ws_id}/datasets/{ds_id}/executeQueries"
        ).mock(return_value=httpx.Response(200, json=SAMPLE_PBI_DAX_RESULT))
        with _mock_pbi_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "query_dataset",
                    {"workspace_id": ws_id, "dataset_id": ds_id, "dax_query": "EVALUATE 'Sales'"},
                )

        text = _get_text(result)
        assert "3 row(s)" in text
        assert "West" in text
        assert "[Region]" in text
        body = json.loads(route.calls[0].request.content)
        assert body["queries"][0]["query"] == "EVALUATE 'Sales'"

    @respx.mock
    async def test_refresh_dataset(self, mcp_server):
        ws_id = "ws-id-001"
        ds_id = "ds-id-001"
        route = respx.post(f"{POWERBI_BASE_URL}/groups/{ws_id}/datasets/{ds_id}/refreshes").mock(
            return_value=httpx.Response(202)
        )
        with _mock_pbi_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "refresh_dataset",
                    {"workspace_id": ws_id, "dataset_id": ds_id},
                )

        text = _get_text(result)
        assert "Refresh triggered" in text
        assert route.called

    @respx.mock
    async def test_query_dataset_my_workspace(self, mcp_server):
        """workspace_id='me' routes to root /datasets/... not /groups/me/..."""
        ds_id = "ds-id-001"
        route = respx.post(f"{POWERBI_BASE_URL}/datasets/{ds_id}/executeQueries").mock(
            return_value=httpx.Response(200, json=SAMPLE_PBI_DAX_RESULT)
        )
        with _mock_pbi_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "query_dataset",
                    {"workspace_id": "me", "dataset_id": ds_id, "dax_query": "EVALUATE 'Sales'"},
                )

        assert "3 row(s)" in _get_text(result)
        assert route.called
        assert "/groups/" not in str(route.calls[0].request.url)

    @respx.mock
    async def test_refresh_dataset_my_workspace(self, mcp_server):
        """workspace_id='me' routes to root /datasets/... not /groups/me/..."""
        ds_id = "ds-id-001"
        route = respx.post(f"{POWERBI_BASE_URL}/datasets/{ds_id}/refreshes").mock(
            return_value=httpx.Response(202)
        )
        with _mock_pbi_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "refresh_dataset",
                    {"workspace_id": "me", "dataset_id": ds_id},
                )

        assert "Refresh triggered" in _get_text(result)
        assert route.called
        assert "/groups/" not in str(route.calls[0].request.url)

    @respx.mock
    async def test_export_report_onedrive_upload_failure_degrades_gracefully(self, mcp_server):
        """If Graph token is missing, export still succeeds with a helpful message."""
        ws_id = "ws-id-001"
        rpt_id = "rpt-id-001"
        export_id = "export-id-001"
        export_location = f"{POWERBI_BASE_URL}/groups/{ws_id}/reports/{rpt_id}/exports/{export_id}"

        respx.post(f"{POWERBI_BASE_URL}/groups/{ws_id}/reports/{rpt_id}/ExportTo").mock(
            return_value=httpx.Response(202, headers={"Location": export_location})
        )
        respx.get(f"{POWERBI_BASE_URL}/groups/{ws_id}/reports/{rpt_id}/exports/{export_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_PBI_EXPORT_SUCCEEDED)
        )
        respx.get(
            f"{POWERBI_BASE_URL}/groups/{ws_id}/reports/{rpt_id}/exports/{export_id}/file"
        ).mock(return_value=httpx.Response(200, content=b"%PDF fake"))
        with (
            _mock_pbi_token(),
            patch("ms_graph_mcp.get_graph_token", side_effect=PermissionError("Not connected.")),
        ):
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "export_report",
                    {"workspace_id": ws_id, "report_id": rpt_id, "export_format": "PDF"},
                )

        text = _get_text(result)
        assert "exported" in text.lower()
        assert "OneDrive" in text
        assert "Microsoft auth" in text
        assert "PDF" in text

    @respx.mock
    async def test_export_report_success(self, mcp_server):
        """Export flow: PBI export → download bytes → upload to OneDrive → return URL."""
        ws_id = "ws-id-001"
        rpt_id = "rpt-id-001"
        export_id = "export-id-001"
        export_location = f"{POWERBI_BASE_URL}/groups/{ws_id}/reports/{rpt_id}/exports/{export_id}"
        fake_pdf = b"%PDF-1.4 fake content"

        # PBI: start export → poll → download
        respx.post(f"{POWERBI_BASE_URL}/groups/{ws_id}/reports/{rpt_id}/ExportTo").mock(
            return_value=httpx.Response(202, headers={"Location": export_location})
        )
        respx.get(f"{POWERBI_BASE_URL}/groups/{ws_id}/reports/{rpt_id}/exports/{export_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_PBI_EXPORT_SUCCEEDED)
        )
        respx.get(
            f"{POWERBI_BASE_URL}/groups/{ws_id}/reports/{rpt_id}/exports/{export_id}/file"
        ).mock(return_value=httpx.Response(200, content=fake_pdf))
        # Graph: upload to OneDrive
        upload_route = respx.put(url__regex=r"/content$").mock(
            return_value=httpx.Response(201, json=SAMPLE_UPLOADED_FILE)
        )

        with _mock_pbi_token(), _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "export_report",
                    {"workspace_id": ws_id, "report_id": rpt_id, "export_format": "PDF"},
                )

        text = _get_text(result)
        assert "exported" in text.lower()
        assert "PDF" in text
        assert "OneDrive" in text
        assert SAMPLE_UPLOADED_FILE["webUrl"] in text
        # Verify correct content-type was sent for PDF
        assert upload_route.calls[0].request.headers["Content-Type"] == "application/pdf"
        # Verify the bytes uploaded match what was downloaded
        assert upload_route.calls[0].request.content == fake_pdf

    async def test_export_report_invalid_format(self, mcp_server):
        with _mock_pbi_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "export_report",
                    {"workspace_id": "ws-1", "report_id": "rpt-1", "export_format": "DOCX"},
                )

        assert "Invalid export_format" in _get_text(result)

    @respx.mock
    async def test_powerbi_auth_error_propagates(self, mcp_server):
        from fastmcp.exceptions import ToolError

        respx.get(f"{POWERBI_BASE_URL}/groups").mock(
            return_value=httpx.Response(
                401, json={"error": {"code": "Unauthorized", "message": "Token expired."}}
            )
        )
        with _mock_pbi_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                with pytest.raises(ToolError, match="Unauthorized"):
                    await client.call_tool("list_powerbi_workspaces", {})


# ---------------------------------------------------------------------------
# Auth
# ---------------------------------------------------------------------------


class TestMCPAuth:
    """Test authentication behavior."""

    async def test_missing_token_raises_error(self, mcp_server):
        from fastmcp import Client
        from fastmcp.exceptions import ToolError

        with patch(
            "ms_graph_mcp.get_graph_token", side_effect=PermissionError("Authorization required.")
        ):
            async with Client(mcp_server) as client:
                with pytest.raises(ToolError, match="Authorization required"):
                    await client.call_tool("list_emails", {})

    async def test_all_tools_require_auth(self, mcp_server):
        """Verify every tool rejects unauthenticated requests."""
        from fastmcp import Client
        from fastmcp.exceptions import ToolError

        # Graph tools — blocked by get_graph_token failing
        graph_tools = [
            ("get_user_profile", {}),
            ("list_emails", {}),
            ("read_email", {"message_id": "fake-id"}),
            ("send_email", {"to": "a@b.com", "subject": "S", "body": "B"}),
            ("list_teams", {}),
            ("list_chats", {}),
            ("read_teams_messages", {"chat_id": "c1"}),
            ("send_teams_message", {"message": "Hi", "chat_id": "c1"}),
            ("get_teams_activity", {}),
            ("list_sharepoint_sites", {}),
            ("list_files", {}),
            ("inspect_file", {"item_id": "x"}),
            ("upload_file", {"filename": "x.txt", "content": "y"}),
            ("manage_file", {"item_id": "x", "new_name": "y"}),
        ]
        with patch(
            "ms_graph_mcp.get_graph_token", side_effect=PermissionError("Authorization required.")
        ):
            async with Client(mcp_server) as client:
                for tool_name, args in graph_tools:
                    with pytest.raises(ToolError, match="Authorization required"):
                        await client.call_tool(tool_name, args)

        # Power BI tools — blocked by get_powerbi_token failing
        pbi_tools = [
            ("list_powerbi_workspaces", {}),
            ("list_powerbi_content", {"workspace_id": "ws-1"}),
            (
                "query_dataset",
                {"workspace_id": "ws-1", "dataset_id": "ds-1", "dax_query": "EVALUATE {1}"},
            ),
            ("refresh_dataset", {"workspace_id": "ws-1", "dataset_id": "ds-1"}),
            ("export_report", {"workspace_id": "ws-1", "report_id": "rpt-1"}),
        ]
        with patch(
            "ms_graph_mcp.get_powerbi_token", side_effect=PermissionError("Authorization required.")
        ):
            async with Client(mcp_server) as client:
                for tool_name, args in pbi_tools:
                    with pytest.raises(ToolError, match="Authorization required"):
                        await client.call_tool(tool_name, args)


# ---------------------------------------------------------------------------
# Desktop JSON tools
# ---------------------------------------------------------------------------


def _structured(result) -> dict:
    """Extract the dict a Desktop JSON tool returned.

    FastMCP surfaces dict returns as structuredContent; fall back to parsing
    the text block for client versions that do not populate it.
    """
    if result.structured_content is not None:
        return result.structured_content
    return json.loads(_get_text(result))


def _mock_missing_connection(connect_url=CONNECT_URL):
    """Patch get_graph_token to raise as JWT mode does for an unconnected user."""
    return patch(
        "ms_graph_mcp.get_graph_token",
        side_effect=MissingProviderConnection(
            provider="microsoft", user_key="u", connect_url=connect_url
        ),
    )


async def _call(mcp_server, name, args=None):
    from fastmcp import Client

    async with Client(mcp_server) as client:
        return await client.call_tool(name, args or {})


class TestPostExchangeScopesKey:
    """The token-exchange shim must persist scopes under the storage key."""

    def test_scope_is_persisted_as_scopes(self):
        from ms_graph_mcp import _microsoft_post_exchange

        out = _microsoft_post_exchange(
            {"access_token": "a", "scope": "Mail.Read Chat.Read", "expires_in": 3600}
        )

        assert out["scopes"] == "Mail.Read Chat.Read"
        assert "scope" not in out
        assert out["access_token"] == "a"
        assert "expires_at" in out

    def test_absent_scope_is_omitted(self):
        from ms_graph_mcp import _microsoft_post_exchange

        out = _microsoft_post_exchange({"access_token": "a"})

        assert "scopes" not in out
        assert "scope" not in out


class TestMCPProfileJson:
    """get_profile_json."""

    @respx.mock
    async def test_maps_graph_fields(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me").mock(
            return_value=httpx.Response(200, json=SAMPLE_USER_PROFILE)
        )
        respx.get(f"{GRAPH_BASE_URL}/me/mailboxSettings").mock(
            return_value=httpx.Response(200, json=SAMPLE_MAILBOX_SETTINGS)
        )
        with _mock_token():
            result = await _call(mcp_server, "get_profile_json")

        assert _structured(result) == {
            "id": "user-id-001",
            "display_name": "Test User",
            "mail": "user@example.com",
            "user_principal_name": "user@example.com",
        }

    async def test_not_connected(self, mcp_server):
        with _mock_missing_connection():
            result = await _call(mcp_server, "get_profile_json")

        assert _structured(result) == {"error": "not_connected", "connect_url": CONNECT_URL}


class TestMCPListMailDelta:
    """list_mail_delta."""

    @respx.mock
    async def test_maps_next_link_and_passes_messages_through_raw(self, mcp_server):
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages/delta").mock(
            return_value=httpx.Response(200, json=SAMPLE_DELTA_PAGE_NEXT)
        )
        with _mock_token():
            result = await _call(mcp_server, "list_mail_delta", {"folder": "inbox"})

        data = _structured(result)
        assert data["messages"] == [SAMPLE_DELTA_MESSAGE]
        assert data["next_cursor"] == SAMPLE_DELTA_NEXT_LINK
        assert data["delta_cursor"] == ""
        assert data["resync"] is False

    @respx.mock
    async def test_tombstones_pass_through_untouched(self, mcp_server):
        """The client's fold logic owns @removed entries — do not filter them here."""
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages/delta").mock(
            return_value=httpx.Response(200, json=SAMPLE_DELTA_PAGE_FINAL)
        )
        with _mock_token():
            result = await _call(mcp_server, "list_mail_delta", {})

        data = _structured(result)
        assert data["messages"] == [SAMPLE_DELTA_TOMBSTONE]
        assert data["messages"][0]["@removed"] == {"reason": "deleted"}
        assert data["next_cursor"] == ""
        assert data["delta_cursor"] == SAMPLE_DELTA_LINK

    @respx.mock
    async def test_cursor_is_requested_verbatim(self, mcp_server):
        route = respx.get(SAMPLE_DELTA_NEXT_LINK).mock(
            return_value=httpx.Response(200, json=SAMPLE_DELTA_PAGE_FINAL)
        )
        with _mock_token():
            result = await _call(mcp_server, "list_mail_delta", {"cursor": SAMPLE_DELTA_NEXT_LINK})

        assert route.call_count == 1
        assert str(route.calls[0].request.url) == SAMPLE_DELTA_NEXT_LINK
        assert _structured(result)["delta_cursor"] == SAMPLE_DELTA_LINK

    @respx.mock
    async def test_expired_cursor_returns_resync(self, mcp_server):
        respx.get(SAMPLE_DELTA_LINK).mock(return_value=httpx.Response(410, json=GRAPH_ERROR_410))
        with _mock_token():
            result = await _call(mcp_server, "list_mail_delta", {"cursor": SAMPLE_DELTA_LINK})

        assert _structured(result) == {
            "messages": [],
            "next_cursor": "",
            "delta_cursor": "",
            "resync": True,
        }

    @respx.mock
    async def test_other_graph_errors_propagate(self, mcp_server):
        """A 500 is the client's "transient, retry later" signal — not a resync."""
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages/delta").mock(
            return_value=httpx.Response(500, json={"error": {"code": "x", "message": "boom"}})
        )
        from fastmcp.exceptions import ToolError

        with _mock_token():
            with pytest.raises(ToolError, match="500"):
                await _call(mcp_server, "list_mail_delta", {})

    async def test_not_connected(self, mcp_server):
        with _mock_missing_connection():
            result = await _call(mcp_server, "list_mail_delta", {})

        assert _structured(result) == {"error": "not_connected", "connect_url": CONNECT_URL}

    async def test_plain_permission_error_has_no_connect_url(self, mcp_server):
        """Laptop (MSAL) mode raises a bare PermissionError with no URL to offer."""
        with patch("ms_graph_mcp.get_graph_token", side_effect=PermissionError("no auth")):
            result = await _call(mcp_server, "list_mail_delta", {})

        assert _structured(result) == {"error": "not_connected", "connect_url": None}


class TestMCPGetMailDetail:
    """get_mail_detail."""

    @respx.mock
    async def test_lowercases_headers_and_keeps_first_occurrence(self, mcp_server):
        route = respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE_DETAIL)
        )
        with _mock_token():
            result = await _call(
                mcp_server, "get_mail_detail", {"message_id": SAMPLE_MESSAGE["id"]}
            )

        data = _structured(result)
        assert data["body_text"] == "Here is the weekly report.\n\nBest,\nAlice"
        assert data["has_attachments"] is True
        assert data["headers"]["message-id"] == "<abc123@example.com>"
        assert data["headers"]["in-reply-to"] == "<parent@example.com>"
        # "Received" appeared twice with different casing; the first one wins.
        assert data["headers"]["received"] == "from mx1.example.com"

        # The attachments ride along on the $expand — no second request.
        assert route.call_count == 1
        assert data["attachment_count"] == 3
        assert data["attachments"][0] == {
            "id": SAMPLE_FILE_ATTACHMENT["id"],
            "name": "report.pdf",
            "content_type": "application/pdf",
            "size": 1_258_291,
            "is_inline": False,
            "content_id": None,
            "kind": "file",
            "source_url": None,
        }
        assert data["attachments"][1]["is_inline"] is True
        assert data["attachments"][1]["content_id"] == "logo@company"
        assert data["attachments"][2]["kind"] == "reference"
        assert data["attachments"][2]["source_url"] == SAMPLE_REFERENCE_ATTACHMENT["sourceUrl"]

    @respx.mock
    async def test_attachment_list_is_capped_but_count_is_true(self, mcp_server):
        """A pathological message lists 50 and still reports how many there are."""
        many = [{**SAMPLE_FILE_ATTACHMENT, "id": f"att-{i}"} for i in range(51)]
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json={**SAMPLE_MESSAGE_DETAIL, "attachments": many})
        )
        with _mock_token():
            result = await _call(
                mcp_server, "get_mail_detail", {"message_id": SAMPLE_MESSAGE["id"]}
            )

        data = _structured(result)
        assert len(data["attachments"]) == 50
        assert data["attachment_count"] == 51
        assert data["attachments"][-1]["id"] == "att-49"

    @respx.mock
    async def test_missing_unique_body_becomes_empty_string(self, mcp_server):
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE_DETAIL_NO_BODY)
        )
        with _mock_token():
            result = await _call(
                mcp_server, "get_mail_detail", {"message_id": SAMPLE_MESSAGE["id"]}
            )

        assert _structured(result) == {
            "body_text": "",
            "headers": {},
            "has_attachments": False,
            "attachments": [],
            "attachment_count": 0,
        }

    async def test_not_connected(self, mcp_server):
        with _mock_missing_connection():
            result = await _call(mcp_server, "get_mail_detail", {"message_id": "m"})

        assert _structured(result)["error"] == "not_connected"


class TestMCPGetMailAttachmentJson:
    """get_mail_attachment_json — the desktop attachment reader."""

    @respx.mock
    async def test_metadata_mode_returns_the_summary_only(self, mcp_server):
        value_route = respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(200, content=b"x")
        )
        respx.get(ATT_FILE_URL).mock(return_value=httpx.Response(200, json=SAMPLE_FILE_ATTACHMENT))
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_FILE_ATTACHMENT["id"],
                    "mode": "metadata",
                },
            )

        assert _structured(result) == {
            "id": SAMPLE_FILE_ATTACHMENT["id"],
            "name": "report.pdf",
            "content_type": "application/pdf",
            "size": 1_258_291,
            "is_inline": False,
            "content_id": None,
            "kind": "file",
            "source_url": None,
        }
        assert not value_route.called

    @respx.mock
    async def test_metadata_mode_adds_the_inner_message_fields(self, mcp_server):
        inner = {
            "subject": "Budget draft",
            "from": {"emailAddress": {"name": "Dana Lee", "address": "dana@example.com"}},
            "receivedDateTime": "2025-12-01T09:00:00Z",
        }

        def _respond(request):
            if "expand" in str(request.url):
                return httpx.Response(200, json={**SAMPLE_ITEM_ATTACHMENT, "item": inner})
            return httpx.Response(200, json=SAMPLE_ITEM_ATTACHMENT)

        respx.get(url__startswith=ATT_ITEM_URL).mock(side_effect=_respond)
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_ITEM_ATTACHMENT["id"],
                    "mode": "metadata",
                },
            )

        data = _structured(result)
        assert data["kind"] == "item"
        assert data["item_subject"] == "Budget draft"
        assert data["item_from"] == "dana@example.com"
        assert data["item_received"] == "2025-12-01T09:00:00Z"

    @respx.mock
    async def test_text_mode_decodes_a_plain_text_attachment(self, mcp_server):
        respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(
                200, content=b"hello there", headers={"Content-Type": "text/plain"}
            )
        )
        respx.get(ATT_FILE_URL).mock(
            return_value=httpx.Response(
                200,
                json={
                    **SAMPLE_FILE_ATTACHMENT,
                    "name": "notes.txt",
                    "contentType": "text/plain",
                    "size": 11,
                },
            )
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_FILE_ATTACHMENT["id"],
                    "mode": "text",
                },
            )

        data = _structured(result)
        assert data["text"] == "hello there"
        assert data["truncated"] is False
        assert "reason" not in data

    @respx.mock
    async def test_text_mode_extracts_a_word_document(self, mcp_server):
        docx = _docx_bytes()
        respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(200, content=docx, headers={"Content-Type": DOCX_MIME})
        )
        respx.get(ATT_FILE_URL).mock(
            return_value=httpx.Response(
                200,
                json={
                    **SAMPLE_FILE_ATTACHMENT,
                    "name": "report.docx",
                    "contentType": DOCX_MIME,
                    "size": len(docx),
                },
            )
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_FILE_ATTACHMENT["id"],
                    "mode": "text",
                },
            )

        data = _structured(result)
        assert "Quarterly Title" in data["text"]
        assert data["content_type"] == DOCX_MIME

    @respx.mock
    async def test_text_mode_on_a_binary_reports_the_reason(self, mcp_server):
        respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(
                200, content=b"\x89PNG\r\n\x1a\n", headers={"Content-Type": "image/png"}
            )
        )
        respx.get(ATT_FILE_URL).mock(
            return_value=httpx.Response(
                200,
                json={
                    **SAMPLE_FILE_ATTACHMENT,
                    "name": "logo.png",
                    "contentType": "image/png",
                    "size": 8,
                },
            )
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_FILE_ATTACHMENT["id"],
                    "mode": "text",
                },
            )

        data = _structured(result)
        assert data["text"] is None
        assert data["reason"] == "binary"
        assert data["truncated"] is False

    @respx.mock
    async def test_text_mode_on_a_link_attachment_reports_reference(self, mcp_server):
        respx.get(ATT_REF_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_REFERENCE_ATTACHMENT)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_REFERENCE_ATTACHMENT["id"],
                    "mode": "text",
                },
            )

        data = _structured(result)
        assert data["text"] is None
        assert data["reason"] == "reference"
        assert data["source_url"] == SAMPLE_REFERENCE_ATTACHMENT["sourceUrl"]

    @respx.mock
    async def test_text_mode_on_an_item_returns_the_inner_body(self, mcp_server):
        inner = {
            "subject": "Budget draft",
            "from": {"emailAddress": {"address": "dana@example.com"}},
            "receivedDateTime": "2025-12-01T09:00:00Z",
            "bodyPreview": "Numbers attached",
            "body": {"contentType": "text", "content": "Numbers attached, see inside."},
        }

        def _respond(request):
            if "expand" in str(request.url):
                return httpx.Response(200, json={**SAMPLE_ITEM_ATTACHMENT, "item": inner})
            return httpx.Response(200, json=SAMPLE_ITEM_ATTACHMENT)

        respx.get(url__startswith=ATT_ITEM_URL).mock(side_effect=_respond)
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_ITEM_ATTACHMENT["id"],
                    "mode": "text",
                },
            )

        data = _structured(result)
        assert data["text"] == "Numbers attached, see inside."
        assert data["truncated"] is False
        assert data["item_subject"] == "Budget draft"

    @respx.mock
    async def test_bytes_mode_returns_the_content(self, mcp_server):
        payload = b"\x89PNG\r\n\x1a\n"
        respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(200, content=payload, headers={"Content-Type": "image/png"})
        )
        respx.get(ATT_FILE_URL).mock(
            return_value=httpx.Response(
                200,
                json={
                    **SAMPLE_FILE_ATTACHMENT,
                    "name": "logo.png",
                    "contentType": "image/png",
                    "size": len(payload),
                },
            )
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_FILE_ATTACHMENT["id"]},
            )

        data = _structured(result)
        assert base64.b64decode(data["content_base64"]) == payload
        assert data["content_type"] == "image/png"
        assert data["kind"] == "file"

    @respx.mock
    async def test_bytes_mode_refuses_oversize_without_downloading(self, mcp_server):
        value_route = respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(200, content=b"x")
        )
        respx.get(ATT_FILE_URL).mock(
            return_value=httpx.Response(200, json={**SAMPLE_FILE_ATTACHMENT, "size": 20_000_000})
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_FILE_ATTACHMENT["id"]},
            )

        assert _structured(result) == {
            "error": "too_large",
            "size": 20_000_000,
            "limit": 10_000_000,
        }
        assert not value_route.called

    @respx.mock
    async def test_bytes_mode_on_a_link_attachment_returns_the_url(self, mcp_server):
        respx.get(ATT_REF_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_REFERENCE_ATTACHMENT)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_REFERENCE_ATTACHMENT["id"]},
            )

        assert _structured(result) == {
            "error": "reference",
            "source_url": SAMPLE_REFERENCE_ATTACHMENT["sourceUrl"],
        }

    @respx.mock
    async def test_bytes_mode_on_an_item_falls_back_to_rfc822(self, mcp_server):
        """An item attachment has no contentType, but its $value is a MIME message."""
        respx.get(f"{ATT_ITEM_URL}/$value").mock(
            return_value=httpx.Response(200, content=b"From: dana@example.com")
        )
        respx.get(ATT_ITEM_URL).mock(
            return_value=httpx.Response(200, json={**SAMPLE_ITEM_ATTACHMENT, "size": 22})
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_ITEM_ATTACHMENT["id"]},
            )

        data = _structured(result)
        assert base64.b64decode(data["content_base64"]) == b"From: dana@example.com"
        assert data["content_type"] == "message/rfc822"

    async def test_invalid_mode_makes_no_request(self, mcp_server):
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {"message_id": ATT_MSG_ID, "attachment_id": "a", "mode": "onedrive"},
            )

        assert _structured(result) == {"error": "invalid_mode"}

    async def test_not_connected(self, mcp_server):
        with _mock_missing_connection():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {"message_id": ATT_MSG_ID, "attachment_id": "a"},
            )

        assert _structured(result) == {"error": "not_connected", "connect_url": CONNECT_URL}

    @respx.mock
    async def test_unknown_attachment_propagates_as_a_tool_error(self, mcp_server):
        respx.get(ATT_FILE_URL).mock(return_value=httpx.Response(404, json=GRAPH_ERROR_404))
        from fastmcp.exceptions import ToolError

        with _mock_token():
            with pytest.raises(ToolError, match="404"):
                await _call(
                    mcp_server,
                    "get_mail_attachment_json",
                    {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_FILE_ATTACHMENT["id"]},
                )


class TestMCPAddDraftAttachmentJson:
    """add_draft_attachment_json — attach bytes to a draft before sending."""

    @respx.mock
    async def test_attaches_and_returns_the_id(self, mcp_server):
        route = respx.post(f"{DRAFT_BASE}/attachments").mock(
            return_value=httpx.Response(201, json=SAMPLE_CREATED_ATTACHMENT)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "add_draft_attachment_json",
                {
                    "draft_id": SAMPLE_DRAFT_MESSAGE["id"],
                    "name": "notes.txt",
                    "content_base64": base64.b64encode(b"hello world").decode("ascii"),
                },
            )

        assert _structured(result) == {"attachment_id": SAMPLE_CREATED_ATTACHMENT["id"]}
        payload = json.loads(route.calls[0].request.content)
        assert payload["@odata.type"] == "#microsoft.graph.fileAttachment"
        assert payload["name"] == "notes.txt"
        assert payload["contentType"] == "text/plain"
        assert base64.b64decode(payload["contentBytes"]) == b"hello world"

    @respx.mock
    async def test_explicit_content_type_wins_over_the_guess(self, mcp_server):
        route = respx.post(f"{DRAFT_BASE}/attachments").mock(
            return_value=httpx.Response(201, json=SAMPLE_CREATED_ATTACHMENT)
        )
        with _mock_token():
            await _call(
                mcp_server,
                "add_draft_attachment_json",
                {
                    "draft_id": SAMPLE_DRAFT_MESSAGE["id"],
                    "name": "notes.txt",
                    "content_base64": base64.b64encode(b"hi").decode("ascii"),
                    "content_type": "text/markdown",
                },
            )

        assert json.loads(route.calls[0].request.content)["contentType"] == "text/markdown"

    async def test_empty_name_is_rejected(self, mcp_server):
        with _mock_token():
            result = await _call(
                mcp_server,
                "add_draft_attachment_json",
                {
                    "draft_id": SAMPLE_DRAFT_MESSAGE["id"],
                    "name": "   ",
                    "content_base64": base64.b64encode(b"hi").decode("ascii"),
                },
            )

        assert _structured(result) == {"error": "empty_name"}

    async def test_invalid_base64_is_rejected(self, mcp_server):
        with _mock_token():
            result = await _call(
                mcp_server,
                "add_draft_attachment_json",
                {
                    "draft_id": SAMPLE_DRAFT_MESSAGE["id"],
                    "name": "notes.txt",
                    "content_base64": "not base64!!",
                },
            )

        assert _structured(result) == {"error": "invalid_base64"}

    async def test_empty_content_is_rejected(self, mcp_server):
        with _mock_token():
            result = await _call(
                mcp_server,
                "add_draft_attachment_json",
                {
                    "draft_id": SAMPLE_DRAFT_MESSAGE["id"],
                    "name": "notes.txt",
                    "content_base64": "",
                },
            )

        assert _structured(result) == {"error": "invalid_base64"}

    async def test_oversize_is_rejected_before_any_request(self, mcp_server, monkeypatch):
        from ms_graph import attachments as attachment_ops

        monkeypatch.setattr(attachment_ops, "MAX_ATTACHMENT_BYTES", 4)
        with _mock_token():
            result = await _call(
                mcp_server,
                "add_draft_attachment_json",
                {
                    "draft_id": SAMPLE_DRAFT_MESSAGE["id"],
                    "name": "notes.txt",
                    "content_base64": base64.b64encode(b"hello world").decode("ascii"),
                },
            )

        assert _structured(result) == {"error": "too_large", "size": 11, "limit": 4}

    async def test_not_connected(self, mcp_server):
        with _mock_missing_connection():
            result = await _call(
                mcp_server,
                "add_draft_attachment_json",
                {
                    "draft_id": SAMPLE_DRAFT_MESSAGE["id"],
                    "name": "notes.txt",
                    "content_base64": base64.b64encode(b"hi").decode("ascii"),
                },
            )

        assert _structured(result) == {"error": "not_connected", "connect_url": CONNECT_URL}


class TestMCPDraftFlow:
    """create_reply_draft_json, update_draft_body, send_draft."""

    @respx.mock
    async def test_create_reply_draft_returns_id_and_link(self, mcp_server):
        respx.post(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(201, json=SAMPLE_REPLY_DRAFT)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "create_reply_draft_json",
                {"message_id": SAMPLE_MESSAGE["id"], "timezone": "America/New_York"},
            )

        assert _structured(result) == {
            "id": "AAMkAGI2draft001=",
            "web_link": "https://outlook.office.com/mail/deeplink/AAMkAGI2draft001",
            "conversation_id": "AAQkAGI2conv001=",
            "internet_message_id": "<draft001@example.com>",
            "subject": "Re: Weekly Report",
            "to": [{"name": "Alice Smith", "address": "alice@example.com"}],
            "cc": [],
        }

    @respx.mock
    async def test_create_reply_draft_retries_once_without_timezone_header(self, mcp_server):
        route = respx.post(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            side_effect=[
                httpx.Response(400, json=GRAPH_ERROR_400),
                httpx.Response(201, json=SAMPLE_REPLY_DRAFT),
            ]
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "create_reply_draft_json",
                {"message_id": SAMPLE_MESSAGE["id"], "timezone": "Bad/Zone"},
            )

        assert route.call_count == 2
        assert _structured(result)["id"] == "AAMkAGI2draft001="

    @respx.mock
    async def test_create_reply_draft_missing_web_link_is_empty(self, mcp_server):
        respx.post(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(201, json={"id": "d1"})
        )
        with _mock_token():
            result = await _call(mcp_server, "create_reply_draft_json", {"message_id": "m1"})

        assert _structured(result) == {
            "id": "d1",
            "web_link": "",
            "conversation_id": None,
            "internet_message_id": None,
            "subject": None,
            "to": [],
            "cc": [],
        }

    @respx.mock
    async def test_update_draft_body_sends_text_content_type(self, mcp_server):
        route = respx.patch(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json=SAMPLE_REPLY_DRAFT)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "update_draft_body",
                {"draft_id": "AAMkAGI2draft001=", "text": "On it."},
            )

        assert _structured(result) == {"ok": True}
        payload = json.loads(route.calls[0].request.content)
        assert payload == {"body": {"contentType": "text", "content": "On it."}}

    @respx.mock
    async def test_send_draft_returns_the_ids_the_sent_copy_carries(self, mcp_server, monkeypatch):
        """The pre-send read is the only chance to learn them, so it happens first."""
        monkeypatch.setattr("ms_graph_mcp._utcnow_iso", lambda: "2026-09-06T12:00:00Z")
        get_route = respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRAFT_FOR_SEND)
        )
        respx.post(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(202)
        )
        with _mock_token():
            result = await _call(mcp_server, "send_draft", {"draft_id": "AAMkAGI2draft001="})

        data = _structured(result)
        assert data == {
            "ok": True,
            "id": "AAMkAGI2draft001=",
            "conversation_id": "AAQkAGI2conv001=",
            "internet_message_id": "<draft001@example.com>",
            "subject": "Re: Weekly Report",
            "to": [{"name": "Alice Smith", "address": "alice@example.com"}],
            # The malformed second cc entry drops out rather than sinking the send.
            "cc": [{"name": "Charlie Brown", "address": "charlie@example.com"}],
            "sent_at": "2026-09-06T12:00:00Z",
        }
        # A sent draft's deep link is dead, so no web_link is offered at all.
        assert "web_link" not in data
        assert _select_of(get_route.calls[0].request) == mail.DRAFT_SEND_SELECT

    @respx.mock
    async def test_send_draft_reads_before_it_sends(self, mcp_server):
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRAFT_FOR_SEND)
        )
        respx.post(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(202)
        )
        with _mock_token():
            await _call(mcp_server, "send_draft", {"draft_id": "AAMkAGI2draft001="})

        methods_and_suffixes = [
            (method, path.split("/me/messages/")[-1]) for method, path in _graph_trail()
        ]
        assert methods_and_suffixes == [
            ("GET", "AAMkAGI2draft001="),
            ("POST", "AAMkAGI2draft001=/send"),
        ]

    @respx.mock
    async def test_send_draft_sends_nothing_when_the_read_fails(self, mcp_server):
        """A failed read is a tool error, and the draft is still there to retry."""
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(404, json=GRAPH_ERROR_404)
        )
        post_route = respx.post(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(202)
        )
        from fastmcp.exceptions import ToolError

        with _mock_token():
            with pytest.raises(ToolError, match="404"):
                await _call(mcp_server, "send_draft", {"draft_id": "AAMkAGI2draft001="})

        assert not post_route.called
        assert not any(method == "POST" for method, _ in _graph_trail())

    async def test_send_draft_not_connected(self, mcp_server):
        with _mock_missing_connection():
            result = await _call(mcp_server, "send_draft", {"draft_id": "d1"})

        assert _structured(result) == {"error": "not_connected", "connect_url": CONNECT_URL}


class TestMCPCreateDraftJson:
    """create_draft_json — a fresh outbound draft, not a reply."""

    URL = f"{GRAPH_BASE_URL}/me/messages"
    EXPECTED = {
        "id": "AAMkAGI2draft888=",
        "web_link": "https://outlook.office.com/mail/deeplink/AAMkAGI2draft888",
        "conversation_id": "AAQkAGI2conv888=",
        "internet_message_id": "<draft888@example.com>",
        "subject": "Lunch?",
        "to": [
            {"name": "Alice Smith", "address": "alice@example.com"},
            {"name": "Bob Jones", "address": "bob@example.com"},
        ],
        "cc": [],
    }

    @respx.mock
    async def test_posts_a_text_draft_and_returns_the_draft_shape(self, mcp_server):
        route = respx.post(self.URL).mock(return_value=httpx.Response(201, json=SAMPLE_NEW_DRAFT))
        with _mock_token():
            result = await _call(
                mcp_server,
                "create_draft_json",
                {
                    "to": "a@example.com, b@example.com",
                    "subject": "Lunch?",
                    "body": "1 < 2, so noon works.",
                    "cc": "c@example.com",
                },
            )

        assert _structured(result) == self.EXPECTED
        payload = json.loads(route.calls[0].request.content)
        # Pinned to text: a typed "<" must reach the recipient as a "<".
        assert payload["body"] == {"contentType": "Text", "content": "1 < 2, so noon works."}
        assert payload["toRecipients"] == [
            {"emailAddress": {"address": "a@example.com"}},
            {"emailAddress": {"address": "b@example.com"}},
        ]
        assert payload["ccRecipients"] == [{"emailAddress": {"address": "c@example.com"}}]
        # An absent bcc is omitted, not sent empty; and the draft is the user's own.
        assert "bccRecipients" not in payload
        assert "from" not in payload

    @respx.mock
    async def test_empty_to_still_creates_a_draft(self, mcp_server):
        """A skeleton the user finishes in Outlook through web_link."""
        route = respx.post(self.URL).mock(return_value=httpx.Response(201, json=SAMPLE_NEW_DRAFT))
        with _mock_token():
            result = await _call(mcp_server, "create_draft_json", {"to": "", "subject": "Draft it"})

        assert _structured(result) == self.EXPECTED
        payload = json.loads(route.calls[0].request.content)
        assert payload["toRecipients"] == []
        assert "ccRecipients" not in payload
        assert "bccRecipients" not in payload

    @respx.mock
    async def test_bcc_is_sent_when_given(self, mcp_server):
        route = respx.post(self.URL).mock(return_value=httpx.Response(201, json=SAMPLE_NEW_DRAFT))
        with _mock_token():
            await _call(
                mcp_server,
                "create_draft_json",
                {"to": "a@example.com", "subject": "s", "bcc": " x@example.com , "},
            )

        payload = json.loads(route.calls[0].request.content)
        assert payload["bccRecipients"] == [{"emailAddress": {"address": "x@example.com"}}]

    async def test_not_connected(self, mcp_server):
        with _mock_missing_connection():
            result = await _call(
                mcp_server, "create_draft_json", {"to": "a@example.com", "subject": "s"}
            )

        assert _structured(result) == {"error": "not_connected", "connect_url": CONNECT_URL}


class TestMCPMarkMailReadJson:
    """mark_mail_read_json."""

    @respx.mock
    async def test_marks_every_id_read(self, mcp_server):
        route = respx.patch(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "mark_mail_read_json",
                {"message_ids": json.dumps(["msg-1", "msg-2"])},
            )

        assert _structured(result) == {"updated": 2, "failed": []}
        assert route.call_count == 2
        assert str(route.calls[0].request.url).endswith("/me/messages/msg-1")
        assert str(route.calls[1].request.url).endswith("/me/messages/msg-2")
        assert json.loads(route.calls[0].request.content) == {"isRead": True}

    @respx.mock
    async def test_is_read_false_marks_unread(self, mcp_server):
        route = respx.patch(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "mark_mail_read_json",
                {"message_ids": json.dumps(["msg-1"]), "is_read": "false"},
            )

        assert _structured(result) == {"updated": 1, "failed": []}
        assert json.loads(route.calls[0].request.content) == {"isRead": False}

    @respx.mock
    async def test_one_missing_message_does_not_sink_the_batch(self, mcp_server):
        """A deleted message is reported in failed; the rest still get patched."""
        route = respx.patch(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            side_effect=[
                httpx.Response(404, json=GRAPH_ERROR_404),
                httpx.Response(200, json=SAMPLE_MESSAGE),
            ]
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "mark_mail_read_json",
                {"message_ids": json.dumps(["gone", "msg-2"])},
            )

        data = _structured(result)
        assert route.call_count == 2
        assert data["updated"] == 1
        assert len(data["failed"]) == 1
        assert data["failed"][0]["id"] == "gone"
        assert "404" in data["failed"][0]["error"]

    @respx.mock
    @pytest.mark.parametrize("bad", ["{oops", '"notanarray"', "[1, 2]"])
    async def test_malformed_message_ids_makes_no_graph_calls(self, mcp_server, bad):
        route = respx.patch(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE)
        )
        with _mock_token():
            result = await _call(mcp_server, "mark_mail_read_json", {"message_ids": bad})

        assert _structured(result) == {
            "updated": 0,
            "failed": [],
            "error": "message_ids must be a JSON array of strings",
        }
        assert route.call_count == 0

    @respx.mock
    async def test_bad_is_read_value_is_an_error_not_a_default(self, mcp_server):
        route = respx.patch(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "mark_mail_read_json",
                {"message_ids": json.dumps(["msg-1"]), "is_read": "yes"},
            )

        assert _structured(result) == {
            "updated": 0,
            "failed": [],
            "error": 'is_read must be "true" or "false"',
        }
        assert route.call_count == 0

    @respx.mock
    async def test_is_read_is_case_insensitive(self, mcp_server):
        route = respx.patch(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "mark_mail_read_json",
                {"message_ids": json.dumps(["msg-1"]), "is_read": "False"},
            )

        assert _structured(result)["updated"] == 1
        assert json.loads(route.calls[0].request.content) == {"isRead": False}

    @respx.mock
    async def test_only_the_first_100_ids_are_processed(self, mcp_server):
        route = respx.patch(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "mark_mail_read_json",
                {"message_ids": json.dumps([f"msg-{i}" for i in range(150)])},
            )

        assert _structured(result) == {"updated": 100, "failed": []}
        assert route.call_count == 100

    async def test_not_connected(self, mcp_server):
        with _mock_missing_connection():
            result = await _call(
                mcp_server, "mark_mail_read_json", {"message_ids": json.dumps(["msg-1"])}
            )

        assert _structured(result) == {"error": "not_connected", "connect_url": CONNECT_URL}


def _fake_jwt(claims: dict) -> str:
    """Build an unsigned JWT whose payload decode_token_claims can read."""
    import base64

    header = base64.urlsafe_b64encode(b'{"alg":"RS256"}').rstrip(b"=").decode()
    payload = base64.urlsafe_b64encode(json.dumps(claims).encode()).rstrip(b"=").decode()
    sig = base64.urlsafe_b64encode(b"s").rstrip(b"=").decode()
    return f"{header}.{payload}.{sig}"


IDENTITY_TOKEN = _fake_jwt({"oid": "user-obj-id", "tid": "tenant-id-123"})

SAMPLE_CHAT_MESSAGE_CREATED = {
    "id": "chat-msg-sent-002",
    "messageType": "message",
    "createdDateTime": "2026-01-06T09:00:00Z",
    "lastModifiedDateTime": "2026-01-06T09:00:00Z",
    "from": {"user": {"id": "user-id-001", "displayName": "Test User"}},
    "body": {"contentType": "text", "content": "on my way"},
}

# What Graph echoes back once the file card is on the message: the body is the
# tag, and the reference attachment carries the name the desktop renders.
SAMPLE_CHAT_MESSAGE_CREATED_WITH_FILE = {
    **SAMPLE_CHAT_MESSAGE_CREATED,
    "body": {
        "contentType": "html",
        "content": f'on my way<attachment id="{TEAMS_UPLOAD_GUID}"></attachment>',
    },
    "attachments": [
        {
            "id": TEAMS_UPLOAD_GUID,
            "contentType": "reference",
            "contentUrl": TEAMS_WEBDAV_URL,
            "name": "notes.txt",
            "thumbnailUrl": None,
        }
    ],
}


class TestMCPMarkChatReadJson:
    """mark_chat_read_json."""

    @respx.mock
    async def test_marks_the_chat_for_the_token_identity(self, mcp_server):
        route = respx.post(f"{GRAPH_BASE_URL}/chats/chat-1on1-001/markChatReadForUser").mock(
            return_value=httpx.Response(204)
        )
        with _mock_token(IDENTITY_TOKEN):
            result = await _call(mcp_server, "mark_chat_read_json", {"chat_id": "chat-1on1-001"})

        assert _structured(result) == {"ok": True}
        assert route.call_count == 1
        assert str(route.calls[0].request.url) == (
            f"{GRAPH_BASE_URL}/chats/chat-1on1-001/markChatReadForUser"
        )
        assert json.loads(route.calls[0].request.content) == {
            "user": {"id": "user-obj-id", "tenantId": "tenant-id-123"}
        }

    @respx.mock
    async def test_empty_chat_id_makes_no_graph_calls(self, mcp_server):
        route = respx.post(url__startswith=f"{GRAPH_BASE_URL}/chats/").mock(
            return_value=httpx.Response(204)
        )
        with _mock_token(IDENTITY_TOKEN):
            result = await _call(mcp_server, "mark_chat_read_json", {"chat_id": "   "})

        assert _structured(result) == {"ok": False, "error": "chat_id must not be empty"}
        assert route.call_count == 0

    @respx.mock
    async def test_token_without_claims_is_no_identity_not_a_blank_call(self, mcp_server):
        """A token we cannot decode would mark the chat for nobody — refuse it."""
        route = respx.post(url__startswith=f"{GRAPH_BASE_URL}/chats/").mock(
            return_value=httpx.Response(204)
        )
        with _mock_token("test-ms-token"):
            result = await _call(mcp_server, "mark_chat_read_json", {"chat_id": "chat-1on1-001"})

        assert _structured(result) == {"ok": False, "error": "no_identity"}
        assert route.call_count == 0

    @respx.mock
    async def test_teams_403_reports_unavailable(self, mcp_server):
        respx.post(f"{GRAPH_BASE_URL}/chats/chat-1on1-001/markChatReadForUser").mock(
            return_value=httpx.Response(403, json=GRAPH_ERROR_403)
        )
        with _mock_token(IDENTITY_TOKEN):
            result = await _call(mcp_server, "mark_chat_read_json", {"chat_id": "chat-1on1-001"})

        assert _structured(result) == {"ok": False, "error": "teams_unavailable"}

    async def test_not_connected(self, mcp_server):
        with _mock_missing_connection():
            result = await _call(mcp_server, "mark_chat_read_json", {"chat_id": "chat-1on1-001"})

        assert _structured(result) == {"error": "not_connected", "connect_url": CONNECT_URL}


class TestMCPSendChatMessageJson:
    """send_chat_message_json."""

    @respx.mock
    async def test_sends_plain_text_and_returns_the_flat_message(self, mcp_server):
        route = respx.post(f"{GRAPH_BASE_URL}/chats/chat-1on1-001/messages").mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_CREATED)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_chat_message_json",
                {"chat_id": "chat-1on1-001", "text": "on my way"},
            )

        assert route.call_count == 1
        assert json.loads(route.calls[0].request.content) == {
            "body": {"contentType": "text", "content": "on my way"}
        }
        assert _structured(result) == {
            "message": {
                "id": "chat-msg-sent-002",
                "message_type": "message",
                "from_user_id": "user-id-001",
                "from_user_display": "Test User",
                "from_application_id": None,
                "body_content": "on my way",
                "body_content_type": "text",
                "mentioned_user_ids": [],
                "created": "2026-01-06T09:00:00Z",
                "last_modified": "2026-01-06T09:00:00Z",
                "attachments": [],
            }
        }

    @respx.mock
    async def test_angle_bracket_is_sent_verbatim_not_as_markup(self, mcp_server):
        """content_type "text" — "auto" would treat a typed < as HTML."""
        route = respx.post(f"{GRAPH_BASE_URL}/chats/chat-1on1-001/messages").mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_CREATED)
        )
        with _mock_token():
            await _call(
                mcp_server,
                "send_chat_message_json",
                {"chat_id": "chat-1on1-001", "text": "a < b"},
            )

        assert json.loads(route.calls[0].request.content) == {
            "body": {"contentType": "text", "content": "a < b"}
        }

    @respx.mock
    async def test_empty_text_makes_no_graph_calls(self, mcp_server):
        route = respx.post(url__startswith=f"{GRAPH_BASE_URL}/chats/").mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_CREATED)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_chat_message_json",
                {"chat_id": "chat-1on1-001", "text": "   "},
            )

        assert _structured(result) == {"message": None, "error": "text must not be empty"}
        assert route.call_count == 0

    @respx.mock
    async def test_empty_chat_id_makes_no_graph_calls(self, mcp_server):
        route = respx.post(url__startswith=f"{GRAPH_BASE_URL}/chats/").mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_CREATED)
        )
        with _mock_token():
            result = await _call(
                mcp_server, "send_chat_message_json", {"chat_id": "", "text": "hello"}
            )

        assert _structured(result) == {"message": None, "error": "chat_id must not be empty"}
        assert route.call_count == 0

    @respx.mock
    async def test_an_attachment_is_uploaded_shared_and_comes_back_on_the_message(self, mcp_server):
        _mock_chat_file_upload()
        post = respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_CREATED_WITH_FILE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_chat_message_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "text": "on my way",
                    "attachments": json.dumps(
                        [
                            {
                                "name": "notes.txt",
                                "content_base64": base64.b64encode(b"hello").decode(),
                            }
                        ]
                    ),
                },
            )

        assert _graph_trail() == [
            ("PUT", "/v1.0/me/drive/root:/Microsoft Teams Chat Files/notes.txt:/content"),
            ("GET", "/v1.0/me/drive/items/teams-upload-001"),
            ("GET", f"/v1.0/chats/{TEAMS_CHAT_ID}/members"),
            ("POST", "/v1.0/me/drive/items/teams-upload-001/invite"),
            ("POST", f"/v1.0/chats/{TEAMS_CHAT_ID}/messages"),
        ]
        payload = json.loads(post.calls[0].request.content)
        assert payload["attachments"][0]["contentUrl"] == TEAMS_WEBDAV_URL
        attachment = _structured(result)["message"]["attachments"][0]
        assert attachment["kind"] == "file"
        assert attachment["name"] == "notes.txt"

    @respx.mock
    async def test_the_sender_from_the_token_is_not_invited(self, mcp_server):
        _mock_chat_file_upload()
        respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_CREATED_WITH_FILE)
        )
        with _mock_token(_token_with_oid("user-id-001")):
            await _call(
                mcp_server,
                "send_chat_message_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "text": "notes attached",
                    "attachments": json.dumps(
                        [
                            {
                                "name": "notes.txt",
                                "content_base64": base64.b64encode(b"hello").decode(),
                            }
                        ]
                    ),
                },
            )

        assert _invite_recipients() == [{"email": "alice@example.com"}]

    @respx.mock
    async def test_the_content_type_is_guessed_from_the_name(self, mcp_server):
        _mock_chat_file_upload()
        respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_CREATED_WITH_FILE)
        )
        with _mock_token():
            await _call(
                mcp_server,
                "send_chat_message_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "text": "on my way",
                    "attachments": json.dumps(
                        [
                            {
                                "name": "notes.txt",
                                "content_base64": base64.b64encode(b"hello").decode(),
                            }
                        ]
                    ),
                },
            )

        assert respx.calls[0].request.headers["Content-Type"] == "text/plain"

    @respx.mock
    async def test_a_file_can_travel_without_any_text(self, mcp_server):
        _mock_chat_file_upload()
        post = respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_CREATED_WITH_FILE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_chat_message_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "text": "",
                    "attachments": json.dumps(
                        [
                            {
                                "name": "notes.txt",
                                "content_base64": base64.b64encode(b"hello").decode(),
                            }
                        ]
                    ),
                },
            )

        assert json.loads(post.calls[0].request.content)["body"]["content"] == (
            f'<attachment id="{TEAMS_UPLOAD_GUID}"></attachment>'
        )
        assert _structured(result)["message"]["id"] == "chat-msg-sent-002"

    @respx.mock
    @pytest.mark.parametrize(
        "attachments,reason",
        [
            ("not json", "attachments must be a JSON array"),
            ('{"name": "a.txt"}', "attachments must be a JSON array"),
            ("[42]", "attachments[0]: not an object"),
            ('[{"content_base64": "aGk="}]', "attachments[0]: missing name"),
            ('[{"name": "  ", "content_base64": "aGk="}]', "attachments[0]: missing name"),
            ('[{"name": "a.txt"}]', "attachments[0]: missing content_base64"),
            ('[{"name": "a.txt", "content_base64": "!!!"}]', "attachments[0]: invalid base64"),
            ('[{"name": "a.txt", "content_base64": ""}]', "attachments[0]: invalid base64"),
            (
                '[{"name": "a.txt", "content_base64": "aGk=", "content_type": 7}]',
                "attachments[0]: content_type must be a string",
            ),
        ],
    )
    async def test_a_bad_attachments_payload_makes_no_graph_call(
        self, mcp_server, attachments, reason
    ):
        route = respx.route().mock(return_value=httpx.Response(200, json={}))
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_chat_message_json",
                {"chat_id": TEAMS_CHAT_ID, "text": "hi", "attachments": attachments},
            )

        assert _structured(result) == {
            "message": None,
            "error": "invalid_attachments",
            "reason": reason,
        }
        assert not route.called

    @respx.mock
    async def test_a_403_on_the_upload_is_a_permanent_scope_error(self, mcp_server):
        respx.put(url__startswith=TEAMS_SEND_UPLOAD_URL).mock(
            return_value=httpx.Response(403, json=GRAPH_ERROR_403)
        )
        post = respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_CREATED_WITH_FILE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_chat_message_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "text": "on my way",
                    "attachments": json.dumps(
                        [
                            {
                                "name": "notes.txt",
                                "content_base64": base64.b64encode(b"hello").decode(),
                            }
                        ]
                    ),
                },
            )

        assert _structured(result) == {"message": None, "error": "files_scope_missing"}
        assert not post.called

    @respx.mock
    async def test_an_empty_attachments_string_sends_exactly_as_before(self, mcp_server):
        route = respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_CREATED)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_chat_message_json",
                {"chat_id": TEAMS_CHAT_ID, "text": "on my way", "attachments": ""},
            )

        assert json.loads(route.calls[0].request.content) == {
            "body": {"contentType": "text", "content": "on my way"}
        }
        assert _structured(result)["message"]["attachments"] == []

    async def test_not_connected(self, mcp_server):
        with _mock_missing_connection():
            result = await _call(
                mcp_server,
                "send_chat_message_json",
                {"chat_id": "chat-1on1-001", "text": "hello"},
            )

        assert _structured(result) == {"error": "not_connected", "connect_url": CONNECT_URL}


class TestMCPChatsPage:
    """list_chats_page and get_chat_members_json."""

    @respx.mock
    async def test_maps_chats_and_tolerates_null_preview(self, mcp_server):
        """last_read_at comes off viewpoint, and is null on a chat without one."""
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/chats").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHATS_PAGE)
        )
        with _mock_token():
            result = await _call(mcp_server, "list_chats_page", {})

        data = _structured(result)
        assert data["next_cursor"] == SAMPLE_CHATS_PAGE_NEXT_LINK
        assert data["chats"] == [
            {
                "id": "chat-1on1-001",
                "topic": None,
                "last_preview_at": "2025-12-15T14:00:00Z",
                "last_read_at": "2025-12-15T14:00:00Z",
            },
            {
                "id": "chat-group-001",
                "topic": "Project Standup",
                "last_preview_at": "2025-12-15T13:00:00Z",
                "last_read_at": "2025-12-15T12:00:00Z",
            },
            {
                "id": "chat-empty-001",
                "topic": "Newly Created",
                "last_preview_at": None,
                "last_read_at": None,
            },
        ]

    @respx.mock
    async def test_cursor_is_requested_verbatim(self, mcp_server):
        route = respx.get(SAMPLE_CHATS_PAGE_NEXT_LINK).mock(
            return_value=httpx.Response(200, json={"value": []})
        )
        with _mock_token():
            result = await _call(
                mcp_server, "list_chats_page", {"cursor": SAMPLE_CHATS_PAGE_NEXT_LINK}
            )

        assert str(route.calls[0].request.url) == SAMPLE_CHATS_PAGE_NEXT_LINK
        assert _structured(result) == {"chats": [], "next_cursor": ""}

    @respx.mock
    async def test_chat_members_maps_user_id_and_display_name(self, mcp_server):
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/chats/").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MEMBERS_RESPONSE)
        )
        with _mock_token():
            result = await _call(mcp_server, "get_chat_members_json", {"chat_id": "chat-1on1-001"})

        assert _structured(result) == {
            "members": [
                {"user_id": "user-id-001", "display_name": "Test User"},
                {"user_id": "user-id-002", "display_name": "Alice Smith"},
            ]
        }

    async def test_not_connected(self, mcp_server):
        with _mock_missing_connection():
            result = await _call(mcp_server, "list_chats_page", {})

        assert _structured(result) == {"error": "not_connected", "connect_url": CONNECT_URL}


class TestMCPChatMessagesPage:
    """list_chat_messages_page."""

    @respx.mock
    async def test_flat_mapping_survives_null_sender_and_body(self, mcp_server):
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/chats/").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGES_PAGE)
        )
        with _mock_token():
            result = await _call(
                mcp_server, "list_chat_messages_page", {"chat_id": "chat-1on1-001"}
            )

        data = _structured(result)
        assert data["next_cursor"] == ""
        user_msg, app_msg, system_msg = data["messages"]

        assert user_msg == {
            "id": "chat-msg-001",
            "message_type": "message",
            "from_user_id": "user-id-002",
            "from_user_display": "Alice Smith",
            "from_application_id": None,
            "body_content": "<p>Sounds good!</p>",
            "body_content_type": "html",
            "mentioned_user_ids": [],
            "created": "2026-01-05T14:00:00Z",
            "last_modified": "2026-01-05T14:05:00Z",
            "attachments": [],
        }
        assert app_msg["from_application_id"] == "app-id-001"
        assert app_msg["from_user_id"] is None
        # from and body are both null on a system event.
        assert system_msg["message_type"] == "systemEventMessage"
        assert system_msg["from_user_id"] is None
        assert system_msg["from_application_id"] is None
        assert system_msg["body_content"] is None

    @respx.mock
    async def test_mentions_flatten_to_user_ids_in_order(self, mcp_server):
        msg_with_mentions = {
            **SAMPLE_CHAT_MESSAGE_FULL,
            "mentions": [
                {
                    "id": 0,
                    "mentionText": "Test User",
                    "mentioned": {"user": {"id": "user-id-001", "displayName": "Test User"}},
                },
                {
                    "id": 1,
                    "mentionText": "Bob Jones",
                    "mentioned": {"user": {"id": "user-id-003", "displayName": "Bob Jones"}},
                },
            ],
        }
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/chats/").mock(
            return_value=httpx.Response(200, json={"value": [msg_with_mentions]})
        )
        with _mock_token():
            result = await _call(
                mcp_server, "list_chat_messages_page", {"chat_id": "chat-1on1-001"}
            )

        assert _structured(result)["messages"][0]["mentioned_user_ids"] == [
            "user-id-001",
            "user-id-003",
        ]

    @respx.mock
    async def test_mentions_without_a_user_are_dropped(self, mcp_server):
        """Channel and tag mentions have no user under mentioned."""
        msg_with_channel_mention = {
            **SAMPLE_CHAT_MESSAGE_FULL,
            "mentions": [
                {
                    "id": 0,
                    "mentionText": "Everyone",
                    "mentioned": {"conversation": {"id": "channel-001", "displayName": "Everyone"}},
                },
                {
                    "id": 1,
                    "mentionText": "Test User",
                    "mentioned": {"user": {"id": "user-id-001", "displayName": "Test User"}},
                },
            ],
        }
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/chats/").mock(
            return_value=httpx.Response(200, json={"value": [msg_with_channel_mention]})
        )
        with _mock_token():
            result = await _call(
                mcp_server, "list_chat_messages_page", {"chat_id": "chat-1on1-001"}
            )

        assert _structured(result)["messages"][0]["mentioned_user_ids"] == ["user-id-001"]

    @respx.mock
    async def test_malformed_mention_entries_are_dropped(self, mcp_server):
        """A mention nests three deep and every level is type-checked. Graph
        sends none of these shapes, but one bad entry must not sink the page —
        the real id at the end is what proves the walk kept going."""
        msg_with_junk_mentions = {
            **SAMPLE_CHAT_MESSAGE_FULL,
            "mentions": [
                "not-a-dict",
                {"id": 0, "mentionText": "No mentioned key"},
                {"id": 1, "mentioned": None},
                {"id": 2, "mentioned": "not-a-dict"},
                {"id": 3, "mentioned": {"user": None}},
                {"id": 4, "mentioned": {"user": {"id": ""}}},
                {"id": 5, "mentioned": {"user": {"id": "user-id-001"}}},
            ],
        }
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/chats/").mock(
            return_value=httpx.Response(200, json={"value": [msg_with_junk_mentions]})
        )
        with _mock_token():
            result = await _call(
                mcp_server, "list_chat_messages_page", {"chat_id": "chat-1on1-001"}
            )

        assert _structured(result)["messages"][0]["mentioned_user_ids"] == ["user-id-001"]

    @respx.mock
    async def test_since_filters_on_the_orderby_property(self, mcp_server):
        route = respx.get(url__startswith=f"{GRAPH_BASE_URL}/chats/").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGES_PAGE)
        )
        with _mock_token():
            await _call(
                mcp_server,
                "list_chat_messages_page",
                {"chat_id": "chat-1on1-001", "since": "2026-01-05T00:00:00Z"},
            )

        query = parse_qs(urlparse(str(route.calls[0].request.url)).query)
        assert query["$filter"][0].split(" ")[0] == query["$orderby"][0].split(" ")[0]


class TestMCPChatMessagePageAttachments:
    """_chat_message_json's attachments list, straight off the page."""

    @respx.mock
    async def test_attachments_flatten_file_image_card_and_junk(self, mcp_server):
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/chats/").mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGES_PAGE_WITH_ATTACHMENTS)
        )
        with _mock_token():
            result = await _call(mcp_server, "list_chat_messages_page", {"chat_id": TEAMS_CHAT_ID})

        file_msg, image_msg, card_msg, junk_msg = _structured(result)["messages"]

        assert file_msg["attachments"] == [
            {
                "id": TEAMS_FILE_ATTACHMENT_ID,
                "kind": "file",
                "name": "roadmap.pptx",
                "content_type": "reference",
                "content_url": TEAMS_FILE_URL,
                "thumbnail_url": None,
                "card_text": None,
            }
        ]
        assert image_msg["attachments"] == [
            {
                "id": TEAMS_HOSTED_ID,
                "kind": "image",
                "name": None,
                "content_type": None,
                "content_url": TEAMS_HOSTED_URL,
                "thumbnail_url": None,
                "card_text": None,
            }
        ]
        assert card_msg["attachments"][0]["card_text"] == "Deploy finished"
        # The malformed entries drop out and the real file at the end survives.
        assert junk_msg["attachments"][-1]["kind"] == "file"
        assert junk_msg["attachments"][-1]["name"] == "roadmap.pptx"


class TestMCPGetChatAttachmentJson:
    """get_chat_attachment_json: bytes, thumbnails, and the permanent errors."""

    @respx.mock
    async def test_file_bytes(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        respx.get(TEAMS_SHARE_CONTENT_URL).mock(
            return_value=httpx.Response(200, content=b"PPTXBYTES")
        )
        respx.get(TEAMS_SHARE_BASE).mock(
            return_value=httpx.Response(200, json=SAMPLE_TEAMS_DRIVE_ITEM)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_chat_attachment_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                },
            )

        assert _structured(result) == {
            "kind": "file",
            "name": "roadmap.pptx",
            "content_type": TEAMS_PPTX_MIME,
            "size": len(b"PPTXBYTES"),
            "content_base64": base64.b64encode(b"PPTXBYTES").decode("ascii"),
        }

    @respx.mock
    async def test_inline_image_bytes(self, mcp_server):
        respx.get(TEAMS_IMAGE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_IMAGE)
        )
        respx.get(TEAMS_HOSTED_VALUE_URL).mock(
            return_value=httpx.Response(
                200, content=PNG_BYTES, headers={"Content-Type": "image/png"}
            )
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_chat_attachment_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "message_id": "chat-msg-image-001",
                    "attachment_id": TEAMS_HOSTED_ID,
                },
            )

        assert _structured(result) == {
            "kind": "image",
            "name": "image-aWQ9eF8wLWN1.png",
            "content_type": "image/png",
            "size": len(PNG_BYTES),
            "content_base64": base64.b64encode(PNG_BYTES).decode("ascii"),
        }

    @respx.mock
    async def test_thumbnail_skips_the_driveitem_metadata(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        respx.get(TEAMS_SHARE_THUMB_URL).mock(
            return_value=httpx.Response(
                200, content=b"THUMB", headers={"Content-Type": "image/jpeg"}
            )
        )
        meta = respx.get(TEAMS_SHARE_BASE).mock(
            return_value=httpx.Response(200, json=SAMPLE_TEAMS_DRIVE_ITEM)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_chat_attachment_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    "thumbnail": "medium",
                },
            )

        assert not meta.called
        assert _structured(result) == {
            "kind": "file",
            "name": "roadmap.pptx",
            "content_type": "image/jpeg",
            "size": 5,
            "content_base64": base64.b64encode(b"THUMB").decode("ascii"),
        }

    @respx.mock
    async def test_no_thumbnail_when_graph_has_none(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        respx.get(TEAMS_SHARE_THUMB_URL).mock(
            return_value=httpx.Response(404, json=GRAPH_ERROR_404)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_chat_attachment_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    "thumbnail": "medium",
                },
            )

        assert _structured(result) == {"error": "no_thumbnail"}

    @respx.mock
    async def test_invalid_thumbnail_is_decided_before_any_request(self, mcp_server):
        route = respx.get(url__startswith=GRAPH_BASE_URL).mock(
            return_value=httpx.Response(200, json={})
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_chat_attachment_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    "thumbnail": "enormous",
                },
            )

        assert _structured(result) == {"error": "invalid_thumbnail"}
        assert not route.called

    @respx.mock
    async def test_unknown_id_is_not_found(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_chat_attachment_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "message_id": "chat-msg-file-001",
                    "attachment_id": "nope",
                },
            )

        assert _structured(result) == {"error": "not_found"}

    @respx.mock
    async def test_a_card_is_not_found_because_it_has_no_bytes(self, mcp_server):
        respx.get(TEAMS_CARD_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_CARD)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_chat_attachment_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "message_id": "chat-msg-card-001",
                    "attachment_id": "card-att-001",
                },
            )

        assert _structured(result) == {"error": "not_found"}

    @respx.mock
    async def test_a_file_without_a_content_url_is_not_found(self, mcp_server):
        msg = {
            **SAMPLE_CHAT_MESSAGE_WITH_FILE,
            "attachments": [{"id": TEAMS_FILE_ATTACHMENT_ID, "contentType": "reference"}],
        }
        respx.get(TEAMS_FILE_MSG_URL).mock(return_value=httpx.Response(200, json=msg))
        share = respx.get(url__startswith=f"{GRAPH_BASE_URL}/shares/").mock(
            return_value=httpx.Response(200, json=SAMPLE_TEAMS_DRIVE_ITEM)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_chat_attachment_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                },
            )

        assert _structured(result) == {"error": "not_found"}
        assert not share.called

    @respx.mock
    async def test_403_on_the_sharing_link_is_access_denied(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        respx.get(TEAMS_SHARE_BASE).mock(return_value=httpx.Response(403, json=GRAPH_ERROR_403))
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_chat_attachment_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                },
            )

        assert _structured(result) == {"error": "access_denied"}

    @respx.mock
    async def test_too_large_is_decided_from_the_driveitem_size(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        content = respx.get(TEAMS_SHARE_CONTENT_URL).mock(
            return_value=httpx.Response(200, content=b"never")
        )
        respx.get(TEAMS_SHARE_BASE).mock(
            return_value=httpx.Response(200, json={**SAMPLE_TEAMS_DRIVE_ITEM, "size": 20_000_000})
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_chat_attachment_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                },
            )

        assert _structured(result) == {
            "error": "too_large",
            "size": 20_000_000,
            "limit": 10_000_000,
        }
        assert not content.called

    @respx.mock
    async def test_a_shared_folder_is_refused(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_CHAT_MESSAGE_WITH_FILE)
        )
        respx.get(TEAMS_SHARE_BASE).mock(
            return_value=httpx.Response(200, json=SAMPLE_DRIVE_ITEM_FOLDER)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_chat_attachment_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                },
            )

        assert _structured(result) == {"error": "is_folder"}

    @respx.mock
    async def test_403_on_the_message_is_teams_unavailable(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(return_value=httpx.Response(403, json=GRAPH_ERROR_403))
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_chat_attachment_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                },
            )

        assert _structured(result) == {"error": "teams_unavailable"}

    async def test_not_connected(self, mcp_server):
        with _mock_missing_connection():
            result = await _call(
                mcp_server,
                "get_chat_attachment_json",
                {
                    "chat_id": TEAMS_CHAT_ID,
                    "message_id": "chat-msg-file-001",
                    "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                },
            )

        assert _structured(result) == {"error": "not_connected", "connect_url": CONNECT_URL}

    @respx.mock
    async def test_an_unknown_message_propagates_as_a_tool_error(self, mcp_server):
        respx.get(TEAMS_FILE_MSG_URL).mock(return_value=httpx.Response(404, json=GRAPH_ERROR_404))
        from fastmcp.exceptions import ToolError

        with _mock_token():
            with pytest.raises(ToolError, match="404"):
                await _call(
                    mcp_server,
                    "get_chat_attachment_json",
                    {
                        "chat_id": TEAMS_CHAT_ID,
                        "message_id": "chat-msg-file-001",
                        "attachment_id": TEAMS_FILE_ATTACHMENT_ID,
                    },
                )


class TestMCPConnectionStatus:
    """connection_status."""

    @respx.mock
    async def test_connected_with_scopes_and_account(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me").mock(
            return_value=httpx.Response(200, json=SAMPLE_USER_PROFILE)
        )
        respx.get(f"{GRAPH_BASE_URL}/me/mailboxSettings").mock(
            return_value=httpx.Response(200, json=SAMPLE_MAILBOX_SETTINGS)
        )
        repo = MagicMock()
        repo.return_value.get_token.return_value = {
            "access_token": "a",
            "scopes": "https://graph.microsoft.com/Mail.Read Chat.Read offline_access",
        }
        with (
            _mock_token(),
            patch("auth.db.repository.TokenRepository", repo),
            patch("auth.resolve_user_key_for_request", return_value="user-key"),
        ):
            result = await _call(mcp_server, "connection_status")

        data = _structured(result)
        assert data["connected"] is True
        # Resource prefixes stripped, names lowercased.
        assert data["scopes"] == ["mail.read", "chat.read", "offline_access"]
        assert data["connect_url"] is None
        assert data["account"]["id"] == "user-id-001"
        assert data["account"]["display_name"] == "Test User"

    @respx.mock
    async def test_legacy_row_with_null_scopes_reports_empty(self, mcp_server):
        """Rows written before the scopes-key fix have nothing recorded."""
        respx.get(f"{GRAPH_BASE_URL}/me").mock(
            return_value=httpx.Response(200, json=SAMPLE_USER_PROFILE)
        )
        respx.get(f"{GRAPH_BASE_URL}/me/mailboxSettings").mock(
            return_value=httpx.Response(200, json=SAMPLE_MAILBOX_SETTINGS)
        )
        repo = MagicMock()
        repo.return_value.get_token.return_value = {"access_token": "a", "scopes": None}
        with (
            _mock_token(),
            patch("auth.db.repository.TokenRepository", repo),
            patch("auth.resolve_user_key_for_request", return_value="user-key"),
        ):
            result = await _call(mcp_server, "connection_status")

        data = _structured(result)
        assert data["connected"] is True
        assert data["scopes"] == []

    async def test_unexpected_auth_error_reports_disconnected(self, mcp_server):
        """A status probe must never surface a tool error.

        The local MSAL path can raise beyond the not-connected contract (e.g.
        AttributeError from a confidential client hitting the device-flow path
        on a stale cache); connection_status maps anything unexpected to a
        plain disconnected report instead of an isError result.
        """
        with patch("ms_graph_mcp.get_graph_token", side_effect=AttributeError("boom")):
            result = await _call(mcp_server, "connection_status")

        data = _structured(result)
        assert data["connected"] is False
        assert data["scopes"] == []
        assert data["connect_url"] is None
        assert data["account"] is None

    @respx.mock
    async def test_scope_lookup_failure_is_swallowed(self, mcp_server):
        """Laptop (MSAL) mode has no token row and no DB — status must still answer."""
        respx.get(f"{GRAPH_BASE_URL}/me").mock(
            return_value=httpx.Response(200, json=SAMPLE_USER_PROFILE)
        )
        respx.get(f"{GRAPH_BASE_URL}/me/mailboxSettings").mock(
            return_value=httpx.Response(200, json=SAMPLE_MAILBOX_SETTINGS)
        )
        repo = MagicMock(side_effect=RuntimeError("no database configured"))
        with (
            _mock_token(),
            patch("auth.db.repository.TokenRepository", repo),
            patch("auth.resolve_user_key_for_request", return_value="user-key"),
        ):
            result = await _call(mcp_server, "connection_status")

        data = _structured(result)
        assert data["connected"] is True
        assert data["scopes"] == []

    @respx.mock
    async def test_unreachable_profile_leaves_account_null(self, mcp_server):
        respx.get(f"{GRAPH_BASE_URL}/me").mock(
            return_value=httpx.Response(500, json={"error": {"code": "x", "message": "boom"}})
        )
        with (
            _mock_token(),
            patch("ms_graph_mcp._stored_graph_scopes", return_value=["mail.read"]),
        ):
            result = await _call(mcp_server, "connection_status")

        data = _structured(result)
        assert data["connected"] is True
        assert data["account"] is None
        assert data["scopes"] == ["mail.read"]

    async def test_not_connected_returns_connect_url(self, mcp_server):
        with _mock_missing_connection():
            result = await _call(mcp_server, "connection_status")

        assert _structured(result) == {
            "connected": False,
            "scopes": [],
            "connect_url": CONNECT_URL,
            "account": None,
        }

    async def test_laptop_mode_permission_error_has_no_connect_url(self, mcp_server):
        with patch("ms_graph_mcp.get_graph_token", side_effect=PermissionError("no auth")):
            result = await _call(mcp_server, "connection_status")

        data = _structured(result)
        assert data["connected"] is False
        assert data["connect_url"] is None


# ---------------------------------------------------------------------------
# Inbox rules
# ---------------------------------------------------------------------------

_RULES_URL = f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messageRules"


class TestMCPInboxRuleTools:
    """Test the manage_inbox_rules tool end-to-end via the in-process client."""

    @respx.mock
    async def test_list_rules(self, mcp_server):
        respx.get(_RULES_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE_RULES_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("manage_inbox_rules", {})

        text = _get_text(result)
        assert "2 rule(s)" in text
        assert "id|displayName|sequence|isEnabled|conditions|actions" in text
        assert "From partner" in text
        assert SAMPLE_MESSAGE_RULE["id"] in text
        # condition/action keys summarised
        assert "senderContains" in text
        assert "moveToFolder" in text

    @respx.mock
    async def test_list_rules_empty(self, mcp_server):
        respx.get(_RULES_URL).mock(return_value=httpx.Response(200, json={"value": []}))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("manage_inbox_rules", {"action": "list"})

        assert "No inbox rules found." in _get_text(result)

    @respx.mock
    async def test_get_rule(self, mcp_server):
        rule_id = SAMPLE_MESSAGE_RULE["id"]
        respx.get(f"{_RULES_URL}/{rule_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE_RULE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_inbox_rules", {"action": "get", "rule_id": rule_id}
                )

        text = _get_text(result)
        assert "From partner" in text
        assert rule_id in text
        # Full JSON includes condition values, not just keys
        assert "adele" in text
        assert "senderContains" in text

    @respx.mock
    async def test_create_rule(self, mcp_server):
        rule = {
            "displayName": "From partner",
            "sequence": 2,
            "actions": {"markAsRead": True},
        }
        route = respx.post(_RULES_URL).mock(
            return_value=httpx.Response(201, json=SAMPLE_MESSAGE_RULE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_inbox_rules",
                    {"action": "create", "options": json.dumps(rule)},
                )

        assert route.called
        payload = json.loads(route.calls[0].request.content)
        assert payload == rule
        text = _get_text(result)
        assert SAMPLE_MESSAGE_RULE["id"] in text
        assert "created" in text.lower()

    @respx.mock
    async def test_update_rule(self, mcp_server):
        rule_id = SAMPLE_MESSAGE_RULE["id"]
        route = respx.patch(f"{_RULES_URL}/{rule_id}").mock(
            return_value=httpx.Response(200, json={**SAMPLE_MESSAGE_RULE, "isEnabled": False})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_inbox_rules",
                    {
                        "action": "update",
                        "rule_id": rule_id,
                        "options": '{"isEnabled": false}',
                    },
                )

        assert route.called
        payload = json.loads(route.calls[0].request.content)
        assert payload == {"isEnabled": False}
        text = _get_text(result)
        assert rule_id in text
        assert "updated" in text.lower()

    @respx.mock
    async def test_delete_rule(self, mcp_server):
        rule_id = SAMPLE_MESSAGE_RULE["id"]
        route = respx.delete(f"{_RULES_URL}/{rule_id}").mock(return_value=httpx.Response(204))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_inbox_rules", {"action": "delete", "rule_id": rule_id}
                )

        assert route.called
        text = _get_text(result)
        assert rule_id in text
        assert "deleted" in text.lower()

    async def test_unknown_action_returns_friendly_message(self, mcp_server):
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("manage_inbox_rules", {"action": "frobnicate"})

        assert "Unknown action" in _get_text(result)

    async def test_missing_rule_id_returns_friendly_message(self, mcp_server):
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                for action in ("get", "update", "delete"):
                    result = await client.call_tool("manage_inbox_rules", {"action": action})
                    assert "rule_id" in _get_text(result)

    async def test_create_without_options_returns_friendly_message(self, mcp_server):
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("manage_inbox_rules", {"action": "create"})

        assert "options" in _get_text(result)

    async def test_update_without_options_returns_friendly_message(self, mcp_server):
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_inbox_rules",
                    {"action": "update", "rule_id": "some-id", "options": "{}"},
                )

        assert "non-empty options" in _get_text(result)

    @respx.mock
    async def test_create_missing_readwrite_surfaces_error(self, mcp_server):
        """403 (ReadWrite scope not granted) surfaces to the caller as a ToolError."""
        from fastmcp.exceptions import ToolError

        respx.post(_RULES_URL).mock(return_value=httpx.Response(403, json=GRAPH_ERROR_403))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                with pytest.raises(ToolError, match="Authorization_RequestDenied"):
                    await client.call_tool(
                        "manage_inbox_rules",
                        {
                            "action": "create",
                            "options": '{"displayName": "x", "sequence": 1, "actions": {"markAsRead": true}}',
                        },
                    )


_FOLDERS_URL = f"{GRAPH_BASE_URL}/me/mailFolders"


class TestMCPFolderTools:
    """Test the manage_mail_folders tool end-to-end via the in-process client."""

    @respx.mock
    async def test_list_folders(self, mcp_server):
        respx.get(_FOLDERS_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_MAIL_FOLDERS_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("manage_mail_folders", {})

        text = _get_text(result)
        assert "2 folder(s)" in text
        assert "id|displayName|childFolderCount|totalItemCount|unreadItemCount" in text
        assert "Projects" in text
        assert SAMPLE_MAIL_FOLDER["id"] in text

    @respx.mock
    async def test_list_folders_empty(self, mcp_server):
        respx.get(_FOLDERS_URL).mock(return_value=httpx.Response(200, json={"value": []}))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("manage_mail_folders", {"action": "list"})

        assert "No folders found." in _get_text(result)

    @respx.mock
    async def test_list_child_folders(self, mcp_server):
        parent_id = SAMPLE_MAIL_FOLDER["id"]
        route = respx.get(f"{_FOLDERS_URL}/{parent_id}/childFolders").mock(
            return_value=httpx.Response(200, json=SAMPLE_MAIL_FOLDERS_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_mail_folders",
                    {"action": "list", "options": json.dumps({"parent_id": parent_id})},
                )

        assert route.called
        assert "2 folder(s)" in _get_text(result)

    @respx.mock
    async def test_get_folder(self, mcp_server):
        folder_id = SAMPLE_MAIL_FOLDER["id"]
        respx.get(f"{_FOLDERS_URL}/{folder_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_MAIL_FOLDER)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_mail_folders", {"action": "get", "folder_id": folder_id}
                )

        text = _get_text(result)
        assert "Projects" in text
        assert folder_id in text
        assert "totalItemCount" in text

    @respx.mock
    async def test_create_folder(self, mcp_server):
        route = respx.post(_FOLDERS_URL).mock(
            return_value=httpx.Response(201, json=SAMPLE_MAIL_FOLDER)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_mail_folders",
                    {"action": "create", "options": json.dumps({"display_name": "Projects"})},
                )

        assert route.called
        payload = json.loads(route.calls[0].request.content)
        assert payload == {"displayName": "Projects"}
        text = _get_text(result)
        assert SAMPLE_MAIL_FOLDER["id"] in text
        assert "created" in text.lower()

    @respx.mock
    async def test_create_child_folder(self, mcp_server):
        parent_id = SAMPLE_MAIL_FOLDER["id"]
        route = respx.post(f"{_FOLDERS_URL}/{parent_id}/childFolders").mock(
            return_value=httpx.Response(201, json=SAMPLE_MAIL_FOLDER)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                await client.call_tool(
                    "manage_mail_folders",
                    {
                        "action": "create",
                        "options": json.dumps({"display_name": "Sub", "parent_id": parent_id}),
                    },
                )

        assert route.called
        payload = json.loads(route.calls[0].request.content)
        assert payload == {"displayName": "Sub"}

    @respx.mock
    async def test_rename_folder(self, mcp_server):
        folder_id = SAMPLE_MAIL_FOLDER["id"]
        route = respx.patch(f"{_FOLDERS_URL}/{folder_id}").mock(
            return_value=httpx.Response(200, json={**SAMPLE_MAIL_FOLDER, "displayName": "Renamed"})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_mail_folders",
                    {
                        "action": "rename",
                        "folder_id": folder_id,
                        "options": json.dumps({"display_name": "Renamed"}),
                    },
                )

        assert route.called
        payload = json.loads(route.calls[0].request.content)
        assert payload == {"displayName": "Renamed"}
        text = _get_text(result)
        assert folder_id in text
        assert "renamed" in text.lower()

    @respx.mock
    async def test_rename_folder_minimal_graph_response_returns_friendly_message(self, mcp_server):
        """Graph returns the updated folder; if it returns an empty dict the tool still succeeds."""
        folder_id = SAMPLE_MAIL_FOLDER["id"]
        respx.patch(f"{_FOLDERS_URL}/{folder_id}").mock(return_value=httpx.Response(200, json={}))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_mail_folders",
                    {
                        "action": "rename",
                        "folder_id": folder_id,
                        "options": json.dumps({"display_name": "Renamed"}),
                    },
                )

        text = _get_text(result)
        assert "renamed" in text.lower()

    @respx.mock
    async def test_delete_folder(self, mcp_server):
        folder_id = SAMPLE_MAIL_FOLDER["id"]
        route = respx.delete(f"{_FOLDERS_URL}/{folder_id}").mock(return_value=httpx.Response(204))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_mail_folders", {"action": "delete", "folder_id": folder_id}
                )

        assert route.called
        text = _get_text(result)
        assert folder_id in text
        assert "deleted" in text.lower()

    async def test_unknown_action_returns_friendly_message(self, mcp_server):
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("manage_mail_folders", {"action": "frobnicate"})

        assert "Unknown action" in _get_text(result)

    async def test_missing_folder_id_returns_friendly_message(self, mcp_server):
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                for action in ("get", "rename", "delete"):
                    result = await client.call_tool("manage_mail_folders", {"action": action})
                    assert "folder_id" in _get_text(result)

    async def test_create_without_display_name_returns_friendly_message(self, mcp_server):
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("manage_mail_folders", {"action": "create"})

        assert "display_name" in _get_text(result)

    @respx.mock
    async def test_create_missing_readwrite_surfaces_error(self, mcp_server):
        """403 (ReadWrite scope not granted) surfaces to the caller as a ToolError."""
        from fastmcp.exceptions import ToolError

        respx.post(_FOLDERS_URL).mock(return_value=httpx.Response(403, json=GRAPH_ERROR_403))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                with pytest.raises(ToolError, match="Authorization_RequestDenied"):
                    await client.call_tool(
                        "manage_mail_folders",
                        {"action": "create", "options": '{"display_name": "x"}'},
                    )

    @respx.mock
    async def test_move_folder(self, mcp_server):
        folder_id = SAMPLE_MAIL_FOLDER["id"]
        dest_id = "AQMkAGfolder-002"
        route = respx.post(f"{_FOLDERS_URL}/{folder_id}/move").mock(
            return_value=httpx.Response(200, json={**SAMPLE_MAIL_FOLDER, "parentFolderId": dest_id})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_mail_folders",
                    {
                        "action": "move",
                        "folder_id": folder_id,
                        "options": json.dumps({"destination_id": dest_id}),
                    },
                )

        assert route.called
        payload = json.loads(route.calls[0].request.content)
        assert payload == {"destinationId": dest_id}
        text = _get_text(result)
        assert folder_id in text
        assert "moved" in text.lower()

    @respx.mock
    async def test_move_folder_minimal_graph_response_returns_friendly_message(self, mcp_server):
        """Graph returns the moved folder; if it returns an empty dict the tool falls back to dest_id."""
        folder_id = SAMPLE_MAIL_FOLDER["id"]
        dest_id = "drafts"
        respx.post(f"{_FOLDERS_URL}/{folder_id}/move").mock(
            return_value=httpx.Response(200, json={})
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_mail_folders",
                    {
                        "action": "move",
                        "folder_id": folder_id,
                        "options": json.dumps({"destination_id": dest_id}),
                    },
                )

        text = _get_text(result)
        assert "moved" in text.lower()
        assert dest_id in text

    async def test_move_missing_folder_id_returns_friendly_message(self, mcp_server):
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool("manage_mail_folders", {"action": "move"})

        assert "folder_id" in _get_text(result)

    async def test_move_missing_destination_returns_friendly_message(self, mcp_server):
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_mail_folders",
                    {"action": "move", "folder_id": SAMPLE_MAIL_FOLDER["id"]},
                )

        assert "destination_id" in _get_text(result)

    # --- input coercion / hardening (commit 2) exercised through the tool -----

    @respx.mock
    async def test_list_non_string_parent_id_is_coerced_not_crashed(self, mcp_server):
        """A numeric parent_id in JSON options is coerced to a string path, not a TypeError."""
        route = respx.get(f"{_FOLDERS_URL}/123/childFolders").mock(
            return_value=httpx.Response(200, json=SAMPLE_MAIL_FOLDERS_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_mail_folders",
                    {"action": "list", "options": json.dumps({"parent_id": 123})},
                )

        assert route.called
        assert "2 folder(s)" in _get_text(result)

    @respx.mock
    async def test_create_non_string_parent_id_is_coerced_not_crashed(self, mcp_server):
        """A numeric parent_id on create is coerced to a string path, not a TypeError."""
        route = respx.post(f"{_FOLDERS_URL}/123/childFolders").mock(
            return_value=httpx.Response(201, json=SAMPLE_MAIL_FOLDER)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                result = await client.call_tool(
                    "manage_mail_folders",
                    {
                        "action": "create",
                        "options": json.dumps({"display_name": "Sub", "parent_id": 123}),
                    },
                )

        assert route.called
        assert "created" in _get_text(result).lower()

    @respx.mock
    async def test_list_top_zero_is_clamped_to_at_least_one(self, mcp_server):
        """top=0 must not reach Graph as $top=0; the tool clamps it to >= 1."""
        route = respx.get(_FOLDERS_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_MAIL_FOLDERS_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                await client.call_tool(
                    "manage_mail_folders",
                    {"action": "list", "options": json.dumps({"top": 0})},
                )

        assert route.called
        assert int(route.calls[0].request.url.params["$top"]) >= 1

    @respx.mock
    async def test_list_include_hidden_adds_query_param(self, mcp_server):
        """include_hidden=true surfaces the Graph includeHiddenFolders query param."""
        route = respx.get(_FOLDERS_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_MAIL_FOLDERS_RESPONSE)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                await client.call_tool(
                    "manage_mail_folders",
                    {"action": "list", "options": json.dumps({"include_hidden": True})},
                )

        assert route.called
        assert route.calls[0].request.url.params["includeHiddenFolders"] == "true"

    @respx.mock
    async def test_list_surfaces_graph_error(self, mcp_server):
        """A 403 on list surfaces to the caller as a ToolError."""
        from fastmcp.exceptions import ToolError

        respx.get(_FOLDERS_URL).mock(return_value=httpx.Response(403, json=GRAPH_ERROR_403))
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                with pytest.raises(ToolError, match="Authorization_RequestDenied"):
                    await client.call_tool("manage_mail_folders", {"action": "list"})

    @respx.mock
    async def test_move_surfaces_graph_error(self, mcp_server):
        """A 404 on move surfaces to the caller as a ToolError."""
        from fastmcp.exceptions import ToolError

        folder_id = SAMPLE_MAIL_FOLDER["id"]
        respx.post(f"{_FOLDERS_URL}/{folder_id}/move").mock(
            return_value=httpx.Response(404, json=GRAPH_ERROR_404)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                with pytest.raises(ToolError, match="ResourceNotFound"):
                    await client.call_tool(
                        "manage_mail_folders",
                        {
                            "action": "move",
                            "folder_id": folder_id,
                            "options": json.dumps({"destination_id": "inbox"}),
                        },
                    )

    @respx.mock
    async def test_delete_surfaces_graph_error(self, mcp_server):
        """A 404 on delete surfaces to the caller as a ToolError."""
        from fastmcp.exceptions import ToolError

        respx.delete(f"{_FOLDERS_URL}/bad-id").mock(
            return_value=httpx.Response(404, json=GRAPH_ERROR_404)
        )
        with _mock_token():
            from fastmcp import Client

            async with Client(mcp_server) as client:
                with pytest.raises(ToolError, match="ResourceNotFound"):
                    await client.call_tool(
                        "manage_mail_folders", {"action": "delete", "folder_id": "bad-id"}
                    )


class TestConnectFlowScopes:
    """The /connect flow and local MSAL login share ONE scope policy.

    Requesting an admin-gated scope walls the whole consent bundle behind
    "Approval required", so the org default must be exactly the consented set;
    the connect flow only adds offline_access, which the code grant needs for
    its refresh token.
    """

    def test_org_connect_request_is_consented_set_plus_offline_access(self):
        import os
        from unittest.mock import patch as env_patch

        from ms_graph.local_auth import CONSENTED_ORG_SCOPES
        from ms_graph_mcp import _ms_graph_scopes

        with env_patch.dict(os.environ, {"MS_TENANT_ID": "t"}, clear=True):
            scopes = _ms_graph_scopes().split()
        assert scopes == [*CONSENTED_ORG_SCOPES, "offline_access"]

    def test_ms_scopes_override_reaches_the_connect_flow(self):
        import os
        from unittest.mock import patch as env_patch

        from ms_graph_mcp import _ms_graph_scopes

        with env_patch.dict(
            os.environ,
            {"MS_TENANT_ID": "t", "MS_SCOPES": "Mail.Read Chat.ReadWrite"},
            clear=True,
        ):
            scopes = _ms_graph_scopes().split()
        assert scopes == ["Mail.Read", "Chat.ReadWrite", "offline_access"]

    def test_offline_access_is_not_duplicated(self):
        import os
        from unittest.mock import patch as env_patch

        from ms_graph_mcp import _ms_graph_scopes

        with env_patch.dict(
            os.environ,
            {"MS_TENANT_ID": "t", "MS_SCOPES": "Mail.Read offline_access"},
            clear=True,
        ):
            assert _ms_graph_scopes().split() == ["Mail.Read", "offline_access"]

    def test_connect_config_resolves_scopes_at_request_time(self):
        """The config must carry the policy FUNCTION, not an import-time call.

        MS_TENANT_ID/MS_SCOPES can load after module import (same reason the
        config's authorize_url/token_url are lambdas over _ms_tenant); a
        captured string would hand an org tenant the consumer wish-list and
        rebuild the "Approval required" wall.
        """
        import os
        from unittest.mock import patch as env_patch

        from ms_graph.local_auth import CONSENTED_ORG_SCOPES
        from ms_graph_mcp import MICROSOFT_CONNECT_CONFIG

        assert callable(MICROSOFT_CONNECT_CONFIG.scopes)
        with env_patch.dict(os.environ, {"MS_TENANT_ID": "t"}, clear=True):
            org = MICROSOFT_CONNECT_CONFIG.resolved_scopes().split()
        with env_patch.dict(os.environ, {"MS_SCOPES": "Mail.Read"}, clear=True):
            overridden = MICROSOFT_CONNECT_CONFIG.resolved_scopes().split()
        assert org == [*CONSENTED_ORG_SCOPES, "offline_access"]
        assert overridden == ["Mail.Read", "offline_access"]


# ---------------------------------------------------------------------------
# External-sender mail policy
# ---------------------------------------------------------------------------

# The sender check reads the parent message by id, percent-encoded, with only
# id/from/sender selected.
SENDER_CHECK_URL = f"{GRAPH_BASE_URL}/me/messages/{quote(ATT_MSG_ID, safe='')}"
SENDER_CHECK_PATH = f"/v1.0/me/messages/{ATT_MSG_ID}"
ATT_EXT_ITEM_URL = f"{ATT_BASE}/{quote(SAMPLE_EXTERNAL_ITEM_ATTACHMENT['id'], safe='')}"
EXTERNAL_MSG_ID = SAMPLE_EXTERNAL_MESSAGE["id"]

# Every mail-touching tool must be classified. Decision 12: the ungated ones
# are writes on ids the caller must already hold (the gated surfaces never hand
# out an external id), draft ids that Exchange rejects on non-drafts, and
# folder metadata that carries no mail.
GATED_MAIL_TOOLS = frozenset(
    {
        "list_emails",
        "read_email",
        "get_email_attachment",
        "send_email",
        "manage_inbox_rules",
        "list_mail_delta",
        "get_mail_detail",
        "get_mail_attachment_json",
        "create_reply_draft_json",
    }
)
DELIBERATELY_UNGATED_MAIL_TOOLS = frozenset(
    {
        "mark_mail_read_json",
        "update_draft_body",
        "send_draft",
        "add_draft_attachment_json",
        "create_draft_json",
        "manage_mail_folders",
    }
)
MAIL_TOOL_WORDS = ("mail", "email", "draft", "inbox")


def _policy_on(monkeypatch, value: str = "example.com") -> None:
    """Turn the policy on. Every sample sender is @example.com, i.e. internal."""
    monkeypatch.setenv(mail_policy.ENV_ALLOWED_SENDER_DOMAINS, value)


def _assert_no_canary(result) -> None:
    """No field of a hidden message may appear in what a tool returned.

    Each conftest canary fixture carries a distinct CANARY-<FIELD> marker, so a
    failure here names the field that leaked.
    """
    text = json.dumps(result, default=str) if isinstance(result, dict) else str(result)
    for marker in ("CANARY-", "mallory", "Mallory"):
        assert marker not in text, f"{marker!r} leaked into a policy-gated result: {text[:400]}"


def _select_of(request) -> str:
    """The $select query parameter of a recorded request."""
    return parse_qs(urlparse(str(request.url)).query).get("$select", [""])[0]


def _sender_checks() -> list[str]:
    """Every request respx saw that was a policy sender check, by path.

    The policy's one probe is identifiable by its $select alone, so this says
    "the policy looked" without depending on which message id it looked at.
    """
    return [
        c.request.url.path
        for c in respx.calls
        if _select_of(c.request) == mail_policy.SENDER_SELECT
    ]


class TestMailSenderPolicy:
    """Every read surface hides mail that arrived from outside the allowlist."""

    # -- list_emails --------------------------------------------------------

    @respx.mock
    async def test_list_emails_hides_external_and_appends_the_notice(self, mcp_server, monkeypatch):
        _policy_on(monkeypatch)
        route = respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages").mock(
            return_value=httpx.Response(
                200, json={"value": [SAMPLE_MESSAGE, SAMPLE_EXTERNAL_MESSAGE]}
            )
        )
        with _mock_token():
            result = await _call(mcp_server, "list_emails", {"top": 10})

        text = _get_text(result)
        assert "1 message(s) in inbox" in text
        assert "alice@example.com" in text
        assert text.endswith(mail_policy.POLICY_NOTICE)
        assert "sender" in _select_of(route.calls[0].request)
        _assert_no_canary(text)

    @respx.mock
    async def test_list_emails_select_asks_for_sender_even_when_off(self, mcp_server):
        """The select is static, so a stored cursor never lacks the field."""
        route = respx.get(f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGES_RESPONSE)
        )
        with _mock_token():
            result = await _call(mcp_server, "list_emails", {})

        select = _select_of(route.calls[0].request)
        assert "from" in select and "sender" in select
        assert mail_policy.POLICY_NOTICE not in _get_text(result)
        _assert_no_canary(_get_text(result))

    @respx.mock
    async def test_list_emails_search_hiding_everything_still_reports_no_results(
        self, mcp_server, monkeypatch
    ):
        _policy_on(monkeypatch)
        respx.get(f"{GRAPH_BASE_URL}/me/messages").mock(
            return_value=httpx.Response(200, json={"value": [SAMPLE_EXTERNAL_MESSAGE]})
        )
        with _mock_token():
            result = await _call(mcp_server, "list_emails", {"query": "CANARY"})

        text = _get_text(result)
        assert text == f'No messages found matching "CANARY".\n{mail_policy.POLICY_NOTICE}'
        _assert_no_canary(text)

    @respx.mock
    async def test_search_notice_is_identical_for_zero_and_many_hidden(
        self, mcp_server, monkeypatch
    ):
        """A hidden count on a $search query would be a content oracle."""
        _policy_on(monkeypatch)
        three = [{**SAMPLE_EXTERNAL_MESSAGE, "id": f"ext-{i}"} for i in range(3)]
        respx.get(f"{GRAPH_BASE_URL}/me/messages").mock(
            side_effect=[
                httpx.Response(200, json={"value": [SAMPLE_EXTERNAL_MESSAGE]}),
                httpx.Response(200, json={"value": three}),
            ]
        )
        with _mock_token():
            one_hidden = _get_text(await _call(mcp_server, "list_emails", {"query": "CANARY"}))
            many_hidden = _get_text(await _call(mcp_server, "list_emails", {"query": "CANARY"}))

        assert one_hidden == many_hidden
        _assert_no_canary(one_hidden)

    @respx.mock
    async def test_list_emails_filters_a_shared_mailbox_too(self, mcp_server, monkeypatch):
        _policy_on(monkeypatch)
        respx.get(f"{GRAPH_BASE_URL}/users/support@example.com/mailFolders/inbox/messages").mock(
            return_value=httpx.Response(
                200, json={"value": [SAMPLE_MESSAGE, SAMPLE_EXTERNAL_MESSAGE]}
            )
        )
        with _mock_token():
            result = await _call(
                mcp_server, "list_emails", {"mailbox": "support@example.com", "top": 10}
            )

        text = _get_text(result)
        assert "1 message(s) in inbox" in text
        assert text.endswith(mail_policy.POLICY_NOTICE)
        _assert_no_canary(text)

    # -- read_email ---------------------------------------------------------

    @respx.mock
    async def test_read_email_refuses_external_before_marking_or_listing(
        self, mcp_server, monkeypatch
    ):
        _policy_on(monkeypatch)
        respx.get(f"{GRAPH_BASE_URL}/me/messages/{EXTERNAL_MSG_ID}").mock(
            return_value=httpx.Response(200, json=SAMPLE_EXTERNAL_MESSAGE)
        )
        patch_route = respx.patch(f"{GRAPH_BASE_URL}/me/messages/{EXTERNAL_MSG_ID}").mock(
            return_value=httpx.Response(200, json=SAMPLE_EXTERNAL_MESSAGE)
        )
        attachments_route = respx.get(
            url__startswith=f"{GRAPH_BASE_URL}/me/messages/{EXTERNAL_MSG_ID}/attachments"
        ).mock(return_value=httpx.Response(200, json=SAMPLE_ATTACHMENTS_RESPONSE))

        with _mock_token():
            result = await _call(
                mcp_server,
                "read_email",
                {"message_id": EXTERNAL_MSG_ID, "options": '{"mark_as_read": true}'},
            )

        text = _get_text(result)
        assert text == mail_policy.EXTERNAL_SENDER_TEXT
        assert not patch_route.called
        assert not attachments_route.called
        assert _graph_trail() == [("GET", f"/v1.0/me/messages/{EXTERNAL_MSG_ID}")]
        _assert_no_canary(text)

    @respx.mock
    async def test_read_email_refuses_mail_sent_on_behalf_of_an_insider(
        self, mcp_server, monkeypatch
    ):
        """from is internal, but an outside service pressed send."""
        _policy_on(monkeypatch)
        msg_id = SAMPLE_ONBEHALF_MESSAGE["id"]
        respx.get(f"{GRAPH_BASE_URL}/me/messages/{msg_id}").mock(
            return_value=httpx.Response(200, json=SAMPLE_ONBEHALF_MESSAGE)
        )
        with _mock_token():
            result = await _call(mcp_server, "read_email", {"message_id": msg_id})

        text = _get_text(result)
        assert text == mail_policy.EXTERNAL_SENDER_TEXT
        _assert_no_canary(text)

    @respx.mock
    async def test_read_email_internal_message_is_unchanged(self, mcp_server, monkeypatch):
        _policy_on(monkeypatch)
        respx.get(f"{GRAPH_BASE_URL}/me/messages/{ATT_MSG_ID}").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE)
        )
        with _mock_token():
            result = await _call(mcp_server, "read_email", {"message_id": ATT_MSG_ID})

        text = _get_text(result)
        assert "Weekly Report" in text
        assert "Here is the weekly report" in text
        _assert_no_canary(text)

    # -- get_email_attachment ----------------------------------------------

    @respx.mock
    async def test_get_email_attachment_refuses_before_touching_the_attachment(
        self, mcp_server, monkeypatch
    ):
        _policy_on(monkeypatch)
        value_route = respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(200, content=b"x")
        )
        meta_route = respx.get(ATT_FILE_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_FILE_ATTACHMENT)
        )
        respx.get(SENDER_CHECK_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_SENDER_ONLY_EXTERNAL)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_FILE_ATTACHMENT["id"]},
            )

        text = _get_text(result)
        assert text == mail_policy.EXTERNAL_SENDER_TEXT
        assert not meta_route.called
        assert not value_route.called
        assert _graph_trail() == [("GET", SENDER_CHECK_PATH)]
        _assert_no_canary(text)

    @respx.mock
    async def test_get_email_attachment_checks_the_mailbox_it_will_read(
        self, mcp_server, monkeypatch
    ):
        """A /me check before a /users/{mailbox} read would be the wrong check."""
        _policy_on(monkeypatch)
        route = respx.get(url__startswith=f"{GRAPH_BASE_URL}/users/support@example.com/").mock(
            return_value=httpx.Response(200, json=SAMPLE_SENDER_ONLY_EXTERNAL)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_FILE_ATTACHMENT["id"],
                    "mailbox": "support@example.com",
                },
            )

        assert _get_text(result) == mail_policy.EXTERNAL_SENDER_TEXT
        assert route.calls[0].request.url.path.startswith(
            "/v1.0/users/support@example.com/messages/"
        )
        _assert_no_canary(_get_text(result))

    @respx.mock
    @pytest.mark.parametrize("mode", ["text", "base64", "onedrive"])
    async def test_get_email_attachment_refuses_an_external_attached_message(
        self, mcp_server, monkeypatch, mode
    ):
        """An attached message is judged by the same rule as a message."""
        _policy_on(monkeypatch)
        value_route = respx.get(f"{ATT_EXT_ITEM_URL}/$value").mock(
            return_value=httpx.Response(200, content=b"raw-eml")
        )

        def _respond(request):
            if "expand" in str(request.url):
                return httpx.Response(200, json=SAMPLE_EXTERNAL_ITEM_ATTACHMENT)
            return httpx.Response(200, json=SAMPLE_EXTERNAL_ITEM_ATTACHMENT_META)

        respx.get(url__startswith=ATT_EXT_ITEM_URL).mock(side_effect=_respond)
        respx.get(SENDER_CHECK_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_SENDER_ONLY_INTERNAL)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_EXTERNAL_ITEM_ATTACHMENT["id"],
                    "mode": mode,
                },
            )

        text = _get_text(result)
        assert text == mail_policy.EXTERNAL_SENDER_TEXT
        assert not value_route.called
        _assert_no_canary(text)

    @respx.mock
    async def test_get_email_attachment_internal_costs_one_extra_request(
        self, mcp_server, monkeypatch
    ):
        _policy_on(monkeypatch)
        respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(
                200, content=b"hello there", headers={"Content-Type": "text/plain"}
            )
        )
        respx.get(ATT_FILE_URL).mock(
            return_value=httpx.Response(
                200,
                json={
                    **SAMPLE_FILE_ATTACHMENT,
                    "name": "notes.txt",
                    "contentType": "text/plain",
                    "size": 11,
                },
            )
        )
        respx.get(SENDER_CHECK_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_SENDER_ONLY_INTERNAL)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_FILE_ATTACHMENT["id"]},
            )

        text = _get_text(result)
        assert "hello there" in text
        trail = _graph_trail()
        assert trail[0] == ("GET", SENDER_CHECK_PATH)
        assert len(trail) == 3
        _assert_no_canary(text)

    @respx.mock
    async def test_policy_off_issues_no_extra_requests(self, mcp_server):
        """Off must be byte-for-byte the old behaviour, including request count."""
        respx.get(f"{ATT_FILE_URL}/$value").mock(
            return_value=httpx.Response(
                200, content=b"hello there", headers={"Content-Type": "text/plain"}
            )
        )
        respx.get(ATT_FILE_URL).mock(
            return_value=httpx.Response(
                200,
                json={
                    **SAMPLE_FILE_ATTACHMENT,
                    "name": "notes.txt",
                    "contentType": "text/plain",
                    "size": 11,
                },
            )
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_email_attachment",
                {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_FILE_ATTACHMENT["id"]},
            )

        assert "hello there" in _get_text(result)
        assert len(_graph_trail()) == 2
        _assert_no_canary(_get_text(result))

    # -- forwarding an attachment out of an external message ----------------

    @respx.mock
    async def test_send_email_refuses_to_forward_an_external_attachment(
        self, mcp_server, monkeypatch
    ):
        _policy_on(monkeypatch)
        attachments_route = respx.get(url__startswith=ATT_BASE).mock(
            return_value=httpx.Response(200, json=SAMPLE_FILE_ATTACHMENT)
        )
        send_route = respx.post(f"{GRAPH_BASE_URL}/me/sendMail").mock(
            return_value=httpx.Response(202)
        )
        respx.get(SENDER_CHECK_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_SENDER_ONLY_EXTERNAL)
        )
        spec = {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_FILE_ATTACHMENT["id"]}
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_email",
                {
                    "to": "bob@example.com",
                    "subject": "fwd",
                    "body": "see attached",
                    "options": json.dumps({"attachments": [spec]}),
                },
            )

        text = _get_text(result)
        assert text == f"attachments[0]: {mail_policy.EXTERNAL_SENDER_TEXT}"
        assert not attachments_route.called
        assert not send_route.called
        assert all(method != "POST" for method, _ in _graph_trail())
        _assert_no_canary(text)

    @respx.mock
    async def test_forward_spec_mailbox_is_the_mailbox_that_is_checked(
        self, mcp_server, monkeypatch
    ):
        _policy_on(monkeypatch)
        route = respx.get(url__startswith=f"{GRAPH_BASE_URL}/users/support@example.com/").mock(
            return_value=httpx.Response(200, json=SAMPLE_SENDER_ONLY_EXTERNAL)
        )
        spec = {
            "message_id": ATT_MSG_ID,
            "attachment_id": SAMPLE_FILE_ATTACHMENT["id"],
            "mailbox": "support@example.com",
        }
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_email",
                {
                    "to": "bob@example.com",
                    "subject": "fwd",
                    "body": "see attached",
                    "options": json.dumps({"attachments": [spec]}),
                },
            )

        assert _get_text(result) == f"attachments[0]: {mail_policy.EXTERNAL_SENDER_TEXT}"
        assert route.calls[0].request.url.path.startswith(
            "/v1.0/users/support@example.com/messages/"
        )
        _assert_no_canary(_get_text(result))

    @respx.mock
    async def test_send_teams_message_refuses_to_forward_an_external_attachment(
        self, mcp_server, monkeypatch
    ):
        """Teams is out of scope, but it must not become a laundering route."""
        _policy_on(monkeypatch)
        post_route = respx.post(TEAMS_CHAT_MSGS).mock(
            return_value=httpx.Response(201, json=SAMPLE_CHAT_MESSAGE_SENT)
        )
        respx.get(SENDER_CHECK_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_SENDER_ONLY_EXTERNAL)
        )
        spec = {"message_id": ATT_MSG_ID, "attachment_id": SAMPLE_FILE_ATTACHMENT["id"]}
        with _mock_token():
            result = await _call(
                mcp_server,
                "send_teams_message",
                {
                    "message": "fyi",
                    "chat_id": TEAMS_CHAT_ID,
                    "options": json.dumps({"attachments": [spec]}),
                },
            )

        text = _get_text(result)
        assert text == f"attachments[0]: {mail_policy.EXTERNAL_SENDER_TEXT}"
        assert not post_route.called
        _assert_no_canary(text)

    # -- Desktop JSON surfaces ---------------------------------------------

    @respx.mock
    async def test_list_mail_delta_hides_external_and_keeps_tombstones(
        self, mcp_server, monkeypatch
    ):
        _policy_on(monkeypatch)
        page = {
            "@odata.nextLink": SAMPLE_DELTA_NEXT_LINK,
            "value": [
                SAMPLE_DELTA_MESSAGE,
                SAMPLE_EXTERNAL_DELTA_MESSAGE,
                SAMPLE_DELTA_TOMBSTONE,
            ],
        }
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages/delta").mock(
            return_value=httpx.Response(200, json=page)
        )
        with _mock_token():
            result = await _call(mcp_server, "list_mail_delta", {})

        data = _structured(result)
        assert data["messages"] == [SAMPLE_DELTA_MESSAGE, SAMPLE_DELTA_TOMBSTONE]
        assert set(data) == {"messages", "next_cursor", "delta_cursor", "resync"}
        assert data["next_cursor"] == SAMPLE_DELTA_NEXT_LINK
        assert data["delta_cursor"] == ""
        assert data["resync"] is False
        _assert_no_canary(data)

    @respx.mock
    async def test_get_mail_detail_returns_only_the_error(self, mcp_server, monkeypatch):
        _policy_on(monkeypatch)
        route = respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json=SAMPLE_EXTERNAL_MESSAGE_DETAIL)
        )
        with _mock_token():
            result = await _call(
                mcp_server, "get_mail_detail", {"message_id": SAMPLE_EXTERNAL_MESSAGE["id"]}
            )

        data = _structured(result)
        assert data == {"error": mail_policy.EXTERNAL_SENDER_ERROR}
        select = _select_of(route.calls[0].request).split(",")
        assert "isDraft" in select
        assert "from" in select and "sender" in select
        _assert_no_canary(data)

    @respx.mock
    @pytest.mark.parametrize("mode", ["metadata", "text", "bytes"])
    async def test_get_mail_attachment_json_refuses_an_external_parent(
        self, mcp_server, monkeypatch, mode
    ):
        _policy_on(monkeypatch)
        attachments_route = respx.get(url__startswith=ATT_BASE).mock(
            return_value=httpx.Response(200, json=SAMPLE_FILE_ATTACHMENT)
        )
        respx.get(SENDER_CHECK_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_SENDER_ONLY_EXTERNAL)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_FILE_ATTACHMENT["id"],
                    "mode": mode,
                },
            )

        data = _structured(result)
        assert data == {"error": mail_policy.EXTERNAL_SENDER_ERROR}
        assert not attachments_route.called
        _assert_no_canary(data)

    @respx.mock
    @pytest.mark.parametrize("mode", ["metadata", "text", "bytes"])
    async def test_get_mail_attachment_json_refuses_an_external_attached_message(
        self, mcp_server, monkeypatch, mode
    ):
        _policy_on(monkeypatch)
        value_route = respx.get(f"{ATT_EXT_ITEM_URL}/$value").mock(
            return_value=httpx.Response(200, content=b"raw-eml")
        )

        def _respond(request):
            if "expand" in str(request.url):
                return httpx.Response(200, json=SAMPLE_EXTERNAL_ITEM_ATTACHMENT)
            return httpx.Response(200, json=SAMPLE_EXTERNAL_ITEM_ATTACHMENT_META)

        respx.get(url__startswith=ATT_EXT_ITEM_URL).mock(side_effect=_respond)
        respx.get(SENDER_CHECK_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_SENDER_ONLY_INTERNAL)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {
                    "message_id": ATT_MSG_ID,
                    "attachment_id": SAMPLE_EXTERNAL_ITEM_ATTACHMENT["id"],
                    "mode": mode,
                },
            )

        data = _structured(result)
        assert data == {"error": mail_policy.EXTERNAL_SENDER_ERROR}
        assert not value_route.called
        _assert_no_canary(data)

    @respx.mock
    async def test_create_reply_draft_json_never_posts_create_reply(self, mcp_server, monkeypatch):
        """Graph would build a draft quoting the original, from = the user."""
        _policy_on(monkeypatch)
        post_route = respx.post(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(201, json=SAMPLE_REPLY_DRAFT)
        )
        respx.get(SENDER_CHECK_URL).mock(
            return_value=httpx.Response(200, json=SAMPLE_SENDER_ONLY_EXTERNAL)
        )
        with _mock_token():
            result = await _call(mcp_server, "create_reply_draft_json", {"message_id": ATT_MSG_ID})

        data = _structured(result)
        assert data == {"error": mail_policy.EXTERNAL_SENDER_ERROR}
        assert not post_route.called
        assert _graph_trail() == [("GET", SENDER_CHECK_PATH)]
        _assert_no_canary(data)

    # -- outbound stays ungated --------------------------------------------

    @respx.mock
    async def test_create_draft_json_is_not_gated(self, mcp_server, monkeypatch):
        """Composing the user's own mail reads nothing that arrived from anyone."""
        _policy_on(monkeypatch)
        respx.post(f"{GRAPH_BASE_URL}/me/messages").mock(
            return_value=httpx.Response(201, json=SAMPLE_NEW_DRAFT)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "create_draft_json",
                {"to": "a@example.com", "subject": "Lunch?", "body": "noon"},
            )

        data = _structured(result)
        assert data["id"] == SAMPLE_NEW_DRAFT["id"]
        assert not _sender_checks()

    @respx.mock
    async def test_send_draft_is_not_gated(self, mcp_server, monkeypatch):
        """The pre-send read selects no from/sender, so there is nothing to gate."""
        _policy_on(monkeypatch)
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json=SAMPLE_DRAFT_FOR_SEND)
        )
        respx.post(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(202)
        )
        with _mock_token():
            result = await _call(mcp_server, "send_draft", {"draft_id": "AAMkAGI2draft001="})

        assert _structured(result)["ok"] is True
        assert not _sender_checks()

    # -- forwarding inbox rules --------------------------------------------

    @respx.mock
    async def test_forwarding_rule_is_refused_before_any_request(self, mcp_server, monkeypatch):
        _policy_on(monkeypatch)
        with _mock_token():
            result = await _call(
                mcp_server,
                "manage_inbox_rules",
                {"action": "create", "options": json.dumps(SAMPLE_FORWARDING_RULE)},
            )

        text = _get_text(result)
        assert text == mail_policy.FORWARDING_RULE_TEXT
        assert _graph_trail() == []
        _assert_no_canary(text)

    @respx.mock
    async def test_updating_a_rule_to_redirect_is_refused(self, mcp_server, monkeypatch):
        _policy_on(monkeypatch)
        changes = {"actions": {"redirectTo": [{"emailAddress": {"address": "x@evil.example.net"}}]}}
        with _mock_token():
            result = await _call(
                mcp_server,
                "manage_inbox_rules",
                {
                    "action": "update",
                    "rule_id": SAMPLE_MESSAGE_RULE["id"],
                    "options": json.dumps(changes),
                },
            )

        text = _get_text(result)
        assert text == mail_policy.FORWARDING_RULE_TEXT
        assert _graph_trail() == []
        _assert_no_canary(text)

    @respx.mock
    async def test_non_forwarding_rules_still_work_while_the_policy_is_on(
        self, mcp_server, monkeypatch
    ):
        """Move/copy/markAsRead rules do not change from, so they stay allowed."""
        _policy_on(monkeypatch)
        rule = {"displayName": "Filed", "sequence": 2, "actions": {"moveToFolder": "AQMkAG"}}
        route = respx.post(_RULES_URL).mock(
            return_value=httpx.Response(201, json=SAMPLE_MESSAGE_RULE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "manage_inbox_rules",
                {"action": "create", "options": json.dumps(rule)},
            )

        assert route.called
        assert "created" in _get_text(result).lower()
        _assert_no_canary(_get_text(result))

    @respx.mock
    async def test_forwarding_rule_is_allowed_while_the_policy_is_off(self, mcp_server):
        """The gate closes the policy bypass; it is not an independent control."""
        route = respx.post(_RULES_URL).mock(
            return_value=httpx.Response(201, json=SAMPLE_MESSAGE_RULE)
        )
        with _mock_token():
            result = await _call(
                mcp_server,
                "manage_inbox_rules",
                {"action": "create", "options": json.dumps(SAMPLE_FORWARDING_RULE)},
            )

        assert route.called
        assert "created" in _get_text(result).lower()
        _assert_no_canary(_get_text(result))

    # -- connection_status --------------------------------------------------

    @respx.mock
    @pytest.mark.parametrize("value, expected", [("example.com", True), (None, False)])
    async def test_connection_status_reports_the_policy_state(
        self, mcp_server, monkeypatch, value, expected
    ):
        if value:
            _policy_on(monkeypatch, value)
        respx.get(f"{GRAPH_BASE_URL}/me").mock(
            return_value=httpx.Response(500, json={"error": {"code": "x", "message": "boom"}})
        )
        with (
            _mock_token(),
            patch("ms_graph_mcp._stored_graph_scopes", return_value=[]),
        ):
            result = await _call(mcp_server, "connection_status")

        data = _structured(result)
        assert data["connected"] is True
        assert data["mail_policy"] == {"enabled": expected}
        _assert_no_canary(data)

    @respx.mock
    async def test_connection_status_never_crashes_on_a_bad_allowlist(self, monkeypatch):
        """A status probe must answer even when the config would stop the pod.

        Called directly rather than through the client, because the lifespan
        refuses to start the server at all on a malformed allowlist.
        """
        monkeypatch.setenv(mail_policy.ENV_ALLOWED_SENDER_DOMAINS, "not a domain")
        respx.get(f"{GRAPH_BASE_URL}/me").mock(
            return_value=httpx.Response(500, json={"error": {"code": "x", "message": "boom"}})
        )
        from ms_graph_mcp import connection_status

        with (
            _mock_token(),
            patch("ms_graph_mcp._stored_graph_scopes", return_value=[]),
        ):
            data = await connection_status()

        assert data["connected"] is True
        assert data["mail_policy"] == {"enabled": True, "error": "invalid_config"}
        _assert_no_canary(data)

    # -- misconfiguration fails closed --------------------------------------

    @respx.mock
    async def test_bad_config_fails_read_email_before_any_graph_call(self, monkeypatch):
        """The tool itself fails closed, independently of the boot check.

        Called directly rather than through the client: the in-process client
        runs the lifespan, which would refuse the allowlist first and prove
        nothing about the tool.
        """
        _policy_on(monkeypatch, "*")
        respx.get(f"{GRAPH_BASE_URL}/me/messages/{ATT_MSG_ID}").mock(
            return_value=httpx.Response(200, json=SAMPLE_MESSAGE)
        )
        from ms_graph_mcp import read_email

        with _mock_token(), pytest.raises(mail_policy.MailPolicyConfigError, match=r"'\*'"):
            await read_email(ATT_MSG_ID)

        assert _graph_trail() == []

    @respx.mock
    async def test_bad_config_is_not_mistaken_for_a_missing_connection(self, monkeypatch):
        """A Desktop JSON tool raises on a bad allowlist instead of returning a dict."""
        _policy_on(monkeypatch, "*")
        respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/messages/").mock(
            return_value=httpx.Response(200, json=SAMPLE_EXTERNAL_MESSAGE_DETAIL)
        )
        from ms_graph_mcp import get_mail_detail

        with _mock_token(), pytest.raises(mail_policy.MailPolicyConfigError):
            await get_mail_detail(SAMPLE_EXTERNAL_MESSAGE["id"])

    async def test_lifespan_raises_on_bad_config(self, monkeypatch):
        """A pod that cannot parse its allowlist must never become ready."""
        monkeypatch.setenv(mail_policy.ENV_ALLOWED_SENDER_DOMAINS, "not a domain")
        from ms_graph_mcp import _lifespan

        with pytest.raises(mail_policy.MailPolicyConfigError):
            async with _lifespan(None):
                pass

    async def test_lifespan_starts_with_a_valid_allowlist(self, monkeypatch):
        monkeypatch.setenv(mail_policy.ENV_ALLOWED_SENDER_DOMAINS, "example.com")
        from ms_graph_mcp import _lifespan

        async with _lifespan(None):
            pass

    # -- the invariant itself -----------------------------------------------

    async def test_every_mail_tool_is_classified(self, mcp_server):
        """A new mail tool must be gated, or deliberately listed as ungated."""
        from fastmcp import Client

        async with Client(mcp_server) as client:
            names = {t.name for t in await client.list_tools()}

        classified = GATED_MAIL_TOOLS | DELIBERATELY_UNGATED_MAIL_TOOLS
        mail_tools = {n for n in names if any(word in n.lower() for word in MAIL_TOOL_WORDS)}
        assert mail_tools - classified == set()
        # And every name we classified is still registered, so the lists cannot
        # rot into a false sense of coverage.
        assert classified - names == set()


class TestCursorGuard:
    """A caller-supplied cursor that is not a Graph URL is refused with no request."""

    EVIL = "https://evil.example/v1.0/me/messages/delta?$deltatoken=x"

    @respx.mock
    @pytest.mark.parametrize(
        ("tool", "args"),
        [
            ("list_mail_delta", {"cursor": EVIL}),
            ("list_chats_page", {"cursor": EVIL}),
            ("list_chat_messages_page", {"chat_id": "19:chat", "cursor": EVIL}),
        ],
    )
    async def test_non_graph_cursor_is_refused(self, mcp_server, tool, args):
        with _mock_token():
            result = await _call(mcp_server, tool, args)
            assert _structured(result) == {"error": "invalid_cursor"}
            assert _graph_trail() == []

    @respx.mock
    async def test_graph_cursor_still_works(self, mcp_server):
        with _mock_token():
            cursor = f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages/delta?$deltatoken=abc"
            respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/mailFolders/inbox/messages/delta").mock(
                return_value=httpx.Response(200, json={"value": [], "@odata.deltaLink": cursor})
            )
            result = await _call(mcp_server, "list_mail_delta", {"cursor": cursor})
            assert _structured(result)["delta_cursor"] == cursor


class TestDraftsUnderPolicy:
    """Drafts have no from; the policy must still show the user's own drafts."""

    @respx.mock
    async def test_list_emails_drafts_folder_keeps_drafts(self, mcp_server, monkeypatch):
        _policy_on(monkeypatch)
        with _mock_token():
            route = respx.get(
                url__startswith=f"{GRAPH_BASE_URL}/me/mailFolders/drafts/messages"
            ).mock(
                return_value=httpx.Response(
                    200, json={"value": [SAMPLE_UNSENT_DRAFT, SAMPLE_EXTERNAL_MESSAGE]}
                )
            )
            result = await _call(mcp_server, "list_emails", {"folder": "drafts"})
            text = _get_text(result)
            _assert_no_canary(text)
            assert "DRAFT-SUBJECT" in text
            assert "isDraft" in _select_of(route.calls[0].request).split(",")

    @respx.mock
    async def test_read_email_shows_a_draft(self, mcp_server, monkeypatch):
        _policy_on(monkeypatch)
        with _mock_token():
            respx.get(f"{GRAPH_BASE_URL}/me/messages/{SAMPLE_UNSENT_DRAFT['id']}").mock(
                return_value=httpx.Response(200, json=SAMPLE_UNSENT_DRAFT)
            )
            result = await _call(
                mcp_server, "read_email", {"message_id": SAMPLE_UNSENT_DRAFT["id"]}
            )
            text = _get_text(result)
            assert "DRAFT-BODY" in text
            assert mail_policy.EXTERNAL_SENDER_TEXT not in text

    @respx.mock
    async def test_attachment_check_admits_a_draft(self, mcp_server, monkeypatch):
        """The id check's $select carries isDraft, so a draft's attachment is readable."""
        _policy_on(monkeypatch)
        with _mock_token():
            draft_id = SAMPLE_UNSENT_DRAFT["id"]
            respx.get(f"{GRAPH_BASE_URL}/me/messages/{draft_id}").mock(
                return_value=httpx.Response(200, json={"id": draft_id, "isDraft": True})
            )
            respx.get(url__startswith=f"{GRAPH_BASE_URL}/me/messages/{draft_id}/attachments/").mock(
                return_value=httpx.Response(200, json=SAMPLE_FILE_ATTACHMENT)
            )
            result = await _call(
                mcp_server,
                "get_mail_attachment_json",
                {
                    "message_id": draft_id,
                    "attachment_id": SAMPLE_FILE_ATTACHMENT["id"],
                    "mode": "metadata",
                },
            )
            data = _structured(result)
            assert "error" not in data
            assert data["id"] == SAMPLE_FILE_ATTACHMENT["id"]
